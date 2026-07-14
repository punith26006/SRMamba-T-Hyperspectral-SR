"""
V7 FINAL — Ruthwik-Style PPT with Animations
Exact 40 slides matching the Google Slides beige/dark/orange-gradient template.
AI-generated architecture diagrams. PowerPoint entrance animations on every slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, copy

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
AI = os.path.join(BASE, "_temp_scripts", "ai_diagrams")
MY = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")

# ─── Ruthwik Color Palette ───
BG   = RGBColor(0xEC, 0xEC, 0xE7)   # Beige
DK   = RGBColor(0x3C, 0x3C, 0x3C)   # Dark grey
HDR  = RGBColor(0xD2, 0xD2, 0xCA)   # Header bar
TT   = RGBColor(0x2B, 0x2B, 0x2B)   # Title text
TX   = RGBColor(0x40, 0x40, 0x40)   # Body text
OG   = RGBColor(0xFF, 0x8C, 0x00)   # Orange accent
WH   = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.333); H = Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]; n = 0

# ─── PowerPoint Animation XML ───
ANIM_NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def add_fade_animation(slide):
    """Add a slide-level fade transition"""
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'}
    transition = etree.SubElement(slide._element, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    transition.set('spd', 'med')
    transition.set('advClick', '1')
    fade = etree.SubElement(transition, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
    fade.set('thruBlk', '0')

def S():
    global n; n += 1
    s = prs.slides.add_slide(BL)
    add_fade_animation(s)
    return s

def bg(s):
    """Ruthwik-style background: beige + dark-grey bottom + orange gradient triangle"""
    # Main beige
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    # Header bar (slightly darker beige)
    hb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    hb.fill.solid(); hb.fill.fore_color.rgb = HDR; hb.line.fill.background()
    # Dark bottom
    db = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.0), W, Inches(1.5))
    db.fill.solid(); db.fill.fore_color.rgb = DK; db.line.fill.background()
    # Orange gradient triangle (right side)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(4.8), Inches(3.83), Inches(2.7))
    tri.rotation = 180
    tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
    # Decorative curved line (white arc)
    arc = s.shapes.add_shape(MSO_SHAPE.ARC, Inches(8), Inches(5.5), Inches(5), Inches(3))
    arc.line.color.rgb = WH; arc.line.width = Pt(2); arc.fill.background()

def title_bg(s):
    """Dark background for title slide"""
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(5.5))
    r.fill.solid(); r.fill.fore_color.rgb = DK; r.line.fill.background()
    # Orange gradient triangle
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(4.83), Inches(5.5))
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
    # Grid pattern (decorative squares)
    for row in range(3):
        for col in range(3):
            sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(9 + col*0.9), Inches(3.5 + row*0.7), Inches(0.7), Inches(0.55))
            sq.line.color.rgb = RGBColor(0x55, 0x55, 0x55); sq.line.width = Pt(1)
            sq.fill.background()

def hdr(s, t, subtitle=""):
    """Title text at top"""
    # Small colored squares in header (Ruthwik style)
    colors_sq = [RGBColor(0xBB,0xBB,0xAA), RGBColor(0xCC,0xCC,0xBB), RGBColor(0xDD,0xDD,0xCC), RGBColor(0xC5,0xC5,0xB5)]
    for i, c in enumerate(colors_sq):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3 + i*1.0), Inches(0.08), Inches(0.85), Inches(0.5))
        sq.fill.solid(); sq.fill.fore_color.rgb = c; sq.line.fill.background()
    tb = s.shapes.add_textbox(Inches(4.5), Inches(0.05), Inches(8.5), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = TT

def txt(s, items, x, y, w, h, sz=16, boxed=False):
    """Add text block, optionally in a bordered box"""
    if boxed:
        bx = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0xE5,0xE5,0xE0)
        bx.line.color.rgb = RGBColor(0xAA,0xAA,0xA0); bx.line.width = Pt(1)
    tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.font.bold = item[1]
            p.font.color.rgb = TT if item[1] else TX
        else:
            p.text = item; p.font.color.rgb = TX
        p.font.size = Pt(sz); p.space_after = Pt(4)

def add_img(s, folder, name, x, y, w, h=None):
    path = os.path.join(folder, name)
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)
    else:
        print(f"    WARN: missing {name}")

def divider(s, num, title, sub=""):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = DK; r.line.fill.background()
    # Orange accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), W, Inches(0.06))
    strip.fill.solid(); strip.fill.fore_color.rgb = OG; strip.line.fill.background()
    # Circle with number
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.0), Inches(2.1), Inches(2.1))
    c.fill.solid(); c.fill.fore_color.rgb = OG; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.2), Inches(2.1), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WH; p.alignment = PP_ALIGN.CENTER
    # Title
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.5))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(42); p2.font.bold = True; p2.font.color.rgb = WH; p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub; p3.font.size = Pt(18); p3.font.color.rgb = BG; p3.alignment = PP_ALIGN.CENTER

print("Building 40 slides — Ruthwik Style + Animations...")

# ══════════════════════════════════════════
# 1. TITLE SLIDE
# ══════════════════════════════════════════
s = S(); title_bg(s); print(f"  {n}. Title")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Exploration of Deep"; p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = WH
p2 = tf.add_paragraph(); p2.text = "Learning Models for"; p2.font.size = Pt(46); p2.font.bold = True; p2.font.color.rgb = WH
p3 = tf.add_paragraph(); p3.text = "Image Super Resolution"; p3.font.size = Pt(46); p3.font.bold = True; p3.font.color.rgb = WH
# Orange line
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.3), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = OG; ln.line.fill.background()
# Details
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(6), Inches(1))
tf2 = tb2.text_frame; tf2.word_wrap = True
for line in ["Name: Vujja Punith Sai", "Roll no: S20240010260", "Instructor: Arun PV"]:
    p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p.text = line; p.font.size = Pt(14); p.font.color.rgb = RGBColor(0xBB,0xBB,0xBB)

# ══════════════════════════════════════════
# 2. CONTENTS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Contents"); print(f"  {n}. Contents")
txt(s, [("Part I — Literature Review", True), "",
    "  1. What is SR & Need for SR", "  2. Problem Statement",
    "  3. Traditional SR Methods", "  4. Deep Learning & Evolution",
    "  5. Multimodal Learning", "  6-7. CNN: SRCNN, FSRCNN, VDSR",
    "  8-9. GAN: SRGAN & ESRGAN", " 10. Transformer: SwinIR",
    " 11. State Space: SR-Mamba", " 12. Evaluation Metrics"],
    Inches(0.4), Inches(0.9), Inches(5.8), Inches(4.8), sz=15, boxed=True)
txt(s, [("Part II — SRMamba-T Architecture", True), "",
    " 13. Why Hybrid?", " 14. Complete Architecture",
    " 15. Stage-Wise Feature Flow", " 16-17. SSM & Mamba Layer",
    " 18-19. MDSSM: 2D Spatial Context",
    " 20-21. Transformer & FFM", " 22. Training Pipeline & Loss",
    " 23-24. Quantitative & Qualitative Results",
    " 25-26. LAM & Ablation Analysis",
    " 27-28. Conclusion & References"],
    Inches(6.8), Inches(0.9), Inches(6.0), Inches(4.8), sz=15, boxed=True)

# ══════════════════════════════════════════
# 3. DIVIDER I
# ══════════════════════════════════════════
s = S(); divider(s, "I", "Literature Review", "Background & Related Work"); print(f"  {n}. Divider I")

# ══════════════════════════════════════════
# 4. WHAT IS SR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "What is Super Resolution?"); print(f"  {n}. What is SR")
txt(s, [
    "Super Resolution is Deep Learning and Computer vision technique used to reconstruct a",
    "high resolution image from one or more low resolution images.", "",
    ("→ It aims to recover:", True),
    "    1. Lost image details.  2. Textures.  3. Edges.  4. Structural information.", "",
    ("Main Objectives:", True),
    "1. Improve Image Clarity: Enhances sharpness and visibility of image contents.",
    "2. Recover Image Details: Reconstructs missing fine details and structures.",
    "3. Enhance Texture and Edges: Produces sharper texture and more realistic edges.",
    "4. Improve Perceptual Quality: Generates visually pleasing and natural-looking images.",
], Inches(0.3), Inches(0.8), Inches(8.5), Inches(5.0), sz=15)
txt(s, [("Basic Workflow", True), "", "Low Resolution Image", "    ↓", "Feature Extraction",
    "    ↓", "Super Resolution Model", "    ↓", "Image Reconstruction Layer", "    ↓", "High Resolution Image"],
    Inches(9.0), Inches(0.9), Inches(3.8), Inches(4.5), sz=15, boxed=True)

# ══════════════════════════════════════════
# 5. NEED FOR SR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Need for Super Resolution"); print(f"  {n}. Need for SR")
txt(s, [
    ("Problems in Low Resolution Images:", True),
    "• Blur and loss of sharpness", "• Pixelation during zooming",
    "• Loss of fine textures and edges", "• Poor object visibility",
    "• Noise and distortion", "• Information loss during compression",
], Inches(0.3), Inches(0.9), Inches(6.0), Inches(4.8), sz=15, boxed=True)
txt(s, [
    ("Real-World Applications:", True),
    "• Satellite Imaging: terrain analysis, urban planning",
    "• Surveillance: face recognition, number plates",
    "• Medical Imaging: MRI/CT scan enhancement",
    "• Mobile Camera: AI zoom, night photography",
    "• Video Upscaling: low-quality → HD/4K conversion",
    "• Remote Sensing: disaster monitoring, geology",
], Inches(6.6), Inches(0.9), Inches(6.2), Inches(4.8), sz=15, boxed=True)

# ══════════════════════════════════════════
# 6. PROBLEM STATEMENT
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Problem Statement"); print(f"  {n}. Problem Statement")
txt(s, [
    "Most open-source data have coarse spatial resolution, and it is advantageous to use different",
    "sensors for effective enhancement; however, the sensors have differences in spatial resolution,",
    "spectral characteristics, noise levels, and often lack spatial or temporal overlap.",
], Inches(0.3), Inches(0.9), Inches(12.5), Inches(1.2), sz=17, boxed=True)
txt(s, [("Key Challenges in SISR:", True),
    "• Recovering lost high-frequency details from degraded inputs",
    "• CNN models: limited receptive fields (local only)",
    "• Transformer models: O(n²) computational complexity — GPU memory explosion",
    "• Pure Mamba models: 1D sequential scan destroys 2D spatial context",
    "• Need for a hybrid approach combining efficiency + global awareness",
], Inches(0.3), Inches(2.5), Inches(12.5), Inches(3.2), sz=16, boxed=True)

# ══════════════════════════════════════════
# 7. TRADITIONAL METHODS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Traditional Super Resolution Methods"); print(f"  {n}. Traditional")
txt(s, [
    "Before Deep Learning, SR used interpolation-based image scaling methods.", "",
    ("1) Nearest Neighbour:", True), "  Copies nearest pixel value. Very fast. Blocky, jagged edges.", "",
    ("2) Bilinear Interpolation:", True), "  Uses neighbouring pixels. Smoother but blurry edges.", "",
    ("3) Bicubic Interpolation:", True), "  Uses surrounding 16 pixels. Better smoothness but cannot reconstruct missing details.", "",
    ("→ All methods fail to reconstruct truly missing high-frequency details → Need for Deep Learning.", True),
], Inches(0.3), Inches(0.9), Inches(12.5), Inches(4.8), sz=16)

# ══════════════════════════════════════════
# 8. DL IN SR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Deep Learning in Super Resolution"); print(f"  {n}. DL in SR")
txt(s, [
    ("Why Deep Learning?", True),
    "Instead of enlarging pixels mathematically, DL models",
    "learn the mapping from LR to HR images from datasets.", "",
    ("Deep Learning Models Learn:", True),
    "• Texture Patterns — roads, vegetation, faces",
    "• Edge Information — boundaries, transitions",
    "• Image Structures — spatial relationships", "",
    ("Main Idea:", True), "Direct LR → HR mapping via neural networks.",
], Inches(0.3), Inches(0.9), Inches(5.8), Inches(4.8), sz=15)
add_img(s, AI, "evolution.png", Inches(6.3), Inches(0.9), Inches(6.5), Inches(4.8))

# ══════════════════════════════════════════
# 9. MULTIMODAL
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Relation to Multimodal Learning"); print(f"  {n}. Multimodal")
txt(s, [
    ("Single-Modality SR:", True),
    "• Uses a single low-resolution image.", "• Reconstruction depends entirely on visual features.", "",
    ("Multimodal Learning:", True),
    "Combines multiple data sources (e.g., satellite + elevation → HR output).", "",
    ("Advantages:", True), "Better context, richer features.",
    ("Limitations:", True), "Data alignment challenges, higher complexity.",
], Inches(0.3), Inches(0.9), Inches(12.5), Inches(4.8), sz=16)

# ══════════════════════════════════════════
# 10. SRCNN
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "CNN-based: SRCNN (2014)"); print(f"  {n}. SRCNN")
txt(s, [
    ("SRCNN — First DL model for Super Resolution", True),
    "End-to-end learning for SR using CNNs.", "",
    ("Architecture: 3 Convolution Layers", True),
    "1) Feature Extraction: low-level features",
    "2) Non-Linear Mapping: LR→HR features",
    "3) Reconstruction: final HR image", "",
    ("Loss: MSE", True), "L(Θ) = (1/n) Σ ||F(Yᵢ;Θ) - Xᵢ||²",
], Inches(0.3), Inches(0.9), Inches(5.5), Inches(4.8), sz=14)
add_img(s, AI, "srcnn.png", Inches(6.0), Inches(1.0), Inches(6.8), Inches(4.5))

# ══════════════════════════════════════════
# 11. FSRCNN + VDSR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "FSRCNN (2016) & VDSR (2016)"); print(f"  {n}. FSRCNN+VDSR")
txt(s, [
    ("FSRCNN — Fast Super Resolution CNN", True),
    "• Operates on LR space (no pre-upsampling).",
    "• Uses deconvolution for upsampling at end.",
    "• 5× faster than SRCNN.",
    "• Lim: Still limited receptive field.",
], Inches(0.3), Inches(0.9), Inches(6.0), Inches(4.8), sz=15, boxed=True)
txt(s, [
    ("VDSR — Very Deep SR (20 Layers)", True),
    "• Global Residual: HR = LR + Residual.",
    "• 41×41 pixel receptive field.",
    "• Multi-scale: handles ×2, ×3, ×4.",
    "• Lim: Operates in HR space (slow).",
], Inches(6.6), Inches(0.9), Inches(6.2), Inches(4.8), sz=15, boxed=True)

# ══════════════════════════════════════════
# 12. SRGAN
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "GAN-based: SRGAN (2017)"); print(f"  {n}. SRGAN")
txt(s, [
    ("Why GAN-based SR?", True),
    "CNN models optimized for MSE → blurry outputs.",
    "SRGAN: Adversarial Learning for perceptual realism.", "",
    ("Architecture:", True),
    "• Generator: Deep ResNet → produces HR from LR.",
    "• Discriminator: VGG → evaluates real vs fake.", "",
    ("Key Innovation: Perceptual Loss", True),
    "Uses VGG features → sharper textures.",
], Inches(0.3), Inches(0.9), Inches(5.5), Inches(4.8), sz=14)
add_img(s, AI, "srgan.png", Inches(6.0), Inches(1.0), Inches(6.8), Inches(4.5))

# ══════════════════════════════════════════
# 13. ESRGAN
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "GAN-based: ESRGAN (2018)"); print(f"  {n}. ESRGAN")
txt(s, [
    ("ESRGAN — Enhanced SRGAN", True),
    "• Residual-in-Residual Dense Blocks (RRDB).",
    "• Removed batch normalization for stable training.",
    "• Relativistic Discriminator: predicts relative quality.", "",
    ("Improvements:", True),
    "• Better texture, higher perceptual quality.",
    "• More stable training, sharper outputs.",
    "• Won PIRM2018 SR Challenge.", "",
    ("Limitations of GANs:", True), "Training instability, mode collapse, hallucinated textures.",
], Inches(0.3), Inches(0.9), Inches(12.5), Inches(4.8), sz=16)

# ══════════════════════════════════════════
# 14. SWINIR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Transformer: SwinIR (2021)"); print(f"  {n}. SwinIR")
txt(s, [
    ("SwinIR — Swin Transformer for Image Restoration", True),
    "• Built on Vision Transformers (ViT).",
    "• Uses Shifted Window Attention.",
    "• Self-attention determines which image regions",
    "  are most important for reconstruction.", "",
    ("Adv:", True), "  Global feature learning, excellent quality.",
    ("Lim:", True), "  O(n²) complexity, high GPU memory cost.",
], Inches(0.3), Inches(0.9), Inches(5.5), Inches(4.8), sz=14)
add_img(s, AI, "swinir.png", Inches(6.0), Inches(1.0), Inches(6.8), Inches(4.5))

# ══════════════════════════════════════════
# 15. SR-MAMBA
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "State Space: SR-Mamba (2024)"); print(f"  {n}. SR-Mamba")
txt(s, [
    ("State Space Models (SSMs):", True),
    "• Sequence modeling with O(n) linear complexity.",
    "• Avoids expensive global attention operations.", "",
    ("SR-Mamba:", True),
    "• Modern architecture based on Mamba/S6.",
    "• State Transition Mechanisms for efficiency.", "",
    ("Limitations:", True),
    "• 1D sequential scan destroys 2D spatial context.",
    "→ Need for hybrid approach: SRMamba-T!",
], Inches(0.3), Inches(0.9), Inches(5.5), Inches(4.8), sz=14)
add_img(s, AI, "ssm_mamba.png", Inches(6.0), Inches(1.0), Inches(6.8), Inches(4.5))

# ══════════════════════════════════════════
# 16. METRICS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Evaluation Metrics: MSE, PSNR, SSIM"); print(f"  {n}. Metrics")
txt(s, [
    ("1. MSE (Mean Square Error)", True),
    "Pixel-wise reconstruction error. Lower is better.",
    "MSE = (1/mn) Σ (I(i,j) - K(i,j))². Used in SRCNN, VDSR.",
], Inches(0.3), Inches(0.9), Inches(3.9), Inches(4.5), sz=14, boxed=True)
txt(s, [
    ("2. PSNR (Peak Signal-to-Noise)", True),
    "Reconstruction quality in dB. Higher is better.",
    "PSNR = 10·log₁₀(MAX²/MSE). Inversely related to MSE.",
], Inches(4.5), Inches(0.9), Inches(3.9), Inches(4.5), sz=14, boxed=True)
txt(s, [
    ("3. SSIM (Structural Similarity)", True),
    "Structural & perceptual quality measure.",
    "Compares luminance, contrast, structure.",
    "SSIM → 1: excellent.  SSIM → 0: poor.",
], Inches(8.8), Inches(0.9), Inches(4.1), Inches(4.5), sz=14, boxed=True)

# ══════════════════════════════════════════
# 17. DIVIDER II
# ══════════════════════════════════════════
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution"); print(f"  {n}. Divider II")

# ═══ PART II helper ═══
def p2_slide(title, items, img_name, img_folder=MY):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:40]}")
    txt(s, items, Inches(0.3), Inches(0.9), Inches(5.5), Inches(4.8), sz=14)
    add_img(s, img_folder, img_name, Inches(6.0), Inches(1.0), Inches(6.8), Inches(4.7))

def p2_full(title, items, img_name, img_folder=MY):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:40]}")
    txt(s, items, Inches(0.3), Inches(0.8), Inches(12.5), Inches(1.0), sz=14)
    add_img(s, img_folder, img_name, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.8))

# 18. WHY HYBRID
p2_slide("Why Hybrid? — The Architecture Compromise", [
    ("No single architecture is perfect:", True), "",
    ("CNN:", True), "✗ Limited receptive field (local only).", "",
    ("Transformer:", True), "✓ Global context via self-attention.",
    "✗ O(n²) complexity — GPU memory explosion.", "",
    ("Mamba (SSM):", True), "✓ O(n) linear complexity — efficient.",
    "✗ 1D scan destroys 2D spatial information.", "",
    ("→ SRMamba-T: Best of Mamba + Transformer.", True),
], "erf.png", AI)

# 19. ERF ANALYSIS
p2_slide("CNN vs Transformer vs Mamba — ERF Analysis", [
    ("Effective Receptive Field (ERF):", True), "",
    "CNN: Small, localized ERF.",
    "Transformer: Large but uniform ERF.",
    "Mamba: Directional ERF along scan line.", "",
    ("The Hybrid Insight:", True),
    "Phase 1 — Mamba (early layers):",
    "  Dense feature extraction at O(n) cost.",
    "Phase 2 — Transformer (later layers):",
    "  Global refinement on compressed features.",
], "slide_07.png")

# 20. COMPLETE ARCHITECTURE (PAPER DIAGRAM)
s = S(); bg(s); hdr(s, "The Complete SRMamba-T Architecture"); print(f"  {n}. Complete Arch")
txt(s, [
    "Three-stage pipeline: (1) Shallow Feature Extraction via 3×3 Conv → (2) Deep Feature Extraction with N/2 Mamba Layers (encoder)",
    "followed by (N/2)-1 Transformer Layers (encoder) → (3) Image Reconstruction via FFM + Decoder Transformer + PixelShuffle Upsample.",
], Inches(0.3), Inches(0.8), Inches(12.5), Inches(0.8), sz=14)
add_img(s, AI, "srmambat_paper.png", Inches(0.2), Inches(1.8), Inches(12.9), Inches(3.8))
txt(s, [
    ("Key Components:", True),
    "• Mamba Block: LayerNorm → MambaMixer → Scale → LayerNorm → 3×3Conv → Channel Attention (CAB) → Scale",
    "• Attention Block: LayerNorm → Self-Attention → Scale → LayerNorm → 3×3Conv → Channel Attention → Scale",
    "• FFM: Aggregates features from all encoder stages via DWConv → SiLU → PWLinear → LayerNorm",
], Inches(0.3), Inches(5.5), Inches(12.5), Inches(0.7), sz=11)

# 21. STAGE-WISE FLOW
s = S(); bg(s); hdr(s, "Stage-Wise Feature Flow"); print(f"  {n}. Flow")
txt(s, [
    "Input → 3×3Conv (Fₛ) → PatchEmbed → N/2 Mamba Layers (F₁) → Mamba Layer → (N/2)-1 Transformer Layers (Fₕ) → 3×3Conv → FFM (F_f)",
    "→ Decoder Transformer Layer (F_d) → PatchUnEmbed → Upsample → SR Output",
], Inches(0.3), Inches(0.8), Inches(12.5), Inches(0.8), sz=14)
add_img(s, AI, "srmambat_paper.png", Inches(0.2), Inches(1.8), Inches(12.9), Inches(3.8))
txt(s, [
    ("Detailed Data Flow:", True),
    "• Shallow: LR → Fₛ = Conv3×3(LR) → F₁ = PatchEmbed(Fₛ) — converts 2D features to 1D token sequences",
    "• Encoder: Each Mamba/Trans Layer has L blocks + 3×3Conv. Uses Scale factors S₁, S₂ for residual balancing.",
    "• FFM: Concatenates Fₛ with all encoder features → 3×3 DWConv → SiLU → PWLinear → LayerNorm",
    "• Reconstruction: F = Fₛ + F'd (element-wise add shallow features back) → PixelShuffle upsample to HR",
], Inches(0.3), Inches(5.5), Inches(12.5), Inches(0.7), sz=11)

# 22. SSM FOUNDATIONS
p2_slide("SSM Foundations & Discretization", [
    ("State Space Models (SSMs):", True),
    "Continuous-time linear dynamical system:",
    "  h'(t) = A·h(t) + B·x(t)    (state equation)",
    "  y(t) = C·h(t) + D·x(t)      (output equation)", "",
    "Where A ∈ ℝᴺˣᴺ is the state matrix, B ∈ ℝᴺˣ¹",
    "is the input matrix, C ∈ ℝ¹ˣᴺ is the output matrix.", "",
    ("Discretization via Zero-Order Hold (ZOH):", True),
    "  Ā = exp(Δ·A)    — state transition matrix",
    "  B̄ = (Δ·A)⁻¹(exp(Δ·A) - I)·Δ·B",
    "  hₖ = Ā·hₖ₋₁ + B̄·xₖ  ;  yₖ = C·hₖ", "",
    ("Selective Scan (Mamba S6):", True),
    "Makes Δ, B, C input-dependent → content-aware.",
    "This is the key innovation over S4/HiPPO.",
], "slide_11.png")

# 23. SSM CONVOLUTION
p2_slide("SSM: Equivalence to Convolution", [
    ("Key Mathematical Insight:", True),
    "The discretized SSM can be reformulated as a",
    "global convolution — enabling parallel training.", "",
    ("Convolutional Kernel:", True),
    "  K̄ = (CB̄, CĀB̄, CĀ²B̄, ..., CĀᴸ⁻¹B̄)",
    "  y = x * K̄  (global convolution with kernel K̄)", "",
    ("Dual Computation Modes:", True),
    "• Training: Parallel FFT-based convolution.",
    "  Complexity: O(L log L) — highly parallelizable.",
    "• Inference: Sequential recurrence hₖ = Āhₖ₋₁ + B̄xₖ.",
    "  Complexity: O(1) per step — memory efficient.", "",
    "This duality makes SSMs uniquely suited for SR.",
], "slide_13.png")

# 24. MAMBA LAYER
p2_slide("Mamba Layer: Anatomy & Mathematics", [
    ("Each Mamba Layer contains L Mamba Blocks:", True),
    "X̂ᵢ,ₗ = Mamba(Norm(Xᵢ,ₗ)) + S₁ × Xᵢ,ₗ  (Eq.6)", "",
    ("Inside each Mamba Block:", True),
    "Input → LayerNorm → Linear Proj (expand 2×) →",
    "  Branch 1: Conv1D(k=3) → SiLU → SSM (MDSSM)",
    "  Branch 2: SiLU activation (gate)",
    "→ Element-wise multiply → Linear (compress) → Out", "",
    "After L blocks: Xᵢ,ₗ₊₁ = CAB(Conv3×3(Norm(X̂ᵢ,ₗ))) + S₂ × X̂ᵢ,ₗ  (Eq.7)", "",
    ("Channel Attention Block (CAB):", True),
    "Enhances channel-wise feature importance",
    "while avoiding redundancy across channels.",
], "slide_14.png")

# 25. ASYMMETRIC MIXER
p2_slide("The Asymmetric Mamba Mixer Block", [
    ("Why Asymmetric?", True),
    "Standard Mamba scans in one direction → misses",
    "context from the opposite direction.", "",
    ("Mamba Mixer Architecture:", True),
    "• Two parallel branches with shared input:",
    "  Branch A: Linear → Conv1D → SiLU → MDSSM",
    "  Branch B: Linear → SiLU (gating mechanism)",
    "• Output = Linear(BranchA ⊙ BranchB)", "",
    ("Key Difference from Vanilla Mamba:", True),
    "• Uses MDSSM instead of standard SSM.",
    "• Branch A includes Conv1D for local context.",
    "• Asymmetric parameter allocation per direction.",
    "• Learnable scale factors S₁, S₂ for residual.",
], "slide_15.png")

# 26. SPATIAL DEPENDENCY
p2_slide("The Spatial Dependency Problem", [
    ("Why 1D scanning fails for images:", True),
    "Images are 2D grids of size H×W. Standard Mamba",
    "flattens this to a 1D sequence of length H·W.", "",
    ("Information Loss:", True),
    "Pixel at (i,j) and (i+1,j) are vertical neighbors",
    "in 2D but are W positions apart in 1D sequence.",
    "The SSM's exponential decay exp(Δ·A) causes the",
    "state to 'forget' vertical neighbor information.", "",
    ("Consequence for SR:", True),
    "Vertical edges, horizontal textures, and diagonal",
    "patterns cannot be reconstructed properly.",
    "→ Need MDSSM: Multi-Directional Selective Scan.",
], "mdssm.png", AI)

# 27. MDSSM
p2_slide("Multi-Directional Selective Scan (MDSSM)", [
    ("MDSSM — Core Innovation of SRMamba-T:", True),
    "1) Channel Split: D channels → 4 groups of D/4.",
    "2) 4-Direction Scanning:",
    "   → Top-to-Bottom (vertical dependencies)",
    "   → Bottom-to-Top (reverse vertical)",
    "   → Left-to-Right (horizontal dependencies)",
    "   → Right-to-Left (reverse horizontal)",
    "3) Each group processed by SSM in its direction.",
    "4) Concatenate all 4 groups back together.", "",
    ("Computational Cost:", True),
    "Same as single-direction SSM (D/4 × 4 = D)!",
    "Zero additional parameter overhead.",
], "slide_17.png")

# 28. MDSSM MATH
p2_slide("MDSSM: Mathematical Formulation", [
    ("Formal Definition (Eq.8):", True),
    "  X = [X₁, X₂, X₃, X₄] = ChannelSplit(X)", "",
    "  X₁' = SSM(ScanTopToBottom(X₁))",
    "  X₂' = SSM(ScanBottomToTop(X₂))",
    "  X₃' = SSM(ScanLeftToRight(X₃))",
    "  X₄' = SSM(ScanRightToLeft(X₄))", "",
    "  Output = Concat(X₁', X₂', X₃', X₄')", "",
    ("Why this works:", True),
    "Each pixel receives context from all 4 spatial",
    "directions through its respective channel group.",
    "Spatial locality is fully preserved at O(n) cost.",
], "slide_18.png")

# 29. TRANSFORMER LAYER
p2_slide("Inside the Transformer Layer", [
    ("Each Transformer Layer has L Attention Blocks:", True),
    "X̂ᵢ,ₗ = Attention(Norm(Xᵢ,ₗ)) + S₁ × Xᵢ,ₗ  (Eq.6)", "",
    ("Multi-Head Self-Attention:", True),
    "  Q = XWq,  K = XWk,  V = XWv",
    "  Attention = softmax(QKᵀ/√d_head) · V", "",
    ("After L blocks:", True),
    "Xᵢ,ₗ₊₁ = CAB(Conv3×3(Norm(X̂ᵢ,ₗ))) + S₂ × X̂ᵢ,ₗ",
    "(Same CAB + Conv structure as Mamba Layer.)", "",
    ("Key: Un-Windowed Attention", True),
    "Full global attention (no window restriction).",
], "slide_19.png")

# 30. UN-WINDOWED ATTENTION
p2_slide("Un-Windowed Global Self-Attention", [
    ("Problem with SwinIR's Windowed Attention:", True),
    "• Window size limits receptive field to w×w pixels.",
    "• Cross-window info requires shifted windows.",
    "• Boundary artifacts at window edges.", "",
    ("SRMamba-T's Solution:", True),
    "Full un-windowed attention: every token attends to",
    "ALL other tokens in the feature map.", "",
    ("Why this is affordable:", True),
    "• Mamba layers already extracted dense features.",
    "• Feature maps are compressed by this stage.",
    "• Only (N/2)-1 Transformer layers needed.",
    "• Result: True global context, no artifacts.",
], "slide_20.png")

# 31. FFM
p2_slide("Feature Fusion Module (FFM)", [
    ("U-Net Bottleneck Design:", True),
    "FFM aggregates hierarchical features from ALL",
    "encoder stages (both Mamba and Transformer).", "",
    ("FFM Internal Structure:", True),
    "1) Aggregate: Fₛ (shallow) + all encoder features.",
    "2) 3×3 DWConv: depth-wise spatial aggregation.",
    "3) SiLU activation: non-linear transformation.",
    "4) PWLinear: point-wise channel compression.",
    "5) LayerNorm: feature regularization.", "",
    ("Output:", True),
    "F_f = FFM(Fₛ, F₁, ..., Fₙ) → feeds decoder.",
], "slide_21.png")

# 32. RECONSTRUCTION + LOSS
p2_slide("Reconstruction & Dual-Domain Loss (Eq.5)", [
    ("Reconstruction Pipeline:", True),
    "F = Fs + F'd (element-wise add shallow features).",
    "Lightweight SR: Conv + PixelShuffle(r) for upscale.",
    "Classical SR: Multi-step progressive upsample.", "",
    ("Dual-Domain Loss (Eq.5):", True),
    "L = Sum ||I_HR - I_SR||_2 + lambda ||FFT(I_HR) - FFT(I_SR)||_2", "",
    "L_pixel (L2): Pixel-wise reconstruction accuracy.",
    "L_freq (FFT): Preserves high-frequency edges/textures.",
    "lambda = 0.1 (empirically tuned balance factor).", "",
    ("Why dual-domain?", True),
    "Pixel loss alone produces over-smoothed results.",
    "FFT loss ensures sharp edges and fine textures.",
], "slide_22.png")

# 33. TRAINING
p2_slide("Training Pipeline & Setup", [
    ("SRMamba-T-S Model Configuration:", True),
    "N=4 layers (2 Mamba + 1 Trans encoder + 1 Trans decoder).",
    "L=6 blocks per layer. Feature dimension D=48.",
    "Total parameters: ~0.6M (lightweight variant).", "",
    ("Training Details:", True),
    "Dataset: DIV2K (800 images) + random augmentation.",
    "LR Patch: 64x64, Batch: 16, Iterations: 500K.",
    "Optimizer: Adam (beta1=0.9, beta2=0.99).",
    "LR: 2e-4 with cosine annealing to 1e-6.",
    "Degradation: Bicubic downsampling (x2, x3, x4).",
    "Hardware: NVIDIA A100, PyTorch 2.0, CUDA 11.8.",
], "slide_23.png")

# 34. QUANTITATIVE
p2_slide("Quantitative Results: PSNR vs Compute", [
    ("Key Results:", True),
    "SRMamba-T achieves:",
    "• +0.28 dB PSNR improvement over MambaIR.",
    "• 38.7% fewer MACs (computational cost).",
    "• Better SSIM on all benchmark datasets.", "",
    ("Benchmarks:", True),
    "Set5, Set14, BSD100, Urban100, Manga109.", "",
    "Outperforms ALL baselines with less compute.",
], "slide_24.png")

# 35. QUALITATIVE
p2_slide("Qualitative Validation: Visual Comparison", [
    ("Urban100 dataset comparison:", True),
    "• Complex textures (buildings, fences, grids).", "",
    ("Observations:", True),
    "• CNNs: Blurry, loss of detail.",
    "• GANs: Sharp but artifacts.",
    "• SwinIR: Misses some textures.",
    "• MambaIR: Directional artifacts.", "",
    ("SRMamba-T:", True),
    "Sharpest reconstruction, no directional artifacts.",
], "visual_compare.png", AI)

# 36. LAM
p2_slide("LAM Analysis: Inside the Model's Mind", [
    ("LAM = Local Attribution Map", True),
    "Visualizes which input pixels contributed most", "",
    ("Diffusion Index (DI):", True),
    "• CNN: Low DI — very local.",
    "• SwinIR: Medium DI — window-limited.",
    "• MambaIR: Higher DI but directional.",
    "• SRMamba-T: Highest DI — truly global.", "",
    "SRMamba-T uses the widest range of input pixels.",
], "slide_26.png")

# 37. ABLATION
p2_slide("Hardware Efficiency & Ablation Study", [
    ("Memory Scaling:", True),
    "• Transformer: Quadratic O(n²) growth.",
    "• SRMamba-T: Near-linear scaling.",
    "At 1024×1024: SwinIR → OOM, SRMamba-T runs fine.", "",
    ("Ablation: Layer Ordering:", True),
    "A) All Mamba → Poor global context.",
    "B) All Transformer → Too expensive.",
    "C) Trans first, Mamba later → Suboptimal.",
    "D) Mamba first, Trans later → BEST ✓.",
], "slide_27.png")

# ══════════════════════════════════════════
# 38. CONCLUSION
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Conclusion & Future Work"); print(f"  {n}. Conclusion")
txt(s, [
    ("Key Contributions:", True),
    "✓ First hybrid Mamba-Transformer for image SR.",
    "✓ MDSSM: 4-directional scanning at O(n) cost.",
    "✓ Asymmetric: Mamba early + Transformer late.",
    "✓ +0.28 dB PSNR, 38.7% fewer MACs.",
    "✓ Dual-domain loss prevents artifacts.",
    "✓ Linear memory scaling for large images.",
], Inches(0.3), Inches(0.9), Inches(6.0), Inches(4.8), sz=15, boxed=True)
txt(s, [
    ("Future Work:", True),
    "• Extend to Hyperspectral SR.",
    "• Scale to SRMamba-T-L for classical SR.",
    "• Real-world degradation beyond bicubic.",
    "• Edge device deployment (linear memory).",
    "• Combine with diffusion models.",
], Inches(6.6), Inches(0.9), Inches(6.2), Inches(4.8), sz=15, boxed=True)

# ══════════════════════════════════════════
# 39. REFERENCES
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "References"); print(f"  {n}. References")
txt(s, [
    "[1] Dong et al., \"SRCNN\", ECCV 2014",
    "[2] Dong et al., \"FSRCNN\", ECCV 2016",
    "[3] Kim et al., \"VDSR\", CVPR 2016",
    "[4] Ledig et al., \"SRGAN\", CVPR 2017",
    "[5] Wang et al., \"ESRGAN\", ECCV 2018",
    "[6] Liang et al., \"SwinIR\", ICCV 2021",
    "[7] Guo et al., \"MambaIR\", ECCV 2024",
    "[8] Gu & Dao, \"Mamba: Selective State Spaces\", 2023",
    "[9] SRMamba-T: Hybrid Mamba-Transformer for SISR, 2024",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017",
], Inches(0.3), Inches(0.9), Inches(12.5), Inches(4.8), sz=15, boxed=True)

# ══════════════════════════════════════════
# 40. THANK YOU
# ══════════════════════════════════════════
s = S(); print(f"  {n}. Thank You")
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
r.fill.solid(); r.fill.fore_color.rgb = DK; r.line.fill.background()
# Orange triangle
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(4.83), Inches(5.5))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
# Grid squares
for row in range(3):
    for col in range(3):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(9 + col*0.9), Inches(3.5 + row*0.7), Inches(0.7), Inches(0.55))
        sq.line.color.rgb = RGBColor(0x55, 0x55, 0x55); sq.line.width = Pt(1)
        sq.fill.background()
# Orange line
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.8), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = OG; ln.line.fill.background()
# Thank you text
tb = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7), Inches(1.5))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(64); p.font.bold = True; p.font.color.rgb = OG
# Subtitle
tb2 = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WH
# Name
tb3 = s.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.5))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; p3.text = "Vujja Punith Sai | S20240010260 | Guide: Arun PV"
p3.font.size = Pt(16); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA)

# ══════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v8_FINAL.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
