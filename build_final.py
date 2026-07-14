"""
Final comprehensive script:
1. Creates a combined PPT with:
   - 10 condensed literature-review slides from friend's PDF (front)
   - Problem statement after intro
   - User's SRMamba-T architecture slides (cleaned, no redundancy) after
2. Organizes btpdataset into separate folders: PPTs/, Code/, Dataset/
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import shutil

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
FRIEND = os.path.join(BASE, "friend_slides")
MINE = os.path.join(BASE, "slide_images")

def f(n):
    """Friend's slide (rotated)"""
    return os.path.join(FRIEND, f"slide_{n:02d}_rot.png")

def m(n):
    """My slide"""
    return os.path.join(MINE, f"slide_{n:02d}.png")

def add_img(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width)
    else:
        print(f"  WARNING: Missing image: {path}")

def add_title_text(slide, text, size=28):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

def add_body_text(slide, text, left, top, width, height, size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)

def create_ppt():
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank = prs.slide_layouts[6]  # blank layout

    print("=== PART 1: 10 Literature Review Slides (from friend's PDF) ===")

    # ---------- SLIDE 1: Title Slide ----------
    print("  Slide 1: Title")
    slide = prs.slides.add_slide(blank)
    # Dark background box
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Exploration of Deep Learning Models for\nImage Super Resolution"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    lines = [
        "B.Tech Project Presentation",
        "",
        "Presented by: Vujja Punith Sai",
        "Roll No: S20240010260",
        "Instructor / Guide: Arun PV",
        "",
        "Institute of Intelligent Computing"
    ]
    for i, line in enumerate(lines):
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(22) if i == 0 else Pt(20)
        p.font.bold = (i == 0)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC) if i > 0 else RGBColor(0xFF, 0xD7, 0x00)
        p.alignment = PP_ALIGN.CENTER

    # ---------- SLIDE 2: What is SR + Need for SR ----------
    print("  Slide 2: Introduction (What is SR + Need for SR)")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "1. Introduction to Super Resolution")
    # Friend's slide 3 (What is SR) + slide 4 (Need for SR) side by side
    add_img(slide, f(3), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(4), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 3: Problem Statement ----------
    print("  Slide 3: Problem Statement")
    slide = prs.slides.add_slide(blank)
    # Background
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.fill.background()
    
    add_title_text(slide, "2. Problem Statement", size=36)
    
    # Problem statement box
    shape2 = slide.shapes.add_shape(
        1, Inches(1), Inches(2), Inches(11), Inches(3.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape2.line.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape2.line.width = Pt(2)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = (
        "Most open-source data have coarse spatial resolution, and it is advantageous "
        "to use different sensors for effective enhancement; however, the sensors have "
        "differences in spatial resolution, spectral characteristics, noise levels, and "
        "often lack spatial or temporal overlap."
    )
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.line_spacing = Pt(36)

    # ---------- SLIDE 4: Traditional Methods + DL in SR ----------
    print("  Slide 4: Traditional Methods + Deep Learning in SR")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "3. Traditional Methods & Deep Learning in SR")
    add_img(slide, f(5), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(6), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 5: Evolution + Multimodal ----------
    print("  Slide 5: Evolution of SR Models")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "4. Evolution of SR Models & Multimodal Learning")
    add_img(slide, f(7), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(8), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 6: CNN-based Models (SRCNN) ----------
    print("  Slide 6: CNN-based Models (SRCNN)")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "5. CNN-based Models: SRCNN, FSRCNN, VDSR")
    add_img(slide, f(10), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(12), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 7: GAN-based Models ----------
    print("  Slide 7: GAN-based Models (SRGAN, ESRGAN)")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "6. GAN-based Models: SRGAN & ESRGAN")
    add_img(slide, f(19), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(20), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 8: Transformer + SwinIR ----------
    print("  Slide 8: Transformer-based Models (SwinIR)")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "7. Transformer-based Models: SwinIR")
    add_img(slide, f(26), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(27), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 9: State Space Models (SR-Mamba) ----------
    print("  Slide 9: State Space Models (SR-Mamba)")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "8. State Space Models: SR-Mamba")
    add_img(slide, f(29), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(30), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    # ---------- SLIDE 10: Evaluation Metrics ----------
    print("  Slide 10: Evaluation Metrics & Conclusion")
    slide = prs.slides.add_slide(blank)
    add_title_text(slide, "9. Evaluation Metrics (MSE, PSNR, SSIM) & Literature Conclusion")
    add_img(slide, f(32), Inches(0.2), Inches(0.9), Inches(6.4), Inches(6.0))
    add_img(slide, f(33), Inches(6.8), Inches(0.9), Inches(6.4), Inches(6.0))

    print("\n=== PART 2: SRMamba-T Architecture Slides (from my PPT, cleaned) ===")
    # From the user's 31 slides, pick unique non-redundant ones in correct order.
    # Redundant ones removed: slides 2,3,5,8,10,11,12,13,16,23 (covered above or repetitive)
    
    my_slides = [
        # (slide_num, title)
        (9,  "10. SRMamba-T: Complete Architecture Overview"),
        (14, "11. Mamba Block: Anatomy & Mathematics (Eq. 6, 7)"),
        (15, "12. Asymmetric Mamba Mixer Block"),
        (17, "13. Multi-Directional Selective Scan Module (MDSSM)"),
        (18, "14. MDSSM: Preserving 2D Spatial Context"),
        (19, "15. Inside the Transformer Layer"),
        (20, "16. Un-Windowed Global Self-Attention"),
        (21, "17. Feature Fusion Module (FFM)"),
        (22, "18. Adaptive Reconstruction & Dual-Domain Loss"),
        (24, "19. Quantitative Results: PSNR vs Compute"),
        (25, "20. Qualitative Validation: Complex Textures"),
        (26, "21. LAM Analysis: Diffusion Index Comparison"),
        (27, "22. Hardware Efficiency & Memory Scaling"),
        (28, "23. Ablation Study: Layer Ordering"),
    ]
    
    for snum, title in my_slides:
        print(f"  Slide {title.split('.')[0]}: {title}")
        slide = prs.slides.add_slide(blank)
        # Use the full slide image at full size
        add_img(slide, m(snum), Inches(0), Inches(0), prs.slide_width, prs.slide_height)

    # ---------- FINAL SLIDE: Thank You ----------
    print("  Final Slide: Thank You")
    slide = prs.slides.add_slide(blank)
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(2))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "Vujja Punith Sai  |  S20240010260  |  Guide: Arun PV"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    # Save
    ppt_dir = os.path.join(BASE, "PPTs")
    os.makedirs(ppt_dir, exist_ok=True)
    out_path = os.path.join(ppt_dir, "SRMamba_T_BTP_Final.pptx")
    prs.save(out_path)
    print(f"\n>>> Presentation saved to: {out_path}")
    print(f">>> Total slides: {len(prs.slides)}")
    return out_path

def organize_folders():
    """Organize btpdataset into PPTs/, Code/, Dataset/ folders"""
    print("\n=== Organizing Folder Structure ===")
    
    ppts_dir = os.path.join(BASE, "PPTs")
    code_dir = os.path.join(BASE, "Code")
    dataset_dir = os.path.join(BASE, "Dataset")
    
    os.makedirs(ppts_dir, exist_ok=True)
    os.makedirs(code_dir, exist_ok=True)
    os.makedirs(dataset_dir, exist_ok=True)
    
    # Move PPT files
    ppt_files = [
        "btpproject.pptx",
        "SRMamba_T_Final_10_Slides.pptx",
        "SRMamba_T_Final_Corrected.pptx",
        "slidesruthwik.pdf",
        "btpmainpdf.pdf",
    ]
    for pf in ppt_files:
        src = os.path.join(BASE, pf)
        dst = os.path.join(ppts_dir, pf)
        if os.path.exists(src) and not os.path.exists(dst):
            shutil.move(src, dst)
            print(f"  Moved {pf} -> PPTs/")
        elif os.path.exists(dst):
            print(f"  {pf} already in PPTs/")
    
    # Move code files
    code_files = [
        "srmamba_t_kaggle_notebook.py",
    ]
    for cf in code_files:
        src = os.path.join(BASE, cf)
        dst = os.path.join(code_dir, cf)
        if os.path.exists(src) and not os.path.exists(dst):
            shutil.move(src, dst)
            print(f"  Moved {cf} -> Code/")
        elif os.path.exists(dst):
            print(f"  {cf} already in Code/")
    
    # Move dataset files
    dataset_files = [
        "PaviaC_Data.zip",
        "PaviaU_Data.zip",
    ]
    for df in dataset_files:
        src = os.path.join(BASE, df)
        dst = os.path.join(dataset_dir, df)
        if os.path.exists(src) and not os.path.exists(dst):
            shutil.move(src, dst)
            print(f"  Moved {df} -> Dataset/")
        elif os.path.exists(dst):
            print(f"  {df} already in Dataset/")
    
    # Clean up temp scripts
    temp_scripts = [
        "create_10_slide_presentation.py",
        "create_final_corrected_presentation.py",
        "extract_images.py",
        "extract_pdf_images.py",
        "extract_pptx.py",
        "ocr_slides.py",
        "read_pdf.py",
        "rotate_all.py",
        "rotate_slides.py",
        "friend_slides_text.txt",
    ]
    temp_dir = os.path.join(BASE, "_temp_scripts")
    os.makedirs(temp_dir, exist_ok=True)
    for ts in temp_scripts:
        src = os.path.join(BASE, ts)
        dst = os.path.join(temp_dir, ts)
        if os.path.exists(src):
            shutil.move(src, dst)
            print(f"  Moved {ts} -> _temp_scripts/")
    
    # Move slide image folders to a working folder
    work_dir = os.path.join(BASE, "_temp_scripts", "slide_images_work")
    os.makedirs(work_dir, exist_ok=True)
    
    for folder in ["slide_images", "friend_slides"]:
        src = os.path.join(BASE, folder)
        dst = os.path.join(work_dir, folder)
        if os.path.exists(src) and not os.path.exists(dst):
            shutil.move(src, dst)
            print(f"  Moved {folder}/ -> _temp_scripts/slide_images_work/")
    
    print("\n=== Final Folder Structure ===")
    for item in sorted(os.listdir(BASE)):
        full = os.path.join(BASE, item)
        if os.path.isdir(full):
            count = len(os.listdir(full))
            print(f"  [DIR]  {item}/ ({count} items)")
        else:
            size_mb = os.path.getsize(full) / (1024*1024)
            print(f"  [FILE] {item} ({size_mb:.1f} MB)")

if __name__ == "__main__":
    create_ppt()
    organize_folders()
    print("\n=== ALL DONE ===")
