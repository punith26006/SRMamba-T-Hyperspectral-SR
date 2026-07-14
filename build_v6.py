"""
V6 - Ruthwik Style Template (Beige, Dark bottom, orange accent)
Exactly 40 slides, using original high-quality diagrams (no AI gibberish).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
CLEAN = os.path.join(BASE, "_temp_scripts", "clean_diagrams")
ORIG = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")

# Ruthwik Style Colors
BG_COLOR = RGBColor(236, 236, 231)       # Beige/light grey
DARK_COLOR = RGBColor(43, 43, 43)        # Dark grey for bottom
ACCENT_COLOR = RGBColor(255, 140, 0)     # Orange/Yellow accent
TITLE_COLOR = RGBColor(40, 40, 40)       # Dark text for title
TEXT_COLOR = RGBColor(60, 60, 60)        # Dark text for body

W = Inches(13.333); H = Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]; n = 0

def S():
    global n; n += 1; return prs.slides.add_slide(BL)

def bg(s):
    # Base background
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_COLOR; sh.line.fill.background()
    
    # Dark bottom polygon (using a rectangle + triangle to simulate the slanted look)
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), W, Inches(1.0))
    bottom.fill.solid(); bottom.fill.fore_color.rgb = DARK_COLOR; bottom.line.fill.background()
    
    # Accent shape (triangle)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10), Inches(6.0), Inches(3.33), Inches(1.5))
    tri.rotation = 180
    tri.fill.solid(); tri.fill.fore_color.rgb = ACCENT_COLOR; tri.line.fill.background()
    
def hdr(s, t):
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = TITLE_COLOR

def bul(s, items, x, y, w, h, sz=16):
    """Add bullet text inside a subtle bordered box (Ruthwik style)"""
    bx = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(230, 230, 225)
    bx.line.color.rgb = RGBColor(150, 150, 150)
    bx.line.width = Pt(1.5)
    
    tb = s.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text, bold = item[0], item[1]
            p.font.bold = bold
            p.font.color.rgb = TITLE_COLOR if bold else TEXT_COLOR
        else:
            p.text = item; p.font.color.rgb = TEXT_COLOR
        p.font.size = Pt(sz); p.space_after = Pt(5)

def text_only(s, items, x, y, w, h, sz=16):
    """Text without box"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text, bold = item[0], item[1]
            p.font.bold = bold
        else:
            p.text = item
        p.font.color.rgb = TEXT_COLOR
        p.font.size = Pt(sz); p.space_after = Pt(5)

def img(s, folder, name, x, y, w, h=None):
    path = os.path.join(folder, name)
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)
    else:
        print(f"    WARN: missing {name} in {folder}")

def divider(s, num, title, sub=""):
    # Clean divider style
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = DARK_COLOR; sh.line.fill.background()
    
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.3), Inches(2.1), Inches(2.1))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT_COLOR; c.line.fill.background()
    
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.5), Inches(2.1), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER
    
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.5))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = RGBColor(255,255,255); p2.alignment = PP_ALIGN.CENTER
    
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub; p3.font.size = Pt(18); p3.font.color.rgb = BG_COLOR; p3.alignment = PP_ALIGN.CENTER

print("Building 40 slides Ruthwik Style...")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. TITLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); print(f"  {n}. Title")
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11.3), Inches(4))
sh.fill.solid(); sh.fill.fore_color.rgb = DARK_COLOR; sh.line.fill.background()
text_only(s, [("Exploration of Deep Learning Models", True)], Inches(1.5), Inches(1.8), Inches(10), Inches(1), sz=44)
text_only(s, [("for Image Super Resolution", True)], Inches(1.5), Inches(2.5), Inches(10), Inches(1), sz=36)
# Recolor title to white/orange
for p in s.shapes[-2].text_frame.paragraphs: p.font.color.rgb = RGBColor(255,255,255)
for p in s.shapes[-1].text_frame.paragraphs: p.font.color.rgb = ACCENT_COLOR

text_only(s, ["B.Tech Project Presentation", "", "Presented by: Vujja Punith Sai", "Roll No: S20240010260", "Instructor: Arun PV"],
    Inches(1.5), Inches(3.4), Inches(10), Inches(2), sz=18)
for p in s.shapes[-1].text_frame.paragraphs: p.font.color.rgb = BG_COLOR

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. CONTENTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Contents"); print(f"  {n}. Contents")
bul(s, [("Part I — Literature Review", True), "",
    "  1. What is SR & Need for SR", "  2. Problem Statement",
    "  3. Traditional SR Methods", "  4. Deep Learning & Evolution",
    "  5. Multimodal Learning", "  6. CNN: SRCNN & VDSR",
    "  7. GAN: SRGAN & ESRGAN", "  8. Transformer: SwinIR",
    "  9. State Space: SR-Mamba", " 10. Evaluation Metrics"],
    Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.2), sz=15)
bul(s, [("Part II — Our Method: SRMamba-T", True), "",
    " 11. Why Hybrid? (CNN vs Trans vs Mamba)", " 12. Complete Architecture",
    " 13. Stage-Wise Feature Flow", " 14. SSM Foundations",
    " 15. Mamba Layer & Asymmetric Mixer", " 16. MDSSM: 2D Spatial Context",
    " 17. Transformer & FFM", " 18. Training Pipeline & Loss",
    " 19. Results & LAM Analysis", " 20. Conclusion & References"],
    Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.2), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. SECTION DIVIDER I
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); divider(s, "I", "Literature Review", "Background & Related Work"); print(f"  {n}. Divider I")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4. WHAT IS SUPER RESOLUTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "What is Super Resolution?"); print(f"  {n}. What is SR")
bul(s, [
    ("* Definition:", True),
    "SR uses Deep Learning to reconstruct high-resolution images from low-resolution inputs.", "",
    ("* Main Objectives:", True),
    "  1. Improve Image Clarity — sharpness & visibility",
    "  2. Recover Image Details — missing fine structures",
    "  3. Enhance Texture & Edges — realistic edges",
    "  4. Improve Perceptual Quality — natural images",
], Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)
bul(s, [
    ("* Basic SR Workflow:", True), "",
    "   Low Resolution Image", "           ↓",
    "   Feature Extraction", "           ↓",
    "   Super Resolution Model", "           ↓",
    "   Image Reconstruction Layer", "           ↓",
    "   High Resolution Image",
], Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0), sz=16)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5. NEED FOR SR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Need for Super Resolution"); print(f"  {n}. Need for SR")
bul(s, [
    ("* Problems in LR Images:", True),
    "  •  Blur and loss of sharpness",
    "  •  Pixelation during zooming",
    "  •  Loss of fine textures and edges",
    "  •  Poor object visibility",
    "  •  Noise and distortion",
    "  •  Information loss during compression",
], Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)
bul(s, [
    ("* Real-World Applications:", True),
    "  • Satellite Imaging: terrain analysis, urban planning",
    "  • Surveillance: face recognition, number plates",
    "  • Medical Imaging: MRI/CT scan enhancement",
    "  • Mobile Camera: AI zoom, night photography",
    "  • Video Upscaling: low-quality → HD/4K conversion",
    "  • Remote Sensing: land use, disaster monitoring",
], Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6. PROBLEM STATEMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Problem Statement"); print(f"  {n}. Problem Statement")
bul(s, ["Most open-source data have coarse spatial resolution, and it is advantageous to use different",
    "sensors for effective enhancement; however, the sensors have differences in spatial resolution,",
    "spectral characteristics, noise levels, and often lack spatial or temporal overlap."],
    Inches(0.4), Inches(1.2), Inches(12.5), Inches(1.6), sz=18)
bul(s, [("* Key Challenges in SISR:", True),
    "  •  Recovering lost high-frequency details from degraded inputs",
    "  •  CNN models: limited receptive fields (local only)",
    "  •  Transformer models: O(n²) computational complexity — GPU memory explosion",
    "  •  Pure Mamba models: 1D sequential scan destroys 2D spatial context",
    "  •  Need for a hybrid approach combining efficiency + global awareness",
], Inches(0.4), Inches(3.2), Inches(12.5), Inches(3.0), sz=17)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 7. TRADITIONAL METHODS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Traditional Super Resolution Methods"); print(f"  {n}. Traditional")
bul(s, [
    ("* Interpolation Methods:", True),
    "  1) Nearest Neighbour: Copies nearest pixel. Fast but blocky.",
    "  2) Bilinear Interpolation: Uses neighboring pixels. Smoother but blurry.",
    "  3) Bicubic Interpolation: Uses 16 surrounding pixels. Better edge continuity.", "",
    "→ All traditional methods fail to reconstruct truly missing high-frequency details.",
    "→ This necessitates the use of Deep Learning."
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0), sz=18)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 8. DEEP LEARNING IN SR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Deep Learning in Super Resolution"); print(f"  {n}. DL in SR")
bul(s, [
    ("* Why Deep Learning?", True),
    "Instead of mathematically scaling pixels, DL models directly learn",
    "the mapping from LR to HR images from large datasets.", "",
    ("DL Models Learn:", True),
    "  •  Texture Patterns (roads, vegetation, faces)",
    "  •  Edge Information (boundaries, transitions)",
    "  •  Image Structures (spatial relationships)",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=15)
img(s, CLEAN, "evolution.png", Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.0))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 9. MULTIMODAL LEARNING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Relation to Multimodal Learning"); print(f"  {n}. Multimodal")
bul(s, [
    ("* Single-Modality SR:", True),
    "  • Uses a single low-resolution image (visual only).",
    "  • Reconstruction depends entirely on visual features.", "",
    ("* Multimodal Learning:", True),
    "  • Combines multiple data sources/modalities.",
    "  • Example: Satellite Image + Elevation Data → HR Output.",
    "  • Advantages: Better context, richer features.",
    "  • Limitations: Data alignment challenges, complexity.",
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0), sz=17)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 10. SRCNN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "CNN-based Models: SRCNN (2014)"); print(f"  {n}. SRCNN")
bul(s, [
    ("* SRCNN — First DL model for SR", True),
    "•  End-to-end mapping from LR to HR.", "",
    ("* 3 Convolution Layers:", True),
    "  1. Feature Extraction: Extracts low-level features.",
    "  2. Non-Linear Mapping: Maps LR to HR features.",
    "  3. Reconstruction: Reconstructs final HR image.", "",
    ("* Loss Function: MSE", True),
    "  L(Θ) = (1/n) Σ ||F(Yᵢ; Θ) - Xᵢ||²",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=14)
img(s, CLEAN, "srcnn.png", Inches(6.5), Inches(1.2), Inches(6.5), Inches(4.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 11. FSRCNN + VDSR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "CNN Models: FSRCNN (2016) & VDSR (2016)"); print(f"  {n}. FSRCNN+VDSR")
bul(s, [
    ("* FSRCNN — Fast Super Resolution CNN", True),
    "•  Operates directly on LR space (faster).",
    "•  Uses deconvolution layer for upsampling at end.",
    "•  Adv: Smaller model, fast.",
    "•  Lim: Still limited receptive field.",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=15)
bul(s, [
    ("* VDSR — Very Deep SR (20 Layers)", True),
    "•  Key: Global Residual Learning (HR = LR + Residual).",
    "•  Larger receptive field (41×41 pixels).",
    "•  Multi-scale: handles ×2, ×3, ×4.",
    "•  Lim: Operates in HR space (slow).",
], Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 12. SRGAN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "GAN-based Models: SRGAN (2017)"); print(f"  {n}. SRGAN")
bul(s, [
    ("* Why GANs for SR?", True),
    "CNN models optimized for MSE produce blurry outputs.",
    "SRGAN introduced Adversarial Learning for realism.", "",
    ("* Architecture:", True),
    "  • Generator: Deep ResNet → produces HR from LR.",
    "  • Discriminator: VGG-style → evaluates real vs fake.", "",
    ("* Key Innovation: Perceptual Loss", True),
    "  Uses VGG feature maps → Sharper textures.",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=14)
img(s, CLEAN, "srgan.png", Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.0))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 13. ESRGAN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "GAN-based Models: ESRGAN (2018)"); print(f"  {n}. ESRGAN")
bul(s, [
    ("* ESRGAN — Enhanced SRGAN", True),
    "•  Residual-in-Residual Dense Blocks (RRDB).",
    "•  Removed batch normalization for stable training.",
    "•  Relativistic Discriminator: predicts relative quality.", "",
    ("* Improvements over SRGAN:", True),
    "•  Multi-level residual connections.",
    "•  Better texture, higher perceptual quality.",
    "•  More stable training, sharper outputs.", "",
    ("* Limitations of GANs:", True),
    "•  Training instability, artifact hallucination.",
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0), sz=16)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 14. SWINIR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Transformer Models: SwinIR (2021)"); print(f"  {n}. SwinIR")
bul(s, [
    ("* SwinIR — Swin Transformer for SR", True),
    "•  Built on Vision Transformers (ViT).",
    "•  Uses Shifted Window Attention.",
    "•  Self-attention determines which image regions",
    "   are most important for reconstruction.", "",
    ("* Advantages:", True), "  Global feature learning, excellent quality.", "",
    ("* Limitations:", True), "  O(n²) complexity, high memory cost.",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=14)
img(s, CLEAN, "swinir.png", Inches(6.5), Inches(1.2), Inches(6.5), Inches(4.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 15. SR-MAMBA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "State Space Models: SR-Mamba (2024)"); print(f"  {n}. SR-Mamba")
bul(s, [
    ("* State Space Models (SSMs):", True),
    "•  Sequence modeling with O(n) linear complexity.",
    "•  Avoids expensive global attention.", "",
    ("* SR-Mamba:", True),
    "•  Modern SR architecture based on Mamba/S6.",
    "•  State Transition Mechanisms instead of attention.", "",
    ("* Limitations:", True),
    "•  1D sequential scan destroys 2D spatial context.",
    "  → Need for hybrid approach: SRMamba-T!",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.0), sz=14)
img(s, CLEAN, "srmamba.png", Inches(6.5), Inches(1.2), Inches(6.5), Inches(4.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 16. EVALUATION METRICS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Evaluation Metrics"); print(f"  {n}. Metrics")
bul(s, [
    ("* 1. MSE (Mean Square Error)", True),
    "  Pixel-wise reconstruction error. Lower is better.",
    "  Used in SRCNN, VDSR. Produces blurry results.", "",
    ("* 2. PSNR (Peak Signal-to-Noise Ratio)", True),
    "  Reconstruction quality in dB. Higher is better.",
    "  Inversely related to MSE.", "",
    ("* 3. SSIM (Structural Similarity Index)", True),
    "  Measures structural & perceptual quality (luminance, contrast).",
    "  SSIM → 1 means excellent structural similarity.",
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0), sz=17)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 17. SECTION DIVIDER II
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for SISR"); print(f"  {n}. Divider II")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PART II: SRMamba-T 
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def my_slide(title, text_items, img_name):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:40]}")
    bul(s, text_items, Inches(0.4), Inches(1.2), Inches(5.5), Inches(5.1), sz=14)
    # The original user images are wide, so we place them on the right
    img(s, ORIG, img_name, Inches(6.2), Inches(1.3), Inches(6.9))

# 18
my_slide("Why Hybrid? — Architecture Compromise", [
    ("* No single architecture is perfect:", True), "",
    ("CNN:", True), "  ✗ Limited receptive field (local only)", "",
    ("Transformer:", True), "  ✓ Global context via self-attention",
    "  ✗ O(n²) complexity — GPU memory explosion", "",
    ("Mamba (SSM):", True), "  ✓ O(n) linear complexity",
    "  ✗ 1D scan destroys 2D spatial context", "",
    ("→ SRMamba-T combines Mamba's efficiency", True),
    ("   with Transformer's global awareness.", True),
], "slide_04.png")

# 19
my_slide("CNN vs Transformer vs Mamba", [
    ("* Effective Receptive Field (ERF):", True), "",
    "• CNN: Small, localized ERF.",
    "• Transformer: Large but uniform ERF (wasteful).",
    "• Mamba: Directional ERF along scan line.", "",
    ("* The Hybrid Insight:", True),
    "Phase 1 (Mamba): Dense feature extraction at O(n) cost.",
    "Phase 2 (Transformer): Global refinement on compressed",
    "features without memory explosion.",
], "slide_06.png")

# 20
s = S(); bg(s); hdr(s, "The Complete SRMamba-T Architecture"); print(f"  {n}. Complete Arch")
bul(s, [
    "U-Net encoder with asymmetric layer placement: Mamba layers (early) for efficient dense extraction, Transformer layers (late) for global refinement."
], Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.8), sz=15)
img(s, ORIG, "slide_09.png", Inches(0.5), Inches(2.1), Inches(12.3))

# 21
s = S(); bg(s); hdr(s, "Stage-Wise Feature Flow"); print(f"  {n}. Flow")
bul(s, [
    "Input → Shallow Ext. → Stage 1 (Mamba) → Stage 2 (Mixed) → Stage 3 (Trans) → Reconstruction → HR."
], Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.8), sz=15)
img(s, ORIG, "slide_10.png", Inches(0.5), Inches(2.1), Inches(12.3))

# 22
my_slide("SSM Foundations & Discretization", [
    ("* State Space Models (SSMs):", True),
    "Continuous-time linear system:",
    "  h'(t) = A·h(t) + B·x(t)",
    "  y(t)  = C·h(t) + D·x(t)", "",
    ("* Discretization (ZOH):", True),
    "For deep learning, discretized to:",
    "  Ā = exp(Δ·A)",
    "  B̄ = (Δ·A)⁻¹(exp(Δ·A) - I)·Δ·B", "",
    ("* Selective Scan:", True),
    "Makes matrices input-dependent for content-awareness.",
], "slide_11.png")

# 23
my_slide("SSM: Equivalence to Convolution", [
    ("* Key Insight:", True),
    "The discretized SSM can be computed as a global",
    "convolution, enabling highly parallel training.", "",
    ("* Convolutional Form:", True),
    "  K = (CB̄, CĀB̄, CĀ²B̄, ..., CĀᴸB̄)",
    "  y = x * K", "",
    ("* Benefit:", True),
    "Parallel convolution during training (fast).",
    "Sequential updates during inference (memory efficient).",
], "slide_13.png")

# 24
my_slide("Mamba Layer: Anatomy & Mathematics", [
    ("* Mamba Block Structure:", True),
    "Input → LayerNorm → Linear Proj →",
    "  Branch 1: Conv1D → SiLU → SSM",
    "  Branch 2: SiLU activation",
    "→ Element-wise multiply → Linear → Output", "",
    ("* Key Equations:", True),
    "  z = σ(Linear(Norm(x)))",
    "  h = SSM(SiLU(Conv1D(Linear(Norm(x)))))",
    "  y = Linear(h ⊙ z) + x",
], "slide_14.png")

# 25
my_slide("The Asymmetric Mamba Mixer Block", [
    ("* Why Asymmetric?", True),
    "Standard Mamba scans in one direction only.", "",
    ("* The Mixer combines:", True),
    "  •  Forward scan branch",
    "  •  Backward scan branch",
    "  •  Different parameter counts per direction", "",
    "This captures bidirectional dependencies while",
    "maintaining computational efficiency.",
], "slide_15.png")

# 26
my_slide("The Spatial Dependency Problem", [
    ("* Why 1D scanning fails for images:", True),
    "Images are 2D grids — every pixel has spatial",
    "neighbors in all directions (up, down, left, right).", "",
    ("* 1D Flattening Problem:", True),
    "When a 2D image is flattened to 1D sequence:",
    "  •  Vertical neighbors become distant in 1D.",
    "  •  Spatial proximity information is lost.",
    "  •  SSM 'forgets' vertical relationships.",
], "slide_16.png")

# 27
my_slide("Multi-Directional Selective Scan (MDSSM)", [
    ("* MDSSM Solution:", True),
    "1) Channel Split: Split features into 4 groups.",
    "2) 4-Direction Scanning:",
    "   → Top-to-Bottom",
    "   → Bottom-to-Top",
    "   → Left-to-Right",
    "   → Right-to-Left",
    "3) Concatenate: Merge scanned features.", "",
    ("* Result:", True),
    "Preserves 2D spatial context at O(n) cost!",
], "slide_17.png")

# 28
my_slide("MDSSM: Preserving 2D Spatial Context", [
    ("* Channel Split Strategy:", True),
    "Instead of running 4 full SSMs (4× cost), MDSSM",
    "splits channels into 4 groups. Each group runs",
    "one direction → only 1× cost!", "",
    ("* Mathematical Formulation:", True),
    "  C₁,C₂,C₃,C₄ = Split(X, dim=channel)",
    "  Cᵢ' = SSM(Scan_dir_i(Cᵢ))",
    "  Output = Concat(C₁',C₂',C₃',C₄')", "",
    "Parameter overhead is ~0%.",
], "slide_18.png")

# 29
my_slide("Inside the Transformer Layer", [
    ("* Transformer Layer Structure:", True),
    "Input → LayerNorm → Multi-Head Self-Attention",
    "→ Residual → LayerNorm → FFN → Residual → Out", "",
    ("* Self-Attention:", True),
    "  Attention(Q,K,V) = softmax(QKᵀ/√d) · V", "",
    ("* In SRMamba-T:", True),
    "Placed in later stages where features are already",
    "compressed, making the O(n²) cost manageable.",
], "slide_19.png")

# 30
my_slide("Un-Windowed Global Self-Attention", [
    ("* Problem with Windowed Attention:", True),
    "SwinIR uses shifted windows → limits receptive field.", "",
    ("* SRMamba-T Solution:", True),
    "Un-windowed attention: every pixel attends to",
    "every other pixel in the feature map.", "",
    ("* Benefit:", True),
    "True global context without window artifacts.",
    "Provides better texture reconstruction.",
], "slide_20.png")

# 31
my_slide("Feature Fusion Module (FFM)", [
    ("* Purpose:", True),
    "Combines features from Mamba and Transformer.", "",
    ("* FFM Process:", True),
    "  1) Receive Mamba features (local/sequential).",
    "  2) Receive Transformer features (global/spatial).",
    "  3) Concatenate along channel dimension.",
    "  4) Apply 1×1 conv for channel reduction.",
    "  5) Add residual connection.", "",
    "Learns optimal fusion weights.",
], "slide_21.png")

# 32
my_slide("Adaptive Reconstruction & Loss", [
    ("* Reconstruction Pipeline:", True),
    "Features → Conv layers → PixelShuffle → HR image.", "",
    ("* Dual-Domain Loss Function:", True),
    "  L_total = L_pixel + λ · L_freq", "",
    "L_pixel: Standard L1/MSE ensures pixel accuracy.",
    "L_freq: FFT-based frequency domain loss",
    "  L_freq = ||FFT(HR) - FFT(SR)||₁", "",
    "Preserves high-frequency details & sharp edges.",
], "slide_22.png")

# 33
my_slide("Training Pipeline & Setup", [
    ("* Training Configuration:", True),
    "•  Dataset: DIV2K (800 training images)",
    "•  Patch size: 64×64 LR patches",
    "•  Batch size: 16",
    "•  Optimizer: Adam (β₁=0.9, β₂=0.99)",
    "•  Learning rate: 2×10⁻⁴ with cosine decay",
    "•  Total iterations: 500K", "",
    ("* Degradation Model:", True),
    "•  Bicubic downsampling (×2, ×3, ×4)",
], "slide_23.png")

# 34
my_slide("Quantitative Results: PSNR vs Compute", [
    ("* Key Results:", True),
    "SRMamba-T achieves:",
    "  •  +0.28 dB PSNR improvement over MambaIR",
    "  •  38.7% fewer MACs (computational cost)",
    "  •  Better SSIM on all benchmark datasets", "",
    ("* Benchmarks tested:", True),
    "  Set5, Set14, BSD100, Urban100, Manga109", "",
    "Outperforms SRCNN, VDSR, SwinIR, and MambaIR",
    "with significantly less computational cost.",
], "slide_24.png")

# 35
my_slide("Qualitative Validation", [
    ("* Visual comparison on Urban100 dataset:", True),
    "•  Complex textures (buildings, fences, grids)",
    "•  High-frequency detail regions", "",
    ("* Observations:", True),
    "•  CNNs: Blurry, loss of detail.",
    "•  GANs: Sharp but artifacts.",
    "•  SwinIR: Misses some textures.",
    "•  MambaIR: Directional artifacts.", "",
    ("* SRMamba-T: Sharpest, no directional artifacts.", True),
], "slide_25.png")

# 36
my_slide("LAM Analysis: Inside the Model's Mind", [
    ("* LAM = Local Attribution Map", True),
    "Visualizes which input pixels contributed most",
    "to each output pixel.", "",
    ("* Diffusion Index (DI):", True),
    "  •  CNN models: Low DI (very local)",
    "  •  SwinIR: Medium DI (window-limited)",
    "  •  MambaIR: Higher DI but directional",
    "  •  SRMamba-T: Highest DI (truly global)", "",
    "SRMamba-T uses information from the widest",
    "range of input pixels.",
], "slide_26.png")

# 37
my_slide("Hardware Efficiency & Ablation Study", [
    ("* Memory Scaling:", True),
    "Transformer has O(n²) memory growth.",
    "SRMamba-T achieves near-linear scaling.",
    "At 1024×1024, SwinIR runs OOM, SRMamba-T runs fine.", "",
    ("* Ablation: Layer Ordering", True),
    "  A) All Mamba → Poor global context",
    "  B) All Transformer → Too expensive",
    "  C) Trans first, Mamba later → Suboptimal",
    "  D) Mamba first, Trans later → BEST",
], "slide_27.png")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 38. CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Conclusion & Future Work"); print(f"  {n}. Conclusion")
bul(s, [
    ("* Key Contributions:", True),
    "✓ First hybrid Mamba-Transformer for image SR.",
    "✓ MDSSM: 4-directional scanning preserves 2D context at O(n) cost.",
    "✓ Asymmetric placement: Mamba (early) + Transformer (late) = optimal tradeoff.",
    "✓ +0.28 dB PSNR, 38.7% fewer MACs than state-of-the-art.",
    "✓ Dual-domain loss prevents high-frequency artifacts.",
    "✓ Linear memory scaling allows super-resolving large images.",
], Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)
bul(s, [
    ("* Future Work:", True),
    "• Extend to Hyperspectral SR (PaviaC, PaviaU).",
    "• Scale to SRMamba-T-L for classical SR tasks.",
    "• Investigate real-world degradation beyond bicubic.",
    "• Deploy on edge devices leveraging linear memory.",
    "• Combine with diffusion models for perceptual quality.",
], Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 39. REFERENCES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "References"); print(f"  {n}. References")
bul(s, [
    "[1] Dong et al., \"Image Super-Resolution Using Deep Convolutional Networks\" (SRCNN), ECCV 2014",
    "[2] Dong et al., \"Accelerating the Super-Resolution CNN\" (FSRCNN), ECCV 2016",
    "[3] Kim et al., \"Accurate Image SR Using Very Deep Networks\" (VDSR), CVPR 2016",
    "[4] Ledig et al., \"Photo-Realistic SISR Using a GAN\" (SRGAN), CVPR 2017",
    "[5] Wang et al., \"ESRGAN: Enhanced Super-Resolution GANs\", ECCV 2018",
    "[6] Liang et al., \"SwinIR: Image Restoration Using Swin Transformer\", ICCV 2021",
    "[7] Guo et al., \"MambaIR: A Simple Baseline for Image Restoration with SSM\", ECCV 2024",
    "[8] Gu & Dao, \"Mamba: Linear-Time Sequence Modeling with Selective State Spaces\", 2023",
    "[9] SRMamba-T: Hybrid Mamba-Transformer for Single Image Super-Resolution, 2024",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017",
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 40. THANK YOU
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); print(f"  {n}. Thank You")
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11.3), Inches(4))
sh.fill.solid(); sh.fill.fore_color.rgb = DARK_COLOR; sh.line.fill.background()
text_only(s, [("Thank You!", True)], Inches(1), Inches(2), Inches(11.3), Inches(1), sz=60)
for p in s.shapes[-1].text_frame.paragraphs:
    p.font.color.rgb = ACCENT_COLOR; p.alignment = PP_ALIGN.CENTER
text_only(s, ["Questions & Discussion"], Inches(1), Inches(3.2), Inches(11.3), Inches(1), sz=30)
for p in s.shapes[-1].text_frame.paragraphs:
    p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER
text_only(s, ["Vujja Punith Sai | S20240010260 | Guide: Arun PV"], Inches(1), Inches(4.5), Inches(11.3), Inches(1), sz=20)
for p in s.shapes[-1].text_frame.paragraphs:
    p.font.color.rgb = BG_COLOR; p.alignment = PP_ALIGN.CENTER

out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v6_RUTHWIK_STYLE.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
