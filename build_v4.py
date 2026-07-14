"""
EXACTLY 40-SLIDE Professional BTP Presentation
Every slide: Typed text + architecture diagram image (where applicable)
NO full-page screenshot pasting
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
CR = os.path.join(BASE, "_temp_scripts", "cropped_diagrams")

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
ACC2 = RGBColor(0xE8, 0x5D, 0x04)
GOLD = RGBColor(0xFF, 0xBE, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBGR = RGBColor(0xF0, 0xF4, 0xF8)
TDRK = RGBColor(0x1A, 0x1A, 0x2E)
TB = RGBColor(0x33, 0x3D, 0x4D)
GREEN = RGBColor(0x2D, 0xA0, 0x5B)
DARK = RGBColor(0x1B, 0x2A, 0x4A)

W = Inches(13.333); H = Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]; n = 0

def S():
    global n; n += 1; return prs.slides.add_slide(BL)

def bg(s, c=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def hdr(s, t):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.85))
    b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), W, Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE

def bul(s, items, x, y, w, h, sz=14):
    """Add bullet text"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text, bold = item[0], item[1]
            p.font.bold = bold
            if bold: p.font.color.rgb = TDRK
            else: p.font.color.rgb = TB
        else:
            p.text = item; p.font.color.rgb = TB
        p.font.size = Pt(sz); p.space_after = Pt(3)

def img(s, name, x, y, w, h=None):
    path = os.path.join(CR, name)
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)
    else:
        print(f"    WARN: missing {name}")

def divider(s, num, title, sub=""):
    bg(s, NAVY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.3), Inches(2.1), Inches(2.1))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.5), Inches(2.1), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.5))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub; p3.font.size = Pt(16); p3.font.color.rgb = ACCENT; p3.alignment = PP_ALIGN.CENTER

def sub_hdr(s, t, x, y, w):
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.4))
    bx.fill.solid(); bx.fill.fore_color.rgb = DARK; bx.line.fill.background()
    tf = bx.text_frame; tf.paragraphs[0].text = t
    tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

print("Building 40 slides...")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. TITLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s, NAVY); print(f"  {n}. Title")
s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H).fill.solid()
s.shapes[-1].fill.fore_color.rgb = ACCENT; s.shapes[-1].line.fill.background()
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.5), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Exploration of Deep Learning Models"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "for Image Super Resolution"; p2.font.size = Pt(44); p2.font.bold = True; p2.font.color.rgb = ACCENT
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.6), Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
bul(s, [("B.Tech Project Presentation", True), "", ("Presented by:  Vujja Punith Sai", False),
    ("Roll No:  S20240010260", False), ("Instructor:  Arun PV", False)],
    Inches(1.5), Inches(4.0), Inches(10), Inches(3), sz=20)
# Recolor details text to gray
for shape in s.shapes:
    if hasattr(shape, 'text_frame'):
        for p in shape.text_frame.paragraphs:
            try:
                if p.font.color.rgb == TB:
                    p.font.color.rgb = RGBColor(0x99,0x99,0x99)
            except:
                pass

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. CONTENTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Contents"); print(f"  {n}. Contents")
bul(s, [("Part I — Literature Review", True), "",
    "  1-2.   What is SR & Need for SR", "  3.     Problem Statement",
    "  4.     Traditional SR Methods", "  5-6.   Deep Learning & Evolution",
    "  7.     Multimodal Learning", "  8-9.   CNN: SRCNN (Arch + Details)",
    "  10.    CNN: FSRCNN & VDSR", "  11-12. GAN: SRGAN & ESRGAN",
    "  13-14. Transformer: SwinIR", "  15-16. State Space: SR-Mamba",
    "  17.    Evaluation Metrics"],
    Inches(0.4), Inches(1.1), Inches(5.8), Inches(6), sz=14)
bul(s, [("Part II — Our Method: SRMamba-T", True), "",
    "  19.    Why Hybrid?", "  20.    SRMamba-T Architecture",
    "  21.    Stage-Wise Feature Flow", "  22.    SSM Foundations",
    "  23-24. Mamba Layer & MDSSM", "  25.    MDSSM: 2D Spatial Context",
    "  26-27. Transformer & Attention", "  28.    Feature Fusion Module",
    "  29.    Reconstruction & Loss", "  30.    Training Pipeline",
    "  31-32. Quantitative & Qualitative Results",
    "  33-34. LAM & Hardware Analysis", "  35.    Ablation Study",
    "  36-37. Conclusion & References"],
    Inches(6.8), Inches(1.1), Inches(5.8), Inches(6), sz=14)
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.1), Inches(0.03), Inches(6))
dv.fill.solid(); dv.fill.fore_color.rgb = ACCENT; dv.line.fill.background()

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. SECTION DIVIDER PART I
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); divider(s, "I", "Literature Review", "Background & Related Work"); print(f"  {n}. Divider I")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4. WHAT IS SUPER RESOLUTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "What is Super Resolution?"); print(f"  {n}. What is SR")
sub_hdr(s, "Definition & Objectives", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "SR uses Deep Learning to reconstruct high-resolution",
    "images from low-resolution inputs.",
    "", ("It aims to recover:", True),
    "  •  Lost image details & fine textures",
    "  •  Sharp edges & structural information",
    "  •  Overall visual quality",
    "", ("Main Objectives:", True),
    "  1. Improve Image Clarity — sharpness & visibility",
    "  2. Recover Image Details — missing fine structures",
    "  3. Enhance Texture & Edges — realistic edges",
    "  4. Improve Perceptual Quality — natural images",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(5.5), sz=14)
sub_hdr(s, "Basic SR Workflow", Inches(6.6), Inches(1.1), Inches(6.3))
bul(s, [
    "", "   Low Resolution Image", "           ↓",
    "   Feature Extraction", "           ↓",
    "   Super Resolution Model", "           ↓",
    "   Image Reconstruction Layer", "           ↓",
    "   High Resolution Image",
], Inches(6.6), Inches(1.7), Inches(6.3), Inches(5.5), sz=16)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5. NEED FOR SR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Need for Super Resolution"); print(f"  {n}. Need for SR")
sub_hdr(s, "Problems in Low Resolution Images", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "Low-resolution images suffer from several quality limitations:",
    "", "  •  Blur and loss of sharpness",
    "  •  Pixelation during zooming",
    "  •  Loss of fine textures and edges",
    "  •  Poor object visibility",
    "  •  Reduced structural details",
    "  •  Noise and distortion",
    "  •  Information loss during compression",
    "  •  Degraded quality in low-light conditions",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(5.5), sz=14)
sub_hdr(s, "Real-World Applications", Inches(6.6), Inches(1.1), Inches(6.3))
bul(s, [
    ("* Satellite Imaging:", True), "    terrain analysis, urban planning, monitoring",
    ("* Surveillance:", True), "    face recognition, number plate enhancement",
    ("* Medical Imaging:", True), "    MRI/CT scan, disease detection",
    ("* Mobile Camera:", True), "    AI zoom, night photography, sharpening",
    ("* Video Upscaling:", True), "    low-quality → HD/4K conversion",
    ("* Remote Sensing:", True), "    land use, disaster monitoring, geology",
    "", ("Industry: Samsung, Google, Apple use SR", True),
], Inches(6.6), Inches(1.7), Inches(6.3), Inches(5.5), sz=14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6. PROBLEM STATEMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s, LBGR); hdr(s, "Problem Statement"); print(f"  {n}. Problem Statement")
bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(2.2))
bx.fill.solid(); bx.fill.fore_color.rgb = WHITE; bx.line.color.rgb = ACCENT; bx.line.width = Pt(2)
bul(s, ["Most open-source data have coarse spatial resolution, and it is advantageous to use different",
    "sensors for effective enhancement; however, the sensors have differences in spatial resolution,",
    "spectral characteristics, noise levels, and often lack spatial or temporal overlap."],
    Inches(1.2), Inches(1.6), Inches(10.9), Inches(1.6), sz=19)
bul(s, [("Key Challenges in SISR:", True), "",
    "  •  Recovering lost high-frequency details from degraded inputs",
    "  •  Balancing reconstruction quality vs. computational cost",
    "  •  CNN models: limited receptive fields (local only)",
    "  •  Transformer models: O(n²) computational complexity — GPU memory explosion",
    "  •  Pure Mamba models: 1D sequential scan destroys 2D spatial context",
    "  •  Need for a hybrid approach combining efficiency + global awareness",
], Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.4), sz=16)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 7. TRADITIONAL METHODS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Traditional Super Resolution Methods"); print(f"  {n}. Traditional")
bul(s, ["Before Deep Learning, SR used interpolation-based image scaling methods."], Inches(0.4), Inches(1.0), Inches(12), Inches(0.4), sz=14)
for i, (title, pts, clr) in enumerate([
    ("1) Nearest Neighbour", ["Copies nearest pixel value", "Computationally very fast", "Adv: Low cost, fast processing", "Lim: Blocky, jagged edges, poor quality"], ACCENT),
    ("2) Bilinear Interpolation", ["Uses neighbouring pixels", "Smoother than nearest neighbour", "Adv: Better visual quality", "Lim: Loss of sharpness, blurred edges"], ACC2),
    ("3) Bicubic Interpolation", ["Uses surrounding 16 pixels", "Better edge continuity", "Adv: Widely used, improved smoothness", "Lim: Cannot reconstruct missing details"], GREEN)]):
    x = Inches(0.3 + i * 4.2)
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(4.0), Inches(4.8))
    bx.fill.solid(); bx.fill.fore_color.rgb = LBGR; bx.line.color.rgb = clr; bx.line.width = Pt(1.5)
    hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(4.0), Inches(0.45))
    hd.fill.solid(); hd.fill.fore_color.rgb = clr; hd.line.fill.background()
    bul(s, [(title, True)], x+Inches(0.15), Inches(1.75), Inches(3.7), Inches(0.4), sz=14)
    # Recolor title white
    for sh in s.shapes:
        if hasattr(sh, 'text_frame'):
            for p in sh.text_frame.paragraphs:
                if p.text == title: p.font.color.rgb = WHITE
    bul(s, [f"•  {p}" for p in pts], x+Inches(0.15), Inches(2.4), Inches(3.7), Inches(3.8), sz=13)
bul(s, [("→ All interpolation methods fail to reconstruct truly missing high-frequency details → Need for Deep Learning", True)],
    Inches(0.3), Inches(6.7), Inches(12.5), Inches(0.5), sz=13)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 8. DEEP LEARNING IN SR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Deep Learning in Super Resolution"); print(f"  {n}. DL in SR")
sub_hdr(s, "Why Deep Learning?", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "Traditional methods cannot recover missing information.",
    "DL models learn LR→HR mapping from large datasets.", "",
    ("Deep Learning Models Learn:", True),
    "  •  Texture Patterns — roads, vegetation, faces",
    "  •  Edge Information — boundaries, transitions",
    "  •  Image Structures — spatial relationships",
    "  •  Feature Representation — hierarchical features", "",
    ("Main Idea:", True),
    "Instead of enlarging pixels mathematically, DL models",
    "learn the mapping from LR to HR images directly.",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(5.5), sz=14)
sub_hdr(s, "Evolution of SR Architectures", Inches(6.6), Inches(1.1), Inches(6.3))
img(s, "gen_timeline.png", Inches(6.6), Inches(1.7), Inches(6.3), Inches(5.3))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 9. MULTIMODAL LEARNING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Relation to Multimodal Learning"); print(f"  {n}. Multimodal")
sub_hdr(s, "Single-Modality SR", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "Traditional SR systems mainly use:",
    "  •  A single low-resolution image",
    "  •  One data modality (visual only)", "",
    "This is known as Single-Modality Super Resolution.", "",
    "In this approach:",
    "  •  Only visual image information is used",
    "  •  Reconstruction depends entirely on image features",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(5.5), sz=15)
sub_hdr(s, "Multimodal Learning", Inches(6.6), Inches(1.1), Inches(6.3))
bul(s, [
    ("What is Multimodal Learning?", True),
    "Combines multiple data sources & different modalities",
    "to improve learning and prediction quality.", "",
    ("Example:", True),
    "(Satellite Image + Elevation Data) → Model →",
    "Enhanced High Resolution Output", "",
    ("Advantages:", True), "Better context, richer features",
    ("Limitations:", True), "Data alignment challenges, higher complexity",
], Inches(6.6), Inches(1.7), Inches(6.3), Inches(5.5), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 10. SRCNN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "CNN-based Models: SRCNN (2014)"); print(f"  {n}. SRCNN")
bul(s, [
    ("SRCNN — First DL model for Super Resolution", True), "",
    "•  End-to-end learning for SR using CNNs",
    "•  Directly learns mapping LR → HR", "",
    ("Architecture: 3 Convolution Layers", True),
    "  1) Feature Extraction: Extracts low-level features",
    "     (edges, corners, textures) from LR input",
    "  2) Non-Linear Mapping: Maps LR features into",
    "     HR feature representations",
    "  3) Reconstruction: Reconstructs final HR image",
    "     from learned feature maps", "",
    ("Loss Function: MSE", True),
    "  L(Θ) = (1/n) Σ ||F(Yᵢ; Θ) - Xᵢ||²",
], Inches(0.4), Inches(1.1), Inches(5.8), Inches(6), sz=14)
img(s, "gen_srcnn.png", Inches(6.3), Inches(1.2), Inches(6.8), Inches(5.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 11. SRCNN DETAILS + FSRCNN + VDSR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "CNN Models: FSRCNN (2016) & VDSR (2016)"); print(f"  {n}. FSRCNN+VDSR")
sub_hdr(s, "FSRCNN — Fast Super Resolution CNN", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "•  Faster version of SRCNN (5× speedup)",
    "•  Operates on LR space directly (no pre-upsampling)",
    "•  Uses deconvolution layer for upsampling at end",
    "•  5 stages: Feature Extraction → Shrinking →",
    "   Mapping → Expanding → Deconvolution",
    "•  Adv: Smaller model, adaptable to multiple scales",
    "•  Lim: Still limited receptive field",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(4.5), sz=14)
sub_hdr(s, "VDSR — Very Deep SR (20 Layers)", Inches(6.6), Inches(1.1), Inches(6.3))
bul(s, [
    "•  20 convolution layers (3×3 filters, 64 channels)",
    "•  Key: Global Residual Learning",
    "   HR = LR + Learned Residual",
    "•  No pooling layers — preserves spatial resolution",
    "•  Multi-scale: single model handles ×2, ×3, ×4",
    "•  Larger receptive field (41×41 pixels)",
    "•  Lim: Operates in HR space (slow), no skip connections",
], Inches(6.6), Inches(1.7), Inches(6.3), Inches(4.5), sz=14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 12. SRGAN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "GAN-based Models: SRGAN (2017)"); print(f"  {n}. SRGAN")
bul(s, [
    ("Why GAN-based SR was introduced:", True),
    "CNN models (SRCNN, VDSR) optimized for MSE → blurry outputs",
    "SRGAN introduced Adversarial Learning for perceptual realism", "",
    ("Architecture:", True),
    "  Generator: Deep ResNet → produces HR from LR",
    "    LR → Conv → Residual Blocks → Skip → Upsample → HR",
    "  Discriminator: VGG-style → evaluates real vs fake", "",
    ("Key Innovation: Perceptual Loss", True),
    "  Uses VGG feature maps instead of pixel MSE",
    "  → Sharper textures, visually realistic results", "",
    ("Lim: Training instability, artifact generation", False),
], Inches(0.4), Inches(1.1), Inches(5.8), Inches(6), sz=14)
img(s, "gen_srgan.png", Inches(6.3), Inches(1.2), Inches(6.8), Inches(5.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 13. ESRGAN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "GAN-based Models: ESRGAN (2018)"); print(f"  {n}. ESRGAN")
bul(s, [
    ("ESRGAN — Enhanced SRGAN", True), "",
    "•  Residual-in-Residual Dense Blocks (RRDB)",
    "•  Removed batch normalization for stable training",
    "•  Relativistic Discriminator: predicts relative quality",
    "   instead of absolute real/fake", "",
    ("Architecture Improvements:", True),
    "•  RRDB: Multi-level residual connections",
    "•  Dense connections within each residual block",
    "•  Network interpolation between PSNR & GAN models", "",
    ("Advantages over SRGAN:", True),
    "•  Better texture, higher perceptual quality",
    "•  More stable training, sharper outputs",
    "•  Won PIRM2018 SR Challenge", "",
    ("Limitations of GAN-based:", True),
    "•  Training instability, mode collapse, hallucinated textures",
], Inches(0.4), Inches(1.1), Inches(12.5), Inches(6), sz=14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 14. SWINIR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Transformer-based Models: SwinIR"); print(f"  {n}. SwinIR")
bul(s, [
    ("SwinIR — Swin Transformer for Image Restoration", True), "",
    "•  Built on Vision Transformers (ViT)",
    "•  Uses Shifted Window Attention instead of full global",
    "•  Self-attention determines:",
    "   - Which image regions are important",
    "   - How distant features are related",
    "   - How global structures should be reconstructed", "",
    ("Architecture:", True),
    "  Shallow Feature Extraction → Deep Feature (RSTB) →",
    "  HR Reconstruction", "",
    ("Adv:", True), "  Global feature learning, excellent quality",
    ("Lim:", True), "  O(n²) complexity, high GPU memory, slow training",
], Inches(0.4), Inches(1.1), Inches(6.0), Inches(6), sz=14)
img(s, "gen_swinir.png", Inches(6.3), Inches(1.3), Inches(6.8), Inches(5.3))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 15. SR-MAMBA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "State Space Models: SR-Mamba"); print(f"  {n}. SR-Mamba")
bul(s, [
    ("State Space Models (SSMs):", True),
    "•  Sequence modeling with O(n) linear complexity",
    "•  Efficient long-range information propagation",
    "•  Lower memory usage than Transformers",
    "•  Avoids expensive global attention operations", "",
    ("SR-Mamba — Super Resolution Mamba:", True),
    "•  Modern SR architecture based on Mamba/S6",
    "•  State Transition Mechanisms instead of self-attention",
    "•  Updates hidden states sequentially",
    "•  Preserves long-range information", "",
    ("Limitations:", True),
    "•  1D sequential scan destroys 2D spatial context",
    "•  Context loss at pixel boundaries",
    "  → Need for hybrid approach: SRMamba-T!",
], Inches(0.4), Inches(1.1), Inches(6.0), Inches(6), sz=14)
img(s, "gen_ssm.png", Inches(6.3), Inches(1.3), Inches(6.8), Inches(5.3))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 16. EVALUATION METRICS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Evaluation Metrics: MSE, PSNR, SSIM"); print(f"  {n}. Metrics")
for i, (title, pts, clr) in enumerate([
    ("MSE (Mean Square Error)", ["Pixel-wise reconstruction error", "MSE = (1/mn) Σ (I(i,j) - K(i,j))²",
        "Lower MSE = better reconstruction", "Used in SRCNN, FSRCNN, VDSR"], ACCENT),
    ("PSNR (Peak Signal-to-Noise)", ["Reconstruction quality in dB", "PSNR = 10·log₁₀(MAX²/MSE)",
        "Higher PSNR = better fidelity", "Inversely related to MSE"], ACC2),
    ("SSIM (Structural Similarity)", ["Structural & perceptual quality", "Compares luminance, contrast, structure",
        "SSIM → 1: good quality", "SSIM → 0: poor quality"], GREEN)]):
    x = Inches(0.3 + i * 4.2)
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(4.0), Inches(4.5))
    bx.fill.solid(); bx.fill.fore_color.rgb = LBGR; bx.line.color.rgb = clr; bx.line.width = Pt(1.5)
    hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.2), Inches(4.0), Inches(0.45))
    hd.fill.solid(); hd.fill.fore_color.rgb = clr; hd.line.fill.background()
    bul(s, [(title, True)], x+Inches(0.1), Inches(1.25), Inches(3.8), Inches(0.4), sz=13)
    bul(s, [f"•  {p}" for p in pts], x+Inches(0.15), Inches(1.9), Inches(3.7), Inches(3.5), sz=13)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 17. SECTION DIVIDER PART II
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution"); print(f"  {n}. Divider II")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PART II: SRMamba-T slides — typed text LEFT + diagram RIGHT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def arch_slide(title, text_items, diagram_name):
    """Left: typed text, Right: architecture diagram"""
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:40]}")
    bul(s, text_items, Inches(0.4), Inches(1.1), Inches(5.5), Inches(6), sz=13)
    img(s, diagram_name, Inches(6.2), Inches(1.1), Inches(6.9), Inches(6.0))

def arch_slide_full(title, text_items, diagram_name):
    """Text on top, full-width diagram below"""
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:40]}")
    bul(s, text_items, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.8), sz=13)
    img(s, diagram_name, Inches(0.3), Inches(2.9), Inches(12.7), Inches(4.4))

# 18. WHY HYBRID
arch_slide("Why Hybrid? — The Architecture Compromise", [
    ("No single architecture solves all SR challenges:", True), "",
    ("CNN:", True), "  ✗  Limited receptive field (local only)",
    "  ✗  Cannot capture long-range dependencies", "",
    ("Transformer:", True), "  ✓  Global context via self-attention",
    "  ✗  O(n²) complexity — GPU memory explosion", "",
    ("Mamba (SSM):", True), "  ✓  O(n) linear complexity — very efficient",
    "  ✗  1D scan destroys 2D spatial information", "",
    ("→ SRMamba-T: Best of Mamba + Transformer", True),
], "gen_srmambat.png")

# 19. COMPARING ARCHITECTURES + ERF
arch_slide("CNN vs Transformer vs Mamba — ERF Analysis", [
    ("Effective Receptive Field (ERF):", True), "",
    "CNN: Small, localized ERF",
    "  — only sees nearby pixels",
    "Transformer: Large but uniform ERF",
    "  — sees everything but wastes compute",
    "Mamba: Directional ERF along scan line",
    "  — efficient but loses 2D context", "",
    ("The Hybrid Insight:", True),
    "Phase 1 — Mamba (early layers):",
    "  Dense feature extraction at O(n) cost", "",
    "Phase 2 — Transformer (later layers):",
    "  Global refinement on compressed features",
    "  Un-windowed attention without memory explosion",
], "my_07.png")

# 20. COMPLETE ARCHITECTURE
arch_slide_full("The Complete SRMamba-T Architecture", [
    "SRMamba-T uses a U-Net encoder with asymmetric layer placement: Mamba layers (early) for efficient dense feature extraction,",
    "Transformer layers (late) for global spatial refinement. Features flow through FSM modules, Mamba blocks, and Transformer blocks.",
], "gen_srmambat.png")

# 21. STAGE-WISE FEATURE FLOW
arch_slide_full("Stage-Wise Feature Flow", [
    "Input → Shallow Feature Extraction → Stage 1 (Mamba-heavy) → Stage 2 (Mixed) → Stage 3 (Transformer-heavy) → Reconstruction → HR Output.",
    "Each stage uses Feature Sharing Modules (FSM) for cross-stage information flow.",
], "my_10.png")

# 22. SSM FOUNDATIONS
arch_slide("SSM Foundations & Discretization", [
    ("State Space Models (SSMs):", True),
    "Continuous-time linear system:", "",
    "  h'(t) = A·h(t) + B·x(t)",
    "  y(t)  = C·h(t) + D·x(t)", "",
    ("Discretization:", True),
    "For deep learning, continuous SSM must be",
    "discretized using Zero-Order Hold (ZOH):", "",
    "  Ā = exp(Δ·A)",
    "  B̄ = (Δ·A)⁻¹(exp(Δ·A) - I)·Δ·B", "",
    ("Selective Scan (Mamba/S6):", True),
    "Makes A, B, C, Δ input-dependent",
    "→ Content-aware state transitions",
], "my_11.png")

# 23. SSM CONVOLUTION EQUIVALENCE
arch_slide("SSM: Equivalence to Linear Convolution", [
    ("Key Insight:", True),
    "The discretized SSM can be computed as a", 
    "global convolution — enabling parallel training.", "",
    ("Convolutional Form:", True),
    "  K = (CB̄, CĀB̄, CĀ²B̄, ..., CĀᴸB̄)",
    "  y = x * K  (global convolution)", "",
    ("This means:", True),
    "•  Training: Use convolution (parallel, fast)",
    "•  Inference: Use recurrence (sequential, memory-efficient)", "",
    ("Benefit for SR:", True),
    "Images can be processed at training time",
    "using parallel convolution, but at inference",
    "using efficient sequential state updates.",
], "my_13.png")

# 24. MAMBA LAYER ANATOMY
arch_slide("Mamba Layer: Anatomy & Mathematics", [
    ("Mamba Block Structure:", True), "",
    "Input → LayerNorm → Linear Projection →",
    "  Branch 1: Conv1D → SiLU → SSM",
    "  Branch 2: SiLU activation",
    "→ Element-wise multiply → Linear → Output", "",
    ("Key Equations:", True),
    "  z = σ(Linear(Norm(x)))",
    "  h = SSM(SiLU(Conv1D(Linear(Norm(x)))))",
    "  y = Linear(h ⊙ z) + x", "",
    ("The SSM inside processes sequences", True),
    ("with input-dependent state transitions", True),
], "gen_mamba_block.png")

# 24. ASYMMETRIC MAMBA MIXER
arch_slide("The Asymmetric Mamba Mixer Block", [
    ("Why Asymmetric?", True),
    "Standard Mamba scans in one direction only.", "",
    ("The Mixer combines:", True),
    "  •  Forward scan branch",
    "  •  Backward scan branch",
    "  •  Different parameter counts per direction", "",
    "This captures bidirectional dependencies",
    "while maintaining computational efficiency.", "",
    ("Key Innovation:", True),
    "Asymmetric design allows the network to",
    "allocate more capacity to the direction",
    "that provides more useful information.",
], "my_15.png")

# 26. SPATIAL DEPENDENCY PROBLEM
arch_slide("The Spatial Dependency Problem in Mamba", [
    ("Why 1D scanning fails for images:", True), "",
    "Images are 2D grids — every pixel has spatial",
    "neighbors in all directions (up, down, left, right).", "",
    ("1D Flattening Problem:", True),
    "When a 2D image is flattened to 1D sequence:",
    "  •  Adjacent pixels in vertical direction become",
    "     distant in the 1D sequence",
    "  •  Spatial proximity information is lost",
    "  •  The SSM cannot model vertical dependencies", "",
    ("Visual Example:", True),
    "Pixel (2,3) and (3,3) are neighbors in 2D",
    "but may be W positions apart in 1D scan",
    "→ SSM 'forgets' their relationship",
], "gen_mdssm.png")

# 27. MDSSM
arch_slide("Multi-Directional Selective Scan (MDSSM)", [
    ("The Spatial Dependency Problem:", True),
    "Standard Mamba processes 2D images as 1D",
    "sequences → loses spatial relationships.", "",
    ("MDSSM Solution:", True),
    "1) Channel Split: Split features into 4 groups",
    "2) 4-Direction Scanning:",
    "   → Top-to-Bottom",
    "   → Bottom-to-Top",
    "   → Left-to-Right",
    "   → Right-to-Left",
    "3) Concatenate: Merge scanned features", "",
    ("Result:", True),
    "Preserves 2D spatial context at O(n) cost!",
], "gen_mdssm.png")

# 26. MDSSM 2D CONTEXT
arch_slide("MDSSM: Preserving 2D Spatial Context", [
    ("Why 4 directions?", True),
    "Each scan captures dependencies along",
    "one axis. 4 directions ensure every pixel", 
    "receives context from all spatial neighbors.", "",
    ("Channel Split Strategy:", True),
    "Instead of running 4 full SSMs (4× cost),",
    "MDSSM splits channels into 4 groups.",
    "Each group runs one direction → only 1× cost!", "",
    ("Mathematical Formulation:", True),
    "  C₁,C₂,C₃,C₄ = Split(X, dim=channel)",
    "  Cᵢ' = SSM(Scan_direction_i(Cᵢ))",
    "  Output = Concat(C₁',C₂',C₃',C₄')", "",
    "Parameter overhead: ~0% (same total params)",
], "my_18.png")

# 27. TRANSFORMER LAYER
arch_slide("Inside the Transformer Layer", [
    ("Transformer Layer Structure:", True), "",
    "Input → LayerNorm → Multi-Head Self-Attention →",
    "Residual Connection → LayerNorm → FFN →",
    "Residual Connection → Output", "",
    ("Self-Attention Mechanism:", True),
    "  Q = X · Wq  (Queries)",
    "  K = X · Wk  (Keys)", 
    "  V = X · Wv  (Values)",
    "  Attention(Q,K,V) = softmax(QKᵀ/√d) · V", "",
    ("In SRMamba-T:", True),
    "Transformer layers are placed in later stages",
    "where feature maps are already compressed",
    "→ quadratic cost is manageable",
], "my_19.png")

# 28. UN-WINDOWED ATTENTION
arch_slide("Un-Windowed Global Self-Attention", [
    ("Problem with Windowed Attention:", True),
    "SwinIR uses shifted windows → limits the",
    "receptive field to window boundaries.", "",
    ("SRMamba-T Solution:", True),
    "Un-windowed attention: every pixel attends",
    "to every other pixel in the feature map.", "",
    "This is possible because:",
    "  •  Mamba layers already compressed features",
    "  •  Feature maps are smaller at later stages",
    "  •  Memory cost is manageable", "",
    ("Benefit:", True),
    "True global context without window artifacts",
    "→ Better texture reconstruction at boundaries",
], "my_20.png")

# 29. FFM
arch_slide("Feature Fusion Module (FFM)", [
    ("Purpose:", True),
    "Combines features from Mamba and Transformer",
    "branches effectively.", "",
    ("FFM Process:", True),
    "  1) Receive Mamba features (F_mamba)",
    "  2) Receive Transformer features (F_trans)", 
    "  3) Concatenate along channel dimension",
    "  4) Apply 1×1 convolution for channel reduction",
    "  5) Add residual connection", "",
    ("Why needed:", True),
    "Mamba captures sequential/local patterns",
    "Transformer captures global/spatial patterns",
    "FFM learns optimal fusion weights",
], "my_21.png")

# 30. RECONSTRUCTION + LOSS
arch_slide("Adaptive Reconstruction & Dual-Domain Loss", [
    ("Reconstruction Pipeline:", True),
    "Features → Conv layers → PixelShuffle →",
    "Sub-pixel upsampling → Final HR image", "",
    ("Dual-Domain Loss Function:", True),
    "  L_total = L_pixel + λ · L_freq", "",
    "L_pixel: Standard L1/MSE pixel-wise loss",
    "  Ensures pixel accuracy", "",
    "L_freq: FFT-based frequency domain loss",
    "  L_freq = ||FFT(HR) - FFT(SR)||₁",
    "  Preserves high-frequency detail", "",
    ("Why dual-domain?", True),
    "Pixel loss alone → blurry outputs",
    "FFT loss → preserves sharp edges & textures",
], "my_22.png")

# 31. TRAINING PIPELINE
arch_slide("Training Pipeline & Setup", [
    ("Training Configuration:", True), "",
    "•  Dataset: DIV2K (800 training images)",
    "•  Patch size: 64×64 LR patches",
    "•  Batch size: 16",
    "•  Optimizer: Adam (β₁=0.9, β₂=0.99)",
    "•  Learning rate: 2×10⁻⁴ with cosine decay",
    "•  Total iterations: 500K", "",
    ("Degradation Model:", True),
    "•  Bicubic downsampling (×2, ×3, ×4)",
    "•  Data augmentation: random flips, rotations", "",
    ("Hardware:", True),
    "•  NVIDIA A100 GPU, PyTorch framework",
], "my_23.png")

# 32. QUANTITATIVE RESULTS
arch_slide("Quantitative Results: PSNR vs Compute", [
    ("Key Results:", True), "",
    "SRMamba-T achieves:", "",
    "  •  +0.28 dB PSNR improvement over MambaIR",
    "  •  38.7% fewer MACs (computational cost)",
    "  •  Better SSIM on all benchmark datasets", "",
    ("Benchmarks tested:", True),
    "  •  Set5, Set14, BSD100, Urban100, Manga109", "",
    ("Comparison vs:", True),
    "  •  SRCNN, VDSR (CNN-based)",
    "  •  SwinIR (Transformer-based)",
    "  •  MambaIR (Mamba-based)",
    "  •  SRMamba-T outperforms ALL with less compute",
], "my_24.png")

# 33. QUALITATIVE RESULTS
arch_slide("Qualitative Validation: Visual Comparison", [
    ("Visual comparison on Urban100 dataset:", True), "",
    "•  Complex textures (buildings, fences, grids)",
    "•  Fine structural patterns",
    "•  High-frequency detail regions", "",
    ("Observations:", True),
    "•  CNN models: Blurry, loss of detail",
    "•  GAN models: Sharp but artifacts",
    "•  SwinIR: Good but misses some textures",
    "•  MambaIR: Directional artifacts", "",
    ("SRMamba-T:", True),
    "•  Sharpest reconstruction",
    "•  No directional artifacts",
    "•  Best structural preservation",
], "my_25.png")

# 34. LAM ANALYSIS
arch_slide("LAM Analysis: Inside the Model's Mind", [
    ("LAM = Local Attribution Map", True), "",
    "Visualizes which input pixels contributed",
    "most to each output pixel.", "",
    ("Diffusion Index (DI) comparison:", True),
    "  •  CNN models: Low DI — very local attention",
    "  •  SwinIR: Medium DI — window-limited",
    "  •  MambaIR: Higher DI but directional",
    "  •  SRMamba-T: Highest DI — truly global", "",
    ("What this means:", True),
    "SRMamba-T uses information from the widest",
    "range of input pixels → better reconstruction",
    "of complex patterns and structures.",
], "my_26.png")

# 35. HARDWARE + ABLATION
arch_slide("Hardware Efficiency & Ablation Study", [
    ("Memory Scaling:", True),
    "•  Transformer: Quadratic O(n²) memory growth",
    "•  SRMamba-T: Near-linear memory scaling", "",
    "At 1024×1024 resolution:",
    "  SwinIR: Out of memory",
    "  SRMamba-T: Runs comfortably", "",
    ("Ablation: Layer Ordering", True),
    "Tested configurations:",
    "  A) All Mamba → Poor global context",
    "  B) All Transformer → Too expensive",
    "  C) Trans first, Mamba later → Suboptimal",
    "  D) Mamba first, Trans later → BEST ✓", "",
    ("Conclusion: Mamba early + Transformer late", True),
    ("= optimal efficiency + quality tradeoff", True),
], "my_27.png")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 36. CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "Conclusion & Future Work"); print(f"  {n}. Conclusion")
sub_hdr(s, "Key Contributions", Inches(0.4), Inches(1.1), Inches(5.8))
bul(s, [
    "✓  First hybrid Mamba-Transformer for image SR", "",
    "✓  MDSSM: 4-directional scanning preserves",
    "   2D spatial context at O(n) cost", "",
    "✓  Asymmetric placement: Mamba (early) +",
    "   Transformer (late) = optimal tradeoff", "",
    "✓  +0.28 dB PSNR, 38.7% fewer MACs", "",
    "✓  Dual-domain loss prevents artifacts", "",
    "✓  Linear memory scaling for large images",
], Inches(0.4), Inches(1.7), Inches(5.8), Inches(5.5), sz=15)
sub_hdr(s, "Future Work", Inches(6.6), Inches(1.1), Inches(6.3))
bul(s, [
    "•  Extend to Hyperspectral SR (PaviaC, PaviaU)", "",
    "•  Scale to SRMamba-T-L for classical SR", "",
    "•  Multi-scale U-Net with downsampling", "",
    "•  Real-world degradation beyond bicubic", "",
    "•  Edge device deployment (linear memory)", "",
    "•  Combine with diffusion models",
], Inches(6.6), Inches(1.7), Inches(6.3), Inches(5.5), sz=15)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 37. REFERENCES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s); hdr(s, "References"); print(f"  {n}. References")
bul(s, [
    "[1]  Dong et al., \"Image Super-Resolution Using Deep Convolutional Networks\" (SRCNN), ECCV 2014",
    "[2]  Dong et al., \"Accelerating the Super-Resolution CNN\" (FSRCNN), ECCV 2016",
    "[3]  Kim et al., \"Accurate Image SR Using Very Deep Networks\" (VDSR), CVPR 2016",
    "[4]  Ledig et al., \"Photo-Realistic SISR Using a GAN\" (SRGAN), CVPR 2017",
    "[5]  Wang et al., \"ESRGAN: Enhanced Super-Resolution GANs\", ECCV 2018",
    "[6]  Liang et al., \"SwinIR: Image Restoration Using Swin Transformer\", ICCV 2021",
    "[7]  Guo et al., \"MambaIR: A Simple Baseline for Image Restoration with SSM\", ECCV 2024",
    "[8]  Gu & Dao, \"Mamba: Linear-Time Sequence Modeling with Selective State Spaces\", 2023",
    "[9]  SRMamba-T: Hybrid Mamba-Transformer for Single Image Super-Resolution, 2024",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017",
], Inches(0.4), Inches(1.2), Inches(12.3), Inches(5.8), sz=14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 38. THANK YOU
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = S(); bg(s, NAVY); print(f"  {n}. Thank You")
s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H).fill.solid()
s.shapes[-1].fill.fore_color.rgb = ACCENT; s.shapes[-1].line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
tb3 = s.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]
p3.text = "Vujja Punith Sai  |  S20240010260  |  Guide: Arun PV"
p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p3.alignment = PP_ALIGN.CENTER

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v5_FINAL.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
