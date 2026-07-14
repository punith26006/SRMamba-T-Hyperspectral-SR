"""
V9 FINAL — Professor-Ready Presentation
- Standard 4:3 slides (10x7.5) - no wasted widescreen space
- Clean Ruthwik-style design with minimal decorations
- Content fills the full slide area
- Professional Calibri font throughout
- 36 slides, every slide packed with info
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
AI = os.path.join(BASE, "_temp_scripts", "ai_diagrams")
MY = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")
FONT = "Calibri"

# Colors
BG   = RGBColor(0xF5, 0xF5, 0xF0)   # Warm off-white
DK   = RGBColor(0x2D, 0x2D, 0x2D)   # Dark charcoal
HDR  = RGBColor(0x1A, 0x1A, 0x2E)   # Navy header
ACC  = RGBColor(0xE8, 0x6C, 0x00)   # Deep orange accent
TT   = RGBColor(0x1A, 0x1A, 0x1A)   # Near-black title
TX   = RGBColor(0x33, 0x33, 0x33)   # Dark grey body
SUB  = RGBColor(0x66, 0x66, 0x66)   # Subtle grey
WH   = RGBColor(0xFF, 0xFF, 0xFF)

# Standard 4:3 slide
W = Inches(13.333); H = Inches(7.5)
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]; n = 0

def fade(s):
    t = etree.SubElement(s._element, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    t.set('spd', 'med'); t.set('advClick', '1')
    f = etree.SubElement(t, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
    f.set('thruBlk', '0')

def S():
    global n; n += 1
    s = prs.slides.add_slide(BL); fade(s); return s

def bg(s):
    """Clean background: off-white + thin navy header bar + subtle orange bottom accent"""
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    # Thin header bar
    hb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    hb.fill.solid(); hb.fill.fore_color.rgb = HDR; hb.line.fill.background()
    # Thin bottom accent line
    bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), W, Inches(0.15))
    bb.fill.solid(); bb.fill.fore_color.rgb = ACC; bb.line.fill.background()

def hdr(s, title):
    """Title with orange underline"""
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = HDR; p.font.name = FONT
    # Orange underline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.75), Inches(3), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACC; ln.line.fill.background()

def txt(s, items, x, y, w, h, sz=15, boxed=False):
    """Text block — fills given area"""
    if boxed:
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xE8)
        bx.line.color.rgb = RGBColor(0xCC, 0xCC, 0xBB); bx.line.width = Pt(1)
    pad = Inches(0.15) if boxed else 0
    tb = s.shapes.add_textbox(x + pad, y + pad, w - 2*pad, h - 2*pad)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.font.bold = item[1]
            p.font.color.rgb = HDR if item[1] else TX
        else:
            p.text = item; p.font.color.rgb = TX
        p.font.size = Pt(sz); p.font.name = FONT; p.space_after = Pt(3)

def img(s, folder, name, x, y, w, h=None):
    path = os.path.join(folder, name)
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)
    else:
        print(f"    WARN: missing {name}")

def divider(s, num, title, sub=""):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = HDR; r.line.fill.background()
    # Orange accent bar
    ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), W, Inches(0.06))
    ab.fill.solid(); ab.fill.fore_color.rgb = ACC; ab.line.fill.background()
    # Circle with number
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.67), Inches(1.2), Inches(2.0), Inches(2.0))
    c.fill.solid(); c.fill.fore_color.rgb = ACC; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(5.67), Inches(1.35), Inches(2.0), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = WH; p.alignment = PP_ALIGN.CENTER; p.font.name = FONT
    # Title
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(38); p2.font.bold = True; p2.font.color.rgb = WH; p2.alignment = PP_ALIGN.CENTER; p2.font.name = FONT
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub
        p3.font.size = Pt(18); p3.font.color.rgb = RGBColor(0xBB,0xBB,0xBB); p3.alignment = PP_ALIGN.CENTER; p3.font.name = FONT

# Content helper for Part II slides
def p2(title, items, img_name, img_folder=MY):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:42]}")
    txt(s, items, Inches(0.4), Inches(0.95), Inches(5.8), Inches(6.2), sz=14)
    img(s, img_folder, img_name, Inches(6.5), Inches(0.95), Inches(6.5), Inches(6.2))

def p2_wide(title, desc_items, img_name, img_folder=MY, detail_items=None):
    """Full-width diagram with description above and details below"""
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:42]}")
    txt(s, desc_items, Inches(0.4), Inches(0.9), Inches(12.5), Inches(1.1), sz=14)
    img(s, img_folder, img_name, Inches(0.3), Inches(2.1), Inches(12.7), Inches(3.2))
    if detail_items:
        txt(s, detail_items, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.8), sz=12)

print("Building ~36 slides — Professor-Ready...")

# ════════════════════════════════════════════════
# 1. TITLE
# ════════════════════════════════════════════════
s = S(); print(f"  {n}. Title")
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
r.fill.solid(); r.fill.fore_color.rgb = HDR; r.line.fill.background()
# Accent triangle
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9), Inches(0), Inches(4.33), Inches(5))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = ACC; tri.line.fill.background()
# Grid squares
for row in range(3):
    for col in range(3):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5+col*0.8), Inches(3.5+row*0.65), Inches(0.6), Inches(0.5))
        sq.line.color.rgb = RGBColor(0x55,0x55,0x55); sq.line.width = Pt(1); sq.fill.background()
# Title text
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
for line in ["Exploration of Deep", "Learning Models for", "Image Super Resolution"]:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = line; p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = WH; p.font.name = FONT
# Orange line
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = ACC; ln.line.fill.background()
# Details
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6), Inches(2))
tf2 = tb2.text_frame; tf2.word_wrap = True
for line in ["B.Tech Project Presentation", "", "Name: Vujja Punith Sai", "Roll No: S20240010260", "Instructor: Arun PV"]:
    p = tf2.paragraphs[0] if not tf2.paragraphs[0].text else tf2.add_paragraph()
    p.text = line; p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xBB,0xBB,0xBB); p.font.name = FONT

# ════════════════════════════════════════════════
# 2. CONTENTS
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Contents"); print(f"  {n}. Contents")
txt(s, [("Part I  —  Literature Review", True), "",
    "  1. What is Super Resolution & Need for SR",
    "  2. Problem Statement",
    "  3. Traditional SR Methods",
    "  4. Deep Learning in SR & Evolution",
    "  5. Multimodal Learning",
    "  6. CNN Models: SRCNN, FSRCNN, VDSR",
    "  7. GAN Models: SRGAN & ESRGAN",
    "  8. Transformer: SwinIR",
    "  9. State Space: SR-Mamba",
    " 10. Evaluation Metrics (MSE, PSNR, SSIM)"],
    Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.2), sz=15, boxed=True)
txt(s, [("Part II  —  SRMamba-T Architecture", True), "",
    " 11. Why Hybrid? CNN vs Transformer vs Mamba",
    " 12. Complete SRMamba-T Architecture",
    " 13. SSM Foundations & Discretization",
    " 14. Mamba Layer & Asymmetric Mixer",
    " 15. MDSSM: Multi-Directional Selective Scan",
    " 16. Transformer Layer & Feature Fusion (FFM)",
    " 17. Reconstruction, Loss & Training Pipeline",
    " 18. Quantitative & Qualitative Results",
    " 19. LAM Analysis & Ablation Study",
    " 20. Conclusion, Future Work & References"],
    Inches(6.8), Inches(0.95), Inches(6.2), Inches(6.2), sz=15, boxed=True)

# ════════════════════════════════════════════════
# 3. DIVIDER I
# ════════════════════════════════════════════════
s = S(); divider(s, "I", "Literature Review", "Background & Related Work"); print(f"  {n}. Divider I")

# ════════════════════════════════════════════════
# 4. WHAT IS SR + NEED
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "What is Super Resolution?"); print(f"  {n}. What is SR")
txt(s, [
    "Super Resolution is a Deep Learning and Computer Vision technique used to reconstruct a",
    "high-resolution image from one or more low-resolution images.", "",
    ("It aims to recover:", True),
    "  1. Lost image details    2. Textures    3. Edges    4. Structural information", "",
    ("Main Objectives:", True),
    "1. Improve Image Clarity: Enhances sharpness and visibility of image contents.",
    "2. Recover Image Details: Reconstructs missing fine details and structures.",
    "3. Enhance Texture and Edges: Produces sharper texture and more realistic edges.",
    "4. Improve Perceptual Quality: Generates visually pleasing and natural-looking images.",
], Inches(0.4), Inches(0.95), Inches(8.0), Inches(6.2), sz=15)
txt(s, [("Basic SR Workflow", True), "",
    "  Low Resolution Image", "         |", "  Feature Extraction", "         |",
    "  Super Resolution Model", "         |", "  Image Reconstruction", "         |",
    "  High Resolution Image"],
    Inches(8.8), Inches(0.95), Inches(4.2), Inches(6.2), sz=15, boxed=True)

# ════════════════════════════════════════════════
# 5. NEED + APPLICATIONS
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Need for Super Resolution"); print(f"  {n}. Need for SR")
txt(s, [
    ("Problems in Low Resolution Images:", True),
    "  * Blur and loss of sharpness — objects become unrecognizable",
    "  * Pixelation during zooming — jagged edges and block artifacts",
    "  * Loss of fine textures and edges — surface details disappear",
    "  * Poor object visibility — critical features obscured",
    "  * Noise amplification during upscaling",
    "  * Information loss during compression and transmission",
], Inches(0.4), Inches(0.95), Inches(6.2), Inches(6.2), sz=15, boxed=True)
txt(s, [
    ("Real-World Applications:", True),
    "  * Satellite Imaging: terrain analysis, urban planning, crop monitoring",
    "  * Surveillance & Security: face recognition, license plate reading",
    "  * Medical Imaging: MRI/CT scan enhancement, tumor detection",
    "  * Mobile Photography: AI zoom, low-light enhancement",
    "  * Video Upscaling: SD to HD/4K conversion for streaming",
    "  * Remote Sensing: disaster response, geological surveys",
    "  * Astronomy: telescope image enhancement",
], Inches(6.8), Inches(0.95), Inches(6.2), Inches(6.2), sz=15, boxed=True)

# ════════════════════════════════════════════════
# 6. PROBLEM STATEMENT
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Problem Statement"); print(f"  {n}. Problem Statement")
txt(s, [
    ("Context:", True),
    "Most open-source data have coarse spatial resolution, and it is advantageous to use different",
    "sensors for effective enhancement; however, the sensors have differences in spatial resolution,",
    "spectral characteristics, noise levels, and often lack spatial or temporal overlap.", "",
    ("Key Challenges in Single Image Super Resolution (SISR):", True),
    "  * Recovering lost high-frequency details (edges, textures, patterns) from degraded inputs",
    "  * CNN models have limited receptive fields — can only see local neighbourhoods",
    "  * Transformer models suffer from O(n^2) computational complexity — GPU memory explosion",
    "    for high-resolution images (e.g., 1024x1024 requires >40GB VRAM with SwinIR)",
    "  * Pure Mamba models scan images as 1D sequences — destroys 2D spatial context",
    "    (vertical neighbours become distant in the 1D flattened sequence)", "",
    ("Our Goal:", True),
    "Design a hybrid architecture (SRMamba-T) that combines Mamba's O(n) efficiency with",
    "Transformer's global awareness, achieving state-of-the-art PSNR with significantly fewer MACs.",
], Inches(0.4), Inches(0.95), Inches(12.5), Inches(6.2), sz=15)

# ════════════════════════════════════════════════
# 7. TRADITIONAL + DL INTRO
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Traditional Methods & Deep Learning for SR"); print(f"  {n}. Traditional + DL")
txt(s, [
    ("Traditional Interpolation Methods:", True),
    "  1) Nearest Neighbour: Copies nearest pixel. Very fast but blocky, jagged.",
    "  2) Bilinear: Weighted average of 4 neighbours. Smoother but still blurry.",
    "  3) Bicubic: Uses 16 surrounding pixels. Better continuity but cannot recover missing details.", "",
    "All traditional methods FAIL to reconstruct truly missing high-frequency information.",
    "They only interpolate existing pixels — they cannot hallucinate new detail.", "",
    ("Why Deep Learning?", True),
    "Instead of mathematically scaling pixels, DL models directly learn the mapping LR -> HR",
    "from large training datasets. They learn to generate plausible high-frequency details.", "",
    ("What DL Models Learn:", True),
    "  * Texture Patterns — roads, vegetation, fabric, faces",
    "  * Edge Information — object boundaries, transitions",
    "  * Image Structures — spatial relationships, symmetry, repetition",
], Inches(0.4), Inches(0.95), Inches(12.5), Inches(6.2), sz=15)

# ════════════════════════════════════════════════
# 8. EVOLUTION + MULTIMODAL
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Evolution of SR & Multimodal Learning"); print(f"  {n}. Evolution")
txt(s, [
    ("Evolution Timeline:", True),
    "  2014: CNN Era begins — SRCNN (first DL-based SR)",
    "  2016: Deeper CNNs — FSRCNN (fast), VDSR (20 layers, residual learning)",
    "  2017: GAN Era — SRGAN (perceptual loss, adversarial training)",
    "  2018: Enhanced GANs — ESRGAN (RRDB blocks, relativistic discriminator)",
    "  2021: Transformer Era — SwinIR (shifted window attention, global context)",
    "  2023: SSM Era — MambaIR (state space models, O(n) complexity)",
    "  2024: Hybrid Era — SRMamba-T (Mamba + Transformer, best of both)",
], Inches(0.4), Inches(0.95), Inches(5.8), Inches(3.0), sz=14, boxed=True)
txt(s, [
    ("Multimodal Learning:", True),
    "  * Single-Modality SR: Uses only a single LR image (visual only).",
    "  * Multimodal SR: Combines multiple data sources (e.g., satellite + elevation).",
    "  * Our work focuses on single-image SR (SISR) — the hardest variant.",
], Inches(0.4), Inches(4.2), Inches(5.8), Inches(3.0), sz=14, boxed=True)
img(s, AI, "evolution.png", Inches(6.5), Inches(0.95), Inches(6.5), Inches(6.2))

# ════════════════════════════════════════════════
# 9. SRCNN + FSRCNN + VDSR
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "CNN-based Models: SRCNN, FSRCNN, VDSR"); print(f"  {n}. CNN Models")
txt(s, [
    ("I. SRCNN (2014) — First DL model for SR", True),
    "  * End-to-end learning: LR -> HR mapping via 3 conv layers.",
    "  * Layers: Feature Extraction -> Non-Linear Mapping -> Reconstruction.",
    "  * Loss: MSE = (1/n) Sum ||F(Yi;theta) - Xi||^2. Simple but blurry.", "",
    ("II. FSRCNN (2016) — Fast Super Resolution CNN", True),
    "  * Operates on LR space directly (no pre-upsampling needed).",
    "  * Deconvolution at the end for upsampling. 5x faster than SRCNN.",
    "  * Still limited receptive field (local context only).", "",
    ("III. VDSR (2016) — Very Deep SR Network (20 layers)", True),
    "  * Key innovation: Global Residual Learning — HR = LR + Residual.",
    "  * 41x41 pixel receptive field (much larger than SRCNN's 13x13).",
    "  * Multi-scale: single model handles x2, x3, x4 upscaling.",
    "  * Limitation: Operates in HR space (computationally expensive).",
], Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.2), sz=14)
img(s, AI, "srcnn.png", Inches(6.5), Inches(0.95), Inches(6.5), Inches(6.2))

# ════════════════════════════════════════════════
# 10. SRGAN + ESRGAN
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "GAN-based Models: SRGAN & ESRGAN"); print(f"  {n}. GAN Models")
txt(s, [
    ("I. SRGAN (2017) — First GAN for Super Resolution", True),
    "  * Problem: CNN models optimized for MSE produce blurry, over-smoothed outputs.",
    "  * Solution: Adversarial training forces Generator to produce realistic textures.",
    "  * Generator: Deep ResNet with 16 residual blocks + PixelShuffle upsampling.",
    "  * Discriminator: VGG-style network classifies images as Real/Fake.",
    "  * Key Innovation: Perceptual Loss using VGG feature maps -> sharper textures.", "",
    ("II. ESRGAN (2018) — Enhanced SRGAN", True),
    "  * Residual-in-Residual Dense Blocks (RRDB) — multi-level residual connections.",
    "  * Removed Batch Normalization for more stable training.",
    "  * Relativistic Discriminator: predicts relative quality (not just real/fake).",
    "  * Won PIRM2018 SR Challenge. Better texture, higher perceptual quality.", "",
    ("Limitations of GAN-based SR:", True),
    "  * Training instability, mode collapse, hallucinated artifacts.",
    "  * Cannot guarantee pixel-level accuracy (trades fidelity for perception).",
], Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.2), sz=13)
img(s, AI, "srgan.png", Inches(6.5), Inches(0.95), Inches(6.5), Inches(6.2))

# ════════════════════════════════════════════════
# 11. SWINIR + SR-MAMBA
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Transformer & SSM Models: SwinIR, SR-Mamba"); print(f"  {n}. SwinIR + SR-Mamba")
txt(s, [
    ("I. SwinIR (2021) — Swin Transformer for Image Restoration", True),
    "  * Built on Vision Transformers (ViT) with Shifted Window Attention.",
    "  * Self-attention determines which image regions are most relevant.",
    "  * Architecture: Shallow Feature Ext -> RSTB blocks -> HQ Reconstruction.",
    "  * Advantage: Global feature learning, excellent reconstruction quality.",
    "  * Limitation: O(n^2) attention complexity -> high GPU memory cost.", "",
    ("II. SR-Mamba (2024) — State Space Model for SR", True),
    "  * Based on Mamba (Selective State Space Model / S6).",
    "  * Replaces attention with State Transition Mechanisms.",
    "  * O(n) linear complexity — highly efficient for long sequences.",
    "  * Limitation: 1D sequential scan destroys 2D spatial context.",
    "  * Vertical neighbours become distant when image is flattened to 1D.", "",
    ("-> This motivates our hybrid approach: SRMamba-T!", True),
], Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.2), sz=13)
img(s, AI, "swinir.png", Inches(6.5), Inches(0.95), Inches(3.1), Inches(2.9))
img(s, AI, "ssm_mamba.png", Inches(6.5), Inches(4.1), Inches(3.1), Inches(2.9))

# ════════════════════════════════════════════════
# 12. EVALUATION METRICS
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Evaluation Metrics: MSE, PSNR, SSIM"); print(f"  {n}. Metrics")
txt(s, [
    ("1. MSE (Mean Square Error)", True),
    "  MSE = (1/mn) Sum_i,j (I(i,j) - K(i,j))^2",
    "  Pixel-wise reconstruction error. Lower is better.",
    "  Used in SRCNN, VDSR. Simple but produces blurry results.",
], Inches(0.4), Inches(0.95), Inches(3.9), Inches(6.2), sz=14, boxed=True)
txt(s, [
    ("2. PSNR (Peak Signal-to-Noise Ratio)", True),
    "  PSNR = 10 * log10(MAX^2 / MSE)  [in dB]",
    "  Reconstruction quality measure. Higher is better.",
    "  Inversely related to MSE. Standard benchmark metric.",
], Inches(4.6), Inches(0.95), Inches(4.0), Inches(6.2), sz=14, boxed=True)
txt(s, [
    ("3. SSIM (Structural Similarity)", True),
    "  SSIM(x,y) = f(luminance, contrast, structure)",
    "  Measures perceptual/structural quality. Range [0,1].",
    "  SSIM -> 1: excellent.  SSIM -> 0: poor.",
    "  Better correlates with human perception than PSNR.",
], Inches(8.8), Inches(0.95), Inches(4.2), Inches(6.2), sz=14, boxed=True)

# ════════════════════════════════════════════════
# 13. DIVIDER II
# ════════════════════════════════════════════════
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution"); print(f"  {n}. Divider II")

# ════════════════════════════════════════════════
# 14. WHY HYBRID
# ════════════════════════════════════════════════
p2("Why Hybrid? — The Architecture Compromise", [
    ("No single architecture is perfect:", True), "",
    ("CNN (SRCNN, VDSR):", True),
    "  + Fast, simple, good for local features.",
    "  - Limited receptive field (sees only neighbours).",
    "  - Cannot capture long-range dependencies.", "",
    ("Transformer (SwinIR):", True),
    "  + Global context via self-attention.",
    "  - O(n^2) complexity — memory explosion.",
    "  - Impractical for high-resolution inputs.", "",
    ("Mamba / SSM (MambaIR):", True),
    "  + O(n) linear complexity — very efficient.",
    "  - 1D scan destroys 2D spatial information.",
    "  - Directional artifacts in reconstruction.", "",
    ("SRMamba-T: Mamba (early) + Transformer (late)", True),
    "Dense extraction at O(n) cost, then global refinement",
    "on compressed features without memory explosion.",
], "erf.png", AI)

# ════════════════════════════════════════════════
# 15. COMPLETE ARCHITECTURE (PAPER DIAGRAM)
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "The Complete SRMamba-T Architecture"); print(f"  {n}. Complete Arch")
txt(s, [
    "Three-stage pipeline: (1) Shallow Feature Extraction via 3x3 Conv, (2) Deep Feature Extraction with N/2 Mamba Layers",
    "(encoder) followed by (N/2)-1 Transformer Layers (encoder), (3) Image Reconstruction via FFM + Decoder Transformer + PixelShuffle.",
], Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.9), sz=14)
img(s, AI, "srmambat_paper.png", Inches(0.2), Inches(1.9), Inches(12.9), Inches(3.5))
txt(s, [
    ("Key Components (bottom row of diagram):", True),
    "* Mamba Block: LayerNorm -> MambaMixer -> Scale -> LayerNorm -> 3x3Conv -> Channel Attention Block (CAB) -> Scale",
    "* Attention Block: LayerNorm -> Self-Attention -> Scale -> LayerNorm -> 3x3Conv -> Channel Attention -> Scale",
    "* Feature Fusion Module (FFM): Aggregates features from all encoder stages via DWConv -> SiLU -> PWLinear -> LayerNorm",
    "* Encoders use N/2 Mamba Layers (each with L Mamba Blocks) + (N/2)-1 Transformer Layers (each with L Attention Blocks)",
], Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.7), sz=12)

# ════════════════════════════════════════════════
# 16. STAGE-WISE FLOW
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Stage-Wise Feature Flow"); print(f"  {n}. Feature Flow")
txt(s, [
    "LR Input -> 3x3Conv (Fs) -> PatchEmbed -> N/2 Mamba Layers (F1) -> Mamba Layer -> (N/2)-1 Transformer Layers (Fh)",
    "-> 3x3Conv -> FFM (Ff) -> Decoder Transformer Layer (Fd) -> PatchUnEmbed -> Upsample -> SR Output",
], Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.9), sz=14)
img(s, AI, "srmambat_paper.png", Inches(0.2), Inches(1.9), Inches(12.9), Inches(3.5))
txt(s, [
    ("Detailed Data Flow:", True),
    "* Shallow: LR -> Fs = Conv3x3(LR) -> F1 = PatchEmbed(Fs) — converts 2D feature maps to 1D token sequences for Mamba/Transformer",
    "* Encoder: Each Mamba/Transformer Layer has L internal blocks + 3x3Conv + CAB. Uses learnable Scale factors S1, S2 for residual balancing.",
    "* FFM: Concatenates Fs (shallow features) with all encoder outputs -> 3x3 DWConv -> SiLU -> PWLinear -> LayerNorm",
    "* Reconstruction: F = Fs + F'd (element-wise add shallow features back) -> PixelShuffle upsample (r times) to HR resolution",
], Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.7), sz=12)

# ════════════════════════════════════════════════
# 17. SSM FOUNDATIONS
# ════════════════════════════════════════════════
p2("SSM Foundations & Discretization", [
    ("State Space Models (SSMs):", True),
    "Continuous-time linear dynamical system:",
    "  h'(t) = A * h(t) + B * x(t)    [state equation]",
    "  y(t)  = C * h(t) + D * x(t)    [output equation]", "",
    "A is NxN state matrix, B is Nx1 input matrix,",
    "C is 1xN output matrix. h(t) is hidden state.", "",
    ("Discretization via Zero-Order Hold (ZOH):", True),
    "  A_bar = exp(delta * A)",
    "  B_bar = (delta*A)^-1 * (exp(delta*A) - I) * delta*B",
    "  h_k = A_bar * h_{k-1} + B_bar * x_k",
    "  y_k = C * h_k", "",
    ("Selective Scan (Mamba S6):", True),
    "Makes delta, B, C input-dependent (not fixed).",
    "-> Content-aware state transitions. Key innovation",
    "over S4 and HiPPO that enables data-dependent filtering.",
], "slide_11.png")

# ════════════════════════════════════════════════
# 18. SSM CONVOLUTION + MAMBA LAYER
# ════════════════════════════════════════════════
p2("Mamba Layer: SSM Convolution & Block Design", [
    ("SSM Equivalence to Convolution:", True),
    "Discretized SSM can be reformulated as global conv:",
    "  K = (C*B_bar, C*A_bar*B_bar, ..., C*A_bar^{L-1}*B_bar)",
    "  y = x * K  (global convolution with kernel K)", "",
    "Training: Parallel FFT-based conv O(L log L).",
    "Inference: Sequential recurrence O(1) per step.", "",
    ("Mamba Block Structure (each layer has L blocks):", True),
    "X_hat = Mamba(Norm(X)) + S1 * X     [Eq.6]", "",
    "Inside Mamba Block:",
    "  Input -> LayerNorm -> Linear (expand 2x) ->",
    "    Branch A: Conv1D(k=3) -> SiLU -> SSM (MDSSM)",
    "    Branch B: SiLU activation (gate)",
    "  -> Element-wise multiply -> Linear (compress) -> Out", "",
    "After L blocks: X_{l+1} = CAB(Conv3x3(Norm(X_hat))) + S2 * X_hat  [Eq.7]",
], "slide_14.png")

# ════════════════════════════════════════════════
# 19. ASYMMETRIC MIXER + SPATIAL PROBLEM
# ════════════════════════════════════════════════
p2("Asymmetric Mixer & Spatial Dependency Problem", [
    ("The Asymmetric Mamba Mixer:", True),
    "Standard Mamba scans in one direction only.",
    "Mamba Mixer has two parallel branches:",
    "  Branch A: Linear -> Conv1D -> SiLU -> MDSSM",
    "  Branch B: Linear -> SiLU (gating mechanism)",
    "  Output = Linear(BranchA * BranchB)",
    "Uses MDSSM instead of standard SSM.", "",
    ("The Spatial Dependency Problem:", True),
    "Images are 2D grids of size HxW. When flattened",
    "to 1D sequence of length H*W:",
    "  * Pixel at (i,j) and (i+1,j) are vertical neighbours",
    "    in 2D but W positions apart in 1D sequence.",
    "  * SSM's exp(delta*A) decay causes the state to",
    "    'forget' information about vertical neighbours.",
    "  * Vertical edges, horizontal textures, and diagonal",
    "    patterns cannot be reconstructed properly.",
    "-> Solution: MDSSM (Multi-Directional Selective Scan).",
], "mdssm.png", AI)

# ════════════════════════════════════════════════
# 20. MDSSM
# ════════════════════════════════════════════════
p2("MDSSM: Multi-Directional Selective Scan", [
    ("MDSSM — Core Innovation of SRMamba-T (Eq.8):", True),
    "1) Channel Split: D channels -> 4 groups of D/4.",
    "2) 4-Direction Scanning:",
    "   -> Top-to-Bottom (vertical dependencies)",
    "   -> Bottom-to-Top (reverse vertical)",
    "   -> Left-to-Right (horizontal dependencies)",
    "   -> Right-to-Left (reverse horizontal)",
    "3) Each group processed by SSM in its scan direction.",
    "4) Concatenate all 4 groups back together.", "",
    ("Formal Definition:", True),
    "  X = [X1, X2, X3, X4] = ChannelSplit(X)",
    "  Xi' = SSM(Scan_direction_i(Xi))",
    "  Output = Concat(X1', X2', X3', X4')", "",
    ("Computational Cost:", True),
    "Same as single-direction SSM (D/4 x 4 = D)!",
    "Zero additional parameter overhead.",
    "Each pixel gets context from all 4 spatial directions.",
], "slide_17.png")

# ════════════════════════════════════════════════
# 21. TRANSFORMER + FFM
# ════════════════════════════════════════════════
p2("Transformer Layer & Feature Fusion Module", [
    ("Transformer Layer (L Attention Blocks):", True),
    "X_hat = Attention(Norm(X)) + S1 * X     [Eq.6]",
    "Multi-Head Self-Attention:",
    "  Q = X*Wq,  K = X*Wk,  V = X*Wv",
    "  Attn = softmax(Q*K^T / sqrt(d)) * V",
    "After L blocks: X_{l+1} = CAB(Conv(Norm(X_hat))) + S2*X_hat", "",
    ("Un-Windowed Global Attention:", True),
    "Unlike SwinIR (shifted windows), SRMamba-T uses full",
    "un-windowed attention: every token attends to ALL others.",
    "Affordable because Mamba already compressed features.", "",
    ("Feature Fusion Module (FFM):", True),
    "Aggregates hierarchical features from all encoder stages.",
    "  1) Concat Fs (shallow) + all encoder features",
    "  2) 3x3 DWConv -> SiLU -> PWLinear -> LayerNorm",
    "  3) Output Ff -> feeds into decoder Transformer Layer.",
], "slide_19.png")

# ════════════════════════════════════════════════
# 22. RECONSTRUCTION + LOSS + TRAINING
# ════════════════════════════════════════════════
p2("Reconstruction, Loss & Training", [
    ("Reconstruction Pipeline:", True),
    "F = Fs + Fd' (element-wise add shallow features back).",
    "Lightweight SR: Conv + PixelShuffle(r) for upscale.",
    "Classical SR: Multi-step progressive upsample.", "",
    ("Dual-Domain Loss Function (Eq.5):", True),
    "L(theta) = Sum ||I_HR - I_SR||_2 + lambda*||FFT(I_HR) - FFT(I_SR)||_2",
    "lambda = 0.1 (empirically tuned balance factor).",
    "Pixel loss -> accuracy. FFT loss -> sharp edges.", "",
    ("SRMamba-T-S Model Configuration:", True),
    "N=4 layers (2 Mamba + 1 Trans enc + 1 Trans dec).",
    "L=6 blocks/layer. D=48. Params: ~0.6M (lightweight).", "",
    ("Training:", True),
    "DIV2K (800 images), 64x64 LR patches, batch=16.",
    "Adam (beta1=0.9, beta2=0.99), LR=2e-4 cosine->1e-6.",
    "500K iterations. Bicubic degradation (x2, x3, x4).",
    "Hardware: NVIDIA A100, PyTorch 2.0, CUDA 11.8.",
], "slide_22.png")

# ════════════════════════════════════════════════
# 23. QUANTITATIVE RESULTS
# ════════════════════════════════════════════════
p2("Quantitative Results: PSNR vs Compute", [
    ("Key Results:", True),
    "SRMamba-T achieves state-of-the-art performance:", "",
    "  * +0.28 dB PSNR improvement over MambaIR",
    "  * 38.7% fewer MACs (computational cost)",
    "  * Better SSIM on all 5 benchmark datasets", "",
    ("Benchmark Datasets:", True),
    "  Set5, Set14, BSD100, Urban100, Manga109", "",
    ("Performance Comparison (x4 upscaling):", True),
    "  SRCNN:    30.48 dB / 0.8628",
    "  VDSR:     31.35 dB / 0.8838",
    "  SwinIR:   32.44 dB / 0.8983",
    "  MambaIR:  32.67 dB / 0.9001",
    "  SRMamba-T: 32.95 dB / 0.9023  (BEST)", "",
    "Outperforms ALL baselines with significantly less compute.",
], "slide_24.png")

# ════════════════════════════════════════════════
# 24. QUALITATIVE RESULTS
# ════════════════════════════════════════════════
p2("Qualitative Validation: Visual Comparison", [
    ("Tested on Urban100 (complex textures):", True),
    "  * Buildings with grid patterns",
    "  * Fences and wire mesh textures",
    "  * Brick walls and tiled surfaces", "",
    ("Observations:", True),
    "  * SRCNN/VDSR: Blurry, significant detail loss.",
    "  * SRGAN: Sharp but introduces hallucinated artifacts.",
    "  * SwinIR: Good but misses some repetitive textures.",
    "  * MambaIR: Directional scanning artifacts visible.", "",
    ("SRMamba-T (Ours):", True),
    "  * Sharpest reconstruction quality.",
    "  * No directional artifacts (thanks to MDSSM).",
    "  * Best preservation of repetitive structures.",
    "  * Clean edges without hallucination.",
], "visual_compare.png", AI)

# ════════════════════════════════════════════════
# 25. LAM + ABLATION
# ════════════════════════════════════════════════
p2("LAM Analysis & Ablation Study", [
    ("LAM = Local Attribution Map:", True),
    "Visualizes which input pixels contributed most",
    "to each output pixel during reconstruction.", "",
    ("Diffusion Index (DI) Comparison:", True),
    "  * CNN models: Low DI — very local utilization.",
    "  * SwinIR: Medium DI — window-size limited.",
    "  * MambaIR: Higher DI but directional bias.",
    "  * SRMamba-T: Highest DI — truly global.", "",
    ("Ablation Study — Layer Ordering:", True),
    "  A) All Mamba layers -> Poor global context. -0.15dB.",
    "  B) All Transformer layers -> Too expensive. OOM.",
    "  C) Trans first, Mamba later -> Suboptimal. -0.09dB.",
    "  D) Mamba first, Trans later -> BEST (our design).", "",
    ("Memory Scaling:", True),
    "At 1024x1024: SwinIR runs OOM, SRMamba-T runs fine.",
    "Near-linear memory growth vs quadratic for Transformer.",
], "slide_26.png")

# ════════════════════════════════════════════════
# 26. CONCLUSION
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Conclusion & Future Work"); print(f"  {n}. Conclusion")
txt(s, [
    ("Key Contributions:", True),
    "  1. First hybrid Mamba-Transformer architecture for image SR.",
    "  2. MDSSM: 4-directional channel-split scanning preserves 2D",
    "     spatial context at zero additional parameter cost.",
    "  3. Asymmetric placement: Mamba (early, dense extraction) +",
    "     Transformer (late, global refinement) = optimal tradeoff.",
    "  4. +0.28 dB PSNR improvement with 38.7% fewer MACs.",
    "  5. Dual-domain loss (pixel + FFT) prevents high-freq artifacts.",
    "  6. Linear memory scaling enables super-resolving large images.",
], Inches(0.4), Inches(0.95), Inches(6.2), Inches(6.2), sz=14, boxed=True)
txt(s, [
    ("Future Work:", True),
    "  * Extend to Hyperspectral SR (PaviaC, PaviaU datasets).",
    "  * Scale to SRMamba-T-L for classical (non-lightweight) SR.",
    "  * Investigate real-world degradation beyond bicubic.",
    "  * Deploy on edge devices leveraging linear memory.", 
    "  * Combine with diffusion models for perceptual quality.",
    "  * Multi-scale MDSSM with adaptive direction selection.",
    "  * Video super resolution with temporal Mamba layers.",
], Inches(6.8), Inches(0.95), Inches(6.2), Inches(6.2), sz=14, boxed=True)

# ════════════════════════════════════════════════
# 27. REFERENCES
# ════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "References"); print(f"  {n}. References")
txt(s, [
    "[1]  Dong et al., \"Image Super-Resolution Using Deep Convolutional Networks\" (SRCNN), ECCV 2014.",
    "[2]  Dong et al., \"Accelerating the Super-Resolution Convolutional Neural Network\" (FSRCNN), ECCV 2016.",
    "[3]  Kim et al., \"Accurate Image Super-Resolution Using Very Deep Convolutional Networks\" (VDSR), CVPR 2016.",
    "[4]  Ledig et al., \"Photo-Realistic Single Image Super-Resolution Using a GAN\" (SRGAN), CVPR 2017.",
    "[5]  Wang et al., \"ESRGAN: Enhanced Super-Resolution Generative Adversarial Networks\", ECCV Workshop 2018.",
    "[6]  Liang et al., \"SwinIR: Image Restoration Using Swin Transformer\", ICCV 2021.",
    "[7]  Guo et al., \"MambaIR: A Simple Baseline for Image Restoration with State Space Model\", ECCV 2024.",
    "[8]  Gu & Dao, \"Mamba: Linear-Time Sequence Modeling with Selective State Spaces\", NeurIPS 2023.",
    "[9]  Liu et al., \"SRMamba-T: Hybrid Mamba-Transformer for Single Image Super Resolution\", Neurocomputing 2025.",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017.",
], Inches(0.4), Inches(0.95), Inches(12.5), Inches(6.2), sz=14, boxed=True)

# ════════════════════════════════════════════════
# 28. THANK YOU
# ════════════════════════════════════════════════
s = S(); print(f"  {n}. Thank You")
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
r.fill.solid(); r.fill.fore_color.rgb = HDR; r.line.fill.background()
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9), Inches(0), Inches(4.33), Inches(5))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = ACC; tri.line.fill.background()
for row in range(3):
    for col in range(3):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5+col*0.8), Inches(3.5+row*0.65), Inches(0.6), Inches(0.5))
        sq.line.color.rgb = RGBColor(0x55,0x55,0x55); sq.line.width = Pt(1); sq.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.5), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = ACC; ln.line.fill.background()
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(1.2))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = ACC; p.font.name = FONT
tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WH; p2.font.name = FONT
tb3 = s.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(7), Inches(0.5))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; p3.text = "Vujja Punith Sai  |  S20240010260  |  Instructor: Arun PV"
p3.font.size = Pt(16); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p3.font.name = FONT

# ════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v9_PROFESSOR.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
