"""
Professional SRMamba-T BTP Presentation
- Custom template with consistent theme
- All text typed properly (no screenshot pasting)
- Key architecture diagrams placed neatly
- Logical flow: Literature Review → Our Method → Results
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
IMG_DIR = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")
FRIEND_DIR = os.path.join(BASE, "_temp_scripts", "slide_images_work", "friend_slides")

# ── Theme Colors ──
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT    = RGBColor(0x00, 0x96, 0xC7)  # Bright blue accent
ACCENT2   = RGBColor(0xE8, 0x5D, 0x04)  # Orange accent
GOLD      = RGBColor(0xFF, 0xBE, 0x0B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_BODY  = RGBColor(0x33, 0x3D, 0x4D)
TEXT_SUB   = RGBColor(0x6B, 0x7B, 0x8D)
DIVIDER   = RGBColor(0x00, 0x77, 0xB6)

W = Inches(13.333)
H = Inches(7.5)

def img_path(name):
    p = os.path.join(IMG_DIR, name)
    return p if os.path.exists(p) else None

def friend_img(name):
    p = os.path.join(FRIEND_DIR, name)
    return p if os.path.exists(p) else None

# ── Reusable Slide Builders ──

def add_bg(slide, color=WHITE):
    """Full slide background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_header_bar(slide, title_text, subtitle_text=None):
    """Dark header bar with title"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), W, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT
        p2.font.italic = True

def add_content_box(slide, bullets, left, top, width, height, font_size=16, bold_first=False):
    """Add a text box with bullet points"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(bullet, tuple):
            # (text, level, is_bold)
            p.text = bullet[0]
            p.level = bullet[1]
            p.font.bold = bullet[2] if len(bullet) > 2 else False
        else:
            p.text = bullet
        
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(6)
        
        if bold_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = TEXT_DARK
            p.font.size = Pt(font_size + 2)

def add_section_divider(slide, section_num, section_title, subtitle=""):
    """Section divider slide"""
    add_bg(slide, NAVY)
    
    # Section number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.5), Inches(2.2), Inches(2.2))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(2.2), Inches(1.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p3 = tf2.add_paragraph()
        p3.text = subtitle
        p3.font.size = Pt(18)
        p3.font.color.rgb = ACCENT
        p3.alignment = PP_ALIGN.CENTER

def add_two_col_header(slide, left_title, right_title, y=Inches(1.5)):
    """Two column sub-headers"""
    for (text, x) in [(left_title, Inches(0.6)), (right_title, Inches(7.0))]:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.fill.background()
        tf = box.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_image_safe(slide, path, left, top, width, height=None):
    if path and os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width)

# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 1: TITLE SLIDE                                       ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 1: Title")
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)

# Decorative accent bar on left
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), H)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Title
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Exploration of Deep Learning Models"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.text = "for Image Super Resolution"
p2.font.size = Pt(44)
p2.font.bold = True
p2.font.color.rgb = ACCENT

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.8), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

# Details
details = [
    ("B.Tech Project Presentation", Pt(20), GOLD, True),
    ("", Pt(10), WHITE, False),
    ("Presented by:  Vujja Punith Sai", Pt(22), WHITE, False),
    ("Roll No:  S20240010260", Pt(18), TEXT_SUB, False),
    ("Instructor:  Arun PV", Pt(18), TEXT_SUB, False),
]
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(3))
tf2 = txBox2.text_frame
for i, (text, size, color, bold) in enumerate(details):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(4)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 2: CONTENTS / OUTLINE                                ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 2: Contents")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "Contents")

contents_left = [
    ("Part I — Literature Review", 0, True),
    ("1.  Introduction to Super Resolution", 0, False),
    ("2.  Problem Statement", 0, False),
    ("3.  Traditional SR Methods", 0, False),
    ("4.  Deep Learning in SR", 0, False),
    ("5.  CNN-based Models (SRCNN, FSRCNN, VDSR)", 0, False),
    ("6.  GAN-based Models (SRGAN, ESRGAN)", 0, False),
    ("7.  Transformer & State Space Models", 0, False),
    ("8.  Evaluation Metrics", 0, False),
]
contents_right = [
    ("Part II — Our Method: SRMamba-T", 0, True),
    ("9.   SRMamba-T Architecture Overview", 0, False),
    ("10.  Mamba Layer & MDSSM", 0, False),
    ("11.  Transformer Layer & FFM", 0, False),
    ("12.  Reconstruction & Loss Function", 0, False),
    ("13.  Quantitative Results", 0, False),
    ("14.  Qualitative Validation & LAM", 0, False),
    ("15.  Hardware Efficiency & Ablation", 0, False),
    ("16.  Conclusion & Future Work", 0, False),
]
add_content_box(slide, contents_left, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.5), font_size=17)
add_content_box(slide, contents_right, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.5), font_size=17)

# Vertical divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.6), Inches(0.03), Inches(5))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT
div.line.fill.background()


# ╔══════════════════════════════════════════════════════════════╗
# ║  SECTION DIVIDER: PART I                                    ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 3: Section Divider - Part I")
slide = prs.slides.add_slide(blank)
add_section_divider(slide, "I", "Literature Review", "Background & Related Work")


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 4: INTRODUCTION                                      ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 4: Introduction")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "1. Introduction to Super Resolution")

left_bullets = [
    ("What is Super Resolution?", 0, True),
    ("Super Resolution (SR) uses Deep Learning to reconstruct a high-resolution image from one or more low-resolution inputs.", 0, False),
    ("", 0, False),
    ("It aims to recover:", 0, True),
    ("•  Lost image details & fine textures", 0, False),
    ("•  Sharp edges and structural information", 0, False),
    ("•  Overall visual quality", 0, False),
    ("", 0, False),
    ("Basic Workflow:", 0, True),
    ("LR Image → Feature Extraction → SR Model → Reconstruction → HR Image", 0, False),
]

right_bullets = [
    ("Real-World Applications:", 0, True),
    ("•  Satellite Imaging — terrain analysis, environmental monitoring", 0, False),
    ("•  Surveillance — face recognition, number plate enhancement", 0, False),
    ("•  Medical Imaging — MRI/CT scan enhancement", 0, False),
    ("•  Mobile Cameras — AI zoom, night photography", 0, False),
    ("•  Video Upscaling — HD to 4K conversion", 0, False),
    ("•  Remote Sensing — land use, geological mapping", 0, False),
    ("", 0, False),
    ("Industry Usage:", 0, True),
    ("Samsung, Google, Apple use SR in camera systems, satellite processing, and video restoration.", 0, False),
]

add_two_col_header(slide, "What is Super Resolution?", "Real-World Applications")
add_content_box(slide, left_bullets, Inches(0.6), Inches(2.2), Inches(5.8), Inches(5), font_size=15)
add_content_box(slide, right_bullets, Inches(7.0), Inches(2.2), Inches(5.8), Inches(5), font_size=15)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 5: PROBLEM STATEMENT                                 ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 5: Problem Statement")
slide = prs.slides.add_slide(blank)
add_bg(slide, LIGHT_BG)
add_header_bar(slide, "2. Problem Statement")

# Main problem box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(11.3), Inches(2.5))
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = ACCENT
box.line.width = Pt(2)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Most open-source data have coarse spatial resolution, and it is advantageous "
          "to use different sensors for effective enhancement; however, the sensors have "
          "differences in spatial resolution, spectral characteristics, noise levels, and "
          "often lack spatial or temporal overlap.")
p.font.size = Pt(22)
p.font.color.rgb = TEXT_DARK
p.line_spacing = Pt(34)

# Key challenges
challenges = [
    ("Key Challenges in SISR:", 0, True),
    ("•  Recovering lost high-frequency details from degraded inputs", 0, False),
    ("•  Balancing reconstruction quality vs. computational cost", 0, False),
    ("•  Maintaining global context while preserving local pixel fidelity", 0, False),
    ("•  Existing CNN models have limited receptive fields", 0, False),
    ("•  Transformer models suffer from quadratic computational complexity", 0, False),
    ("•  Pure Mamba models lose 2D spatial context due to 1D sequential scanning", 0, False),
]
add_content_box(slide, challenges, Inches(1), Inches(4.6), Inches(11.3), Inches(2.7), font_size=17)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 6: TRADITIONAL METHODS                               ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 6: Traditional Methods")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "3. Traditional Super Resolution Methods")

methods = [
    ("Before deep learning, SR relied on interpolation-based image scaling methods:", 0, False),
]
add_content_box(slide, methods, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6), font_size=17)

# Three method boxes
method_data = [
    ("Nearest Neighbour", ["Copies nearest pixel value", "Fast but blocky appearance", "Jagged edges, poor quality"]),
    ("Bilinear Interpolation", ["Uses neighbouring pixels", "Smoother than nearest", "Loss of sharpness, blurred edges"]),
    ("Bicubic Interpolation", ["Uses 16 surrounding pixels", "Better edge continuity", "Still cannot reconstruct missing details"]),
]

for i, (title, points) in enumerate(method_data):
    x = Inches(0.6 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT if i == 2 else RGBColor(0xDD, 0xDD, 0xDD)
    box.line.width = Pt(1)
    
    # Title
    tbox = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.5))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    bullets = [(f"•  {p}", 0, False) for p in points]
    add_content_box(slide, bullets, x + Inches(0.2), Inches(3.2), Inches(3.4), Inches(3), font_size=15)

# Bottom note
note = [("Limitation: All interpolation methods fail to reconstruct truly missing high-frequency details → Need for Deep Learning", 0, True)]
add_content_box(slide, note, Inches(0.6), Inches(7.0), Inches(12), Inches(0.5), font_size=15)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 7: DEEP LEARNING IN SR + EVOLUTION                   ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 7: Deep Learning + Evolution")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "4. Deep Learning in Super Resolution")

add_two_col_header(slide, "Why Deep Learning?", "Evolution of SR Architectures")

left_content = [
    ("Traditional methods cannot recover missing information.", 0, False),
    ("Deep Learning models learn the LR→HR mapping directly from large datasets.", 0, False),
    ("", 0, False),
    ("What DL Models Learn:", 0, True),
    ("•  Texture patterns (roads, vegetation, faces)", 0, False),
    ("•  Edge information (boundaries, transitions)", 0, False),
    ("•  Spatial structures (relationships between regions)", 0, False),
    ("•  Hierarchical features (low-level → high-level)", 0, False),
]

right_content = [
    ("1. CNN-based → SRCNN, FSRCNN, VDSR", 0, False),
    ("     Local feature extraction, fast but limited receptive field", 0, False),
    ("2. GAN-based → SRGAN, ESRGAN", 0, False),
    ("     Perceptual realism, but training instability", 0, False),
    ("3. Transformer-based → SwinIR", 0, False),
    ("     Global attention, but O(n²) complexity", 0, False),
    ("4. State Space Models → SR-Mamba", 0, False),
    ("     Linear complexity, but 1D scan loses 2D context", 0, False),
    ("5. Hybrid → SRMamba-T (Ours)", 0, True),
    ("     Best of Mamba + Transformer!", 0, True),
]

add_content_box(slide, left_content, Inches(0.6), Inches(2.2), Inches(5.8), Inches(5), font_size=15)
add_content_box(slide, right_content, Inches(7.0), Inches(2.2), Inches(5.8), Inches(5), font_size=15)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 8: CNN MODELS                                        ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 8: CNN Models")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "5. CNN-based Models for Super Resolution")

# Three model cards
models = [
    ("SRCNN (2014)", [
        "First DL model for SR",
        "3 Conv layers: Feature Extraction → Non-Linear Mapping → Reconstruction",
        "End-to-end learning LR→HR",
        "Limitation: Shallow, limited receptive field"
    ]),
    ("FSRCNN (2016)", [
        "Faster version of SRCNN",
        "Operates on LR space directly",
        "Deconv layer for upsampling at the end",
        "5x faster than SRCNN"
    ]),
    ("VDSR (2016)", [
        "Very Deep SR (20 conv layers)",
        "Global residual learning",
        "Learns residual: HR = LR + Residual",
        "Multi-scale: handles ×2, ×3, ×4"
    ]),
]

for i, (title, points) in enumerate(models):
    x = Inches(0.4 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(4.0), Inches(5.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)
    
    tbox = slide.shapes.add_textbox(x + Inches(0.3), Inches(1.7), Inches(3.4), Inches(0.5))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT2
    
    bullets = [(f"•  {p}", 0, False) for p in points]
    add_content_box(slide, bullets, x + Inches(0.3), Inches(2.5), Inches(3.4), Inches(4), font_size=14)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 9: GAN + TRANSFORMER + SSM MODELS                    ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 9: GAN + Transformer + SSM")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "6. GAN, Transformer & State Space Models")

models2 = [
    ("GAN-based: SRGAN & ESRGAN", [
        "Introduced perceptual loss + adversarial training",
        "Generator produces HR, Discriminator evaluates realism",
        "ESRGAN: Residual-in-Residual Dense Blocks (RRDB)",
        "Sharper textures, high-frequency detail recovery",
        "Limitation: Training instability, mode collapse"
    ], ACCENT2),
    ("Transformer: SwinIR", [
        "Uses Vision Transformer with Shifted Window Attention",
        "Global feature learning via self-attention (Q, K, V)",
        "Excellent texture reconstruction quality",
        "Limitation: O(n²) complexity, high GPU memory",
        "Limitation: Struggles with very large input resolutions"
    ], ACCENT),
    ("State Space: SR-Mamba", [
        "Based on Mamba / S6 selective state space models",
        "Linear O(n) complexity — very efficient",
        "Long-range dependency via state transitions",
        "Limitation: 1D sequential scan destroys 2D spatial context",
        "Limitation: Context loss at pixel boundaries"
    ], RGBColor(0x2D, 0xA0, 0x5B)),
]

for i, (title, points, color) in enumerate(models2):
    x = Inches(0.4 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(4.0), Inches(5.7))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    # Colored header inside box
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(4.0), Inches(0.6))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    
    tbox = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.55), Inches(3.6), Inches(0.5))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    bullets = [(f"•  {p}", 0, False) for p in points]
    add_content_box(slide, bullets, x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(4.5), font_size=14)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 10: EVALUATION METRICS                               ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 10: Evaluation Metrics")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "7. Evaluation Metrics")

metrics = [
    ("MSE (Mean Square Error)", [
        "Measures pixel-wise reconstruction error",
        "Formula: MSE = (1/mn) Σ (I(i,j) - K(i,j))²",
        "Lower MSE = better reconstruction",
        "Used in SRCNN, FSRCNN, VDSR"
    ]),
    ("PSNR (Peak Signal-to-Noise Ratio)", [
        "Measures reconstruction quality in dB",
        "Formula: PSNR = 10·log₁₀(MAX²/MSE)",
        "Higher PSNR = better image fidelity",
        "Inversely related to MSE"
    ]),
    ("SSIM (Structural Similarity Index)", [
        "Measures structural & perceptual similarity",
        "Compares: luminance, contrast, structure",
        "SSIM → 1: high similarity (good)",
        "SSIM → 0: poor similarity (bad)"
    ]),
]

for i, (title, points) in enumerate(metrics):
    x = Inches(0.4 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(4.0), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)
    
    tbox = slide.shapes.add_textbox(x + Inches(0.3), Inches(1.7), Inches(3.4), Inches(0.5))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    bullets = [(f"•  {p}", 0, False) for p in points]
    add_content_box(slide, bullets, x + Inches(0.3), Inches(2.5), Inches(3.4), Inches(3.5), font_size=14)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SECTION DIVIDER: PART II                                   ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 11: Section Divider - Part II")
slide = prs.slides.add_slide(blank)
add_section_divider(slide, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution")


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 12: WHY HYBRID? MOTIVATION                           ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 12: Why Hybrid?")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "8. Why Hybrid? The Architecture Compromise", 
               "No single architecture solves all SR challenges")

add_two_col_header(slide, "The Problem with Each Approach", "The SRMamba-T Solution")

left_content = [
    ("CNN:", 0, True),
    ("  ✗  Limited receptive field (local only)", 0, False),
    ("  ✗  Cannot capture long-range dependencies", 0, False),
    ("", 0, False),
    ("Transformer:", 0, True),
    ("  ✓  Global context via self-attention", 0, False),
    ("  ✗  O(n²) complexity — GPU memory explosion", 0, False),
    ("", 0, False),
    ("Mamba (SSM):", 0, True),
    ("  ✓  O(n) linear complexity — very efficient", 0, False),
    ("  ✗  1D sequential scan destroys 2D spatial info", 0, False),
]

right_content = [
    ("Hybrid Insight: Sequence Before Space", 0, True),
    ("", 0, False),
    ("Phase 1 — Mamba (Early Layers):", 0, True),
    ("  Dense feature extraction at linear cost O(n)", 0, False),
    ("  Handles the bulk of heavy computation efficiently", 0, False),
    ("", 0, False),
    ("Phase 2 — Transformer (Later Layers):", 0, True),
    ("  Global refinement on already-compressed features", 0, False),
    ("  Un-windowed attention without memory explosion", 0, False),
    ("", 0, False),
    ("→ Best of both worlds: Efficient + Accurate", 0, True),
]

add_content_box(slide, left_content, Inches(0.6), Inches(2.2), Inches(5.8), Inches(5), font_size=15)
add_content_box(slide, right_content, Inches(7.0), Inches(2.2), Inches(5.8), Inches(5), font_size=15)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 13: ARCHITECTURE OVERVIEW (with diagram)              ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 13: Architecture Overview")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "9. SRMamba-T: Complete Architecture")

# Architecture diagram from user's slide 09
add_image_safe(slide, img_path("slide_09.png"), Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 14: MAMBA LAYER + MAMBA MIXER                        ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 14: Mamba Layer")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "10. Mamba Block: Anatomy & Asymmetric Mixer")

# Left: Mamba Layer anatomy (slide_14)
add_image_safe(slide, img_path("slide_14.png"), Inches(0.2), Inches(1.3), Inches(6.5), Inches(6.0))
# Right: Asymmetric Mixer (slide_15)
add_image_safe(slide, img_path("slide_15.png"), Inches(6.8), Inches(1.3), Inches(6.3), Inches(6.0))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 15: MDSSM                                            ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 15: MDSSM")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "11. Multi-Directional Selective Scan Module (MDSSM)")

# Left: MDSSM concept (slide_17)
add_image_safe(slide, img_path("slide_17.png"), Inches(0.2), Inches(1.3), Inches(6.5), Inches(6.0))
# Right: MDSSM equations (slide_18)
add_image_safe(slide, img_path("slide_18.png"), Inches(6.8), Inches(1.3), Inches(6.3), Inches(6.0))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 16: TRANSFORMER + FFM                                ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 16: Transformer + FFM")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "12. Transformer Layer & Feature Fusion Module")

# Left: Transformer (slide_19)
add_image_safe(slide, img_path("slide_19.png"), Inches(0.2), Inches(1.3), Inches(6.5), Inches(6.0))
# Right: FFM (slide_21)
add_image_safe(slide, img_path("slide_21.png"), Inches(6.8), Inches(1.3), Inches(6.3), Inches(6.0))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 17: RECONSTRUCTION + LOSS                            ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 17: Reconstruction + Loss")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "13. Adaptive Reconstruction & Dual-Domain Loss")

# Full width diagram (slide_22)
add_image_safe(slide, img_path("slide_22.png"), Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 18: QUANTITATIVE RESULTS                             ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 18: Results")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "14. Quantitative Results")

# Left: PSNR chart (slide_24)
add_image_safe(slide, img_path("slide_24.png"), Inches(0.2), Inches(1.3), Inches(6.5), Inches(6.0))
# Right: Hardware (slide_27)
add_image_safe(slide, img_path("slide_27.png"), Inches(6.8), Inches(1.3), Inches(6.3), Inches(6.0))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 19: QUALITATIVE + LAM                                ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 19: Qualitative + LAM")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "15. Qualitative Validation & LAM Analysis")

# Left: Visual comparison (slide_25)
add_image_safe(slide, img_path("slide_25.png"), Inches(0.2), Inches(1.3), Inches(6.5), Inches(6.0))
# Right: LAM (slide_26)
add_image_safe(slide, img_path("slide_26.png"), Inches(6.8), Inches(1.3), Inches(6.3), Inches(6.0))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 20: ABLATION STUDY                                   ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 20: Ablation")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "16. Ablation Study: Layer Ordering")

# Full width (slide_28)
add_image_safe(slide, img_path("slide_28.png"), Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 21: CONCLUSION                                       ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 21: Conclusion")
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header_bar(slide, "17. Conclusion & Future Work")

add_two_col_header(slide, "Key Contributions", "Future Work")

left_content = [
    ("✓  Proposed SRMamba-T: first hybrid Mamba-Transformer architecture for image SR", 0, False),
    ("", 0, False),
    ("✓  MDSSM: Channel-split 4-directional scanning preserves 2D spatial context at O(n) cost", 0, False),
    ("", 0, False),
    ("✓  Asymmetric layer placement: Mamba (early) + Transformer (late) achieves optimal efficiency", 0, False),
    ("", 0, False),
    ("✓  +0.28 dB PSNR improvement over MambaIR with 38.7% fewer MACs", 0, False),
    ("", 0, False),
    ("✓  Dual-domain loss (Pixel + FFT) prevents frequency-domain artifacts", 0, False),
]

right_content = [
    ("•  Extend to Hyperspectral Image Super Resolution (PaviaC, PaviaU datasets)", 0, False),
    ("", 0, False),
    ("•  Scale to larger models (SRMamba-T-L) for classical SR tasks", 0, False),
    ("", 0, False),
    ("•  Explore multi-scale U-Net with downsampling between encoder layers", 0, False),
    ("", 0, False),
    ("•  Investigate real-world degradation models beyond bicubic downsampling", 0, False),
    ("", 0, False),
    ("•  Deploy on edge devices leveraging Mamba's linear memory scaling", 0, False),
]

add_content_box(slide, left_content, Inches(0.6), Inches(2.2), Inches(5.8), Inches(5), font_size=16)
add_content_box(slide, right_content, Inches(7.0), Inches(2.2), Inches(5.8), Inches(5), font_size=16)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 22: THANK YOU                                        ║
# ╚══════════════════════════════════════════════════════════════╝
print("Slide 22: Thank You")
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), H)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.8))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Questions & Discussion"
p2.font.size = Pt(28)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11), Inches(1.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Vujja Punith Sai  |  S20240010260  |  Guide: Arun PV"
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p3.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
ppt_dir = os.path.join(BASE, "PPTs")
os.makedirs(ppt_dir, exist_ok=True)
out_path = os.path.join(ppt_dir, "SRMamba_T_BTP_Final.pptx")
prs.save(out_path)
print(f"\n>>> Saved: {out_path}")
print(f">>> Total slides: {len(prs.slides)}")
print(">>> DONE!")
