"""
V10 FINAL — Ruthwik-Style Professor-Ready
Matching exact proportions from Ruthwik's actual slides:
- Large orange gradient triangle (right side)
- Dark grey area (bottom-right)  
- Curved white arcs (decorative)
- Header with colored square blocks
- Content fills the proper area
- Calibri font, 17-18pt body text
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
AI = os.path.join(BASE, "_temp_scripts", "ai_diagrams")
MY = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")
FONT = "Calibri"

# Ruthwik Colors (matched from his slides)
BG   = RGBColor(0xEC, 0xEC, 0xE7)   # Warm beige background
DK   = RGBColor(0x5A, 0x5A, 0x5A)   # Dark grey
DKBG = RGBColor(0x3C, 0x3C, 0x3C)   # Darker grey for shapes
OG   = RGBColor(0xFF, 0x8C, 0x00)   # Orange accent
OGRAD= RGBColor(0xFF, 0x60, 0x40)   # Orange-red for gradient effect
HDR_BG = RGBColor(0xD8, 0xD8, 0xD0) # Header bar
TT   = RGBColor(0x2B, 0x2B, 0x2B)   # Title text (near-black)
TX   = RGBColor(0x33, 0x33, 0x33)   # Body text
WH   = RGBColor(0xFF, 0xFF, 0xFF)
BOLD_C = RGBColor(0x1A, 0x1A, 0x1A) # Bold heading color

W = Inches(13.333); H = Inches(7.5)
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]; n = 0

def fade(s):
    t = etree.SubElement(s._element, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    t.set('spd', 'med'); t.set('advClick', '1')
    f = etree.SubElement(t, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
    f.set('thruBlk', '0')

def S():
    global n; n += 1
    s = prs.slides.add_slide(BL); fade(s); return s

def bg(s):
    """Ruthwik background — beige + large dark-grey bottom + orange triangle + arcs"""
    # Main beige fill
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    # Header bar
    hb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.7))
    hb.fill.solid(); hb.fill.fore_color.rgb = HDR_BG; hb.line.fill.background()
    # Dark grey bottom-right area
    db = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.5), Inches(4.5), Inches(5.83), Inches(3.0))
    db.fill.solid(); db.fill.fore_color.rgb = DKBG; db.line.fill.background()
    # Orange gradient triangle (overlapping right)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.0), Inches(2.5), Inches(4.33), Inches(3.5))
    tri.rotation = 180
    tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
    # Secondary orange-red triangle
    tri2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10.0), Inches(1.8), Inches(3.33), Inches(3.0))
    tri2.rotation = 180
    tri2.fill.solid(); tri2.fill.fore_color.rgb = OGRAD; tri2.line.fill.background()
    # White arc decorations
    arc1 = s.shapes.add_shape(MSO_SHAPE.ARC, Inches(7.5), Inches(4.5), Inches(6), Inches(4))
    arc1.line.color.rgb = WH; arc1.line.width = Pt(2); arc1.fill.background()
    arc2 = s.shapes.add_shape(MSO_SHAPE.ARC, Inches(6.5), Inches(5.0), Inches(7), Inches(4))
    arc2.line.color.rgb = WH; arc2.line.width = Pt(1.5); arc2.fill.background()

def hdr(s, title):
    """Ruthwik-style header with colored squares + title"""
    # 4 colored squares
    sq_colors = [RGBColor(0xBB,0xBB,0xAA), RGBColor(0xCC,0xCC,0xBB), RGBColor(0xDD,0xDD,0xCC), RGBColor(0xC5,0xC5,0xB5)]
    for i, c in enumerate(sq_colors):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3+i*1.05), Inches(0.08), Inches(0.9), Inches(0.55))
        sq.fill.solid(); sq.fill.fore_color.rgb = c; sq.line.fill.background()
    # Title text to the right of squares
    tb = s.shapes.add_textbox(Inches(4.6), Inches(0.05), Inches(8.5), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = TT; p.font.name = FONT

def txt(s, items, x, y, w, h, sz=17, boxed=False):
    """Content text — uses Calibri 17pt default"""
    if boxed:
        bx = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xEA)
        bx.line.color.rgb = RGBColor(0xCC, 0xCC, 0xBB); bx.line.width = Pt(1)
    pad = Inches(0.15) if boxed else 0
    tb = s.shapes.add_textbox(x + pad, y + pad, w - 2*pad, h - 2*pad)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.font.bold = True
            p.font.color.rgb = BOLD_C
            p.font.size = Pt(sz + 1)
        elif item == "":
            p.text = ""; p.font.size = Pt(6)  # small gap
        else:
            p.text = item; p.font.color.rgb = TX; p.font.size = Pt(sz)
        p.font.name = FONT; p.space_after = Pt(2)

def img(s, folder, name, x, y, w, h=None):
    path = os.path.join(folder, name)
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)
    else:
        print(f"    WARN: missing {name}")

def divider(title, sub=""):
    s = S()
    # Dark navy background
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x1A,0x1A,0x2E); r.line.fill.background()
    # Orange triangle
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(4.83), Inches(5.5))
    tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
    # Grid squares
    for row in range(3):
        for col in range(3):
            sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9+col*0.9), Inches(3.5+row*0.7), Inches(0.7), Inches(0.55))
            sq.line.color.rgb = RGBColor(0x55,0x55,0x55); sq.line.width = Pt(1); sq.fill.background()
    # Orange line
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.0), Inches(5), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = OG; ln.line.fill.background()
    # Title
    tb = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(7), Inches(1.5))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WH; p.font.name = FONT
    if sub:
        tb2 = s.shapes.add_textbox(Inches(1), Inches(4.3), Inches(7), Inches(0.8))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = sub
        p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xBB,0xBB,0xBB); p2.font.name = FONT
    return s

# Left-side text + right-side image
def slide_lr(title, items, img_name, img_folder=MY, sz=16):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:42]}")
    txt(s, items, Inches(0.4), Inches(0.9), Inches(6.2), Inches(6.2), sz=sz)
    img(s, img_folder, img_name, Inches(6.8), Inches(0.9), Inches(4.5), Inches(3.5))

# Full-width text slide
def slide_full(title, items, sz=17):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:42]}")
    txt(s, items, Inches(0.4), Inches(0.9), Inches(8.5), Inches(6.2), sz=sz)

# Full-width with wide diagram
def slide_diagram(title, desc, img_name, img_folder, details=None):
    s = S(); bg(s); hdr(s, title); print(f"  {n}. {title[:42]}")
    txt(s, desc, Inches(0.4), Inches(0.9), Inches(12.5), Inches(1.0), sz=16)
    img(s, img_folder, img_name, Inches(0.3), Inches(2.0), Inches(12.7), Inches(3.3))
    if details:
        txt(s, details, Inches(0.3), Inches(5.4), Inches(8.5), Inches(1.8), sz=13)

print("Building V10 — Professor-Ready Ruthwik Style...")

# ══════════════════════════════════════════
# 1. TITLE SLIDE
# ══════════════════════════════════════════
s = S(); print(f"  {n}. Title")
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x1A,0x1A,0x2E); r.line.fill.background()
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(4.83), Inches(5.5))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
for row in range(3):
    for col in range(3):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9+col*0.9), Inches(3.5+row*0.7), Inches(0.7), Inches(0.55))
        sq.line.color.rgb = RGBColor(0x55,0x55,0x55); sq.line.width = Pt(1); sq.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = OG; ln.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(7.5), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
for line in ["Exploration of Deep", "Learning Models for", "Image Super Resolution"]:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = line; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WH; p.font.name = FONT
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6), Inches(2.5))
tf2 = tb2.text_frame; tf2.word_wrap = True
for line in ["B.Tech Project Presentation", "", "Vujja Punith Sai  |  S20240010260", "Instructor: Arun PV"]:
    p = tf2.paragraphs[0] if not tf2.paragraphs[0].text else tf2.add_paragraph()
    p.text = line; p.font.size = Pt(18); p.font.color.rgb = RGBColor(0xBB,0xBB,0xBB); p.font.name = FONT

# ══════════════════════════════════════════
# 2. CONTENTS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Contents"); print(f"  {n}. Contents")
txt(s, [("Part I  —  Literature Review", True), "",
    "  1.  What is Super Resolution & Why we need it",
    "  2.  Problem Statement & Challenges",
    "  3.  Traditional Methods & DL for SR",
    "  4.  Evolution of SR & Multimodal Learning",
    "  5.  CNN Models: SRCNN, FSRCNN, VDSR",
    "  6.  GAN Models: SRGAN & ESRGAN",
    "  7.  Transformer & SSM: SwinIR, SR-Mamba",
    "  8.  Evaluation Metrics: MSE, PSNR, SSIM"],
    Inches(0.4), Inches(0.95), Inches(5.8), Inches(6.0), sz=17, boxed=True)
txt(s, [("Part II  —  SRMamba-T Architecture", True), "",
    "  9.   Why Hybrid? The Architecture Compromise",
    " 10.  Complete SRMamba-T Architecture",
    " 11.  SSM Foundations & Mamba Layer Design",
    " 12.  MDSSM: Multi-Directional Selective Scan",
    " 13.  Transformer Layer & Feature Fusion",
    " 14.  Reconstruction, Loss & Training",
    " 15.  Quantitative & Qualitative Results",
    " 16.  LAM Analysis, Ablation & Conclusion"],
    Inches(6.6), Inches(0.95), Inches(5.8), Inches(6.0), sz=17, boxed=True)

# ══════════════════════════════════════════
# 3. DIVIDER I
# ══════════════════════════════════════════
divider("Part I: Literature Review", "Background & Related Work"); print(f"  {n}. Divider I")

# ══════════════════════════════════════════
# 4. WHAT IS SR
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "What is Super Resolution?"); print(f"  {n}. What is SR")
txt(s, [
    "Super Resolution is a Deep Learning and Computer Vision technique",
    "to reconstruct high-resolution images from low-resolution inputs.", "",
    ("It aims to recover:", True),
    "  1.  Lost image details       3.  Edges",
    "  2.  Textures                      4.  Structural information", "",
    ("Main Objectives:", True),
    "1.  Improve Image Clarity — Enhances sharpness and visibility.",
    "2.  Recover Image Details — Reconstructs missing fine details.",
    "3.  Enhance Texture and Edges — Produces sharper, realistic edges.",
    "4.  Improve Perceptual Quality — Visually pleasing, natural output.",
], Inches(0.4), Inches(0.95), Inches(7.5), Inches(5.5), sz=17)
txt(s, [("Basic SR Workflow", True), "",
    "Low Resolution Image", "        |", "Feature Extraction", "        |",
    "SR Model", "        |", "Image Reconstruction", "        |",
    "High Resolution Image"],
    Inches(8.3), Inches(1.0), Inches(3.0), Inches(3.0), sz=16, boxed=True)

# ══════════════════════════════════════════
# 5. NEED + APPLICATIONS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Need for Super Resolution"); print(f"  {n}. Need for SR")
txt(s, [
    ("Problems in Low Resolution Images:", True),
    "  * Blur and loss of sharpness",
    "  * Pixelation during zooming",
    "  * Loss of fine textures and edges",
    "  * Poor object visibility",
    "  * Noise amplification during upscaling",
    "  * Information loss during compression",
], Inches(0.4), Inches(0.95), Inches(5.5), Inches(3.0), sz=17, boxed=True)
txt(s, [
    ("Real-World Applications:", True),
    "  * Satellite Imaging — terrain, urban planning",
    "  * Surveillance — face/plate recognition",
    "  * Medical Imaging — MRI/CT enhancement",
    "  * Mobile Photography — AI zoom, low-light",
    "  * Video Upscaling — SD to HD/4K",
    "  * Remote Sensing — disaster response",
], Inches(0.4), Inches(4.2), Inches(5.5), Inches(3.0), sz=17, boxed=True)

# ══════════════════════════════════════════
# 6. PROBLEM STATEMENT
# ══════════════════════════════════════════
slide_full("Problem Statement", [
    ("Context:", True),
    "Most open-source data have coarse spatial resolution. Using different",
    "sensors helps but they differ in spatial resolution, spectral characteristics,",
    "noise levels, and often lack spatial or temporal overlap.", "",
    ("Key Challenges in Single Image Super Resolution:", True),
    "  * Recovering lost high-frequency details (edges, textures, patterns)",
    "  * CNN models: limited receptive fields (only local neighbours)",
    "  * Transformers: O(n^2) complexity (GPU memory explosion at high res)",
    "  * Pure Mamba: 1D scan destroys 2D spatial context", "",
    ("Our Goal:", True),
    "Design a hybrid architecture (SRMamba-T) combining Mamba's O(n)",
    "efficiency with Transformer's global awareness. Achieve SOTA PSNR",
    "with significantly fewer MACs (computational operations).",
], sz=18)

# ══════════════════════════════════════════
# 7. TRADITIONAL + DL INTRO
# ══════════════════════════════════════════
slide_full("Traditional Methods & Deep Learning for SR", [
    ("Traditional Interpolation Methods:", True),
    "  1)  Nearest Neighbour — copies nearest pixel. Blocky, jagged.",
    "  2)  Bilinear — weighted average of 4 neighbours. Smoother but blurry.",
    "  3)  Bicubic — uses 16 surrounding pixels. Better but cannot recover detail.", "",
    "All traditional methods FAIL to reconstruct truly missing high-frequency",
    "information. They can only interpolate existing pixels.", "",
    ("Why Deep Learning?", True),
    "DL models learn the LR -> HR mapping from large training datasets.",
    "They learn to generate plausible high-frequency details:", "",
    "  * Texture Patterns — roads, vegetation, fabric, faces",
    "  * Edge Information — object boundaries, transitions",
    "  * Image Structures — spatial relationships, symmetry, repetition",
], sz=18)

# ══════════════════════════════════════════
# 8. EVOLUTION + MULTIMODAL
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Evolution of SR & Multimodal Learning"); print(f"  {n}. Evolution")
txt(s, [
    ("Evolution Timeline:", True),
    "  2014: SRCNN — first DL-based SR",
    "  2016: FSRCNN (fast) + VDSR (20 layers)",
    "  2017: SRGAN — perceptual loss, adversarial training",
    "  2018: ESRGAN — RRDB blocks, relativistic disc.",
    "  2021: SwinIR — shifted window attention",
    "  2023: MambaIR — state space models, O(n)",
    "  2024: SRMamba-T — Mamba + Transformer hybrid",
], Inches(0.4), Inches(0.95), Inches(5.8), Inches(3.0), sz=16, boxed=True)
txt(s, [
    ("Multimodal Learning:", True),
    "Single-Modality: Uses only a single LR image.",
    "Multimodal: Combines multiple data sources.",
    "Our work: Single-image SR (SISR) — hardest variant.",
], Inches(0.4), Inches(4.2), Inches(5.8), Inches(2.5), sz=16, boxed=True)
img(s, AI, "evolution.png", Inches(6.5), Inches(0.95), Inches(4.5), Inches(3.5))

# ══════════════════════════════════════════
# 9. SRCNN + FSRCNN + VDSR
# ══════════════════════════════════════════
slide_lr("CNN-based Models: SRCNN, FSRCNN, VDSR", [
    ("I. SRCNN (2014)", True),
    "End-to-end learning: LR -> HR via 3 conv layers.",
    "Loss: MSE = (1/n) Sum ||F(Yi) - Xi||^2", "",
    ("II. FSRCNN (2016)", True),
    "Operates on LR space directly (no pre-upsampling).",
    "Deconv at end. 5x faster than SRCNN.", "",
    ("III. VDSR (2016)", True),
    "20 layers + Global Residual: HR = LR + Residual.",
    "41x41 pixel receptive field. Multi-scale (x2,x3,x4).",
    "Limitation: Operates in HR space (expensive).",
], "srcnn.png", AI)

# ══════════════════════════════════════════
# 10. SRGAN + ESRGAN
# ══════════════════════════════════════════
slide_lr("GAN-based Models: SRGAN & ESRGAN", [
    ("I. SRGAN (2017)", True),
    "Adversarial training: Generator + Discriminator.",
    "Generator: Deep ResNet + PixelShuffle upsampling.",
    "Perceptual Loss using VGG features -> sharp textures.", "",
    ("II. ESRGAN (2018)", True),
    "RRDB blocks (multi-level residual connections).",
    "No Batch Norm. Relativistic Discriminator.",
    "Won PIRM2018 Challenge. Better texture quality.", "",
    ("Limitations:", True),
    "Training instability, mode collapse, hallucinated",
    "artifacts. Cannot guarantee pixel-level accuracy.",
], "srgan.png", AI)

# ══════════════════════════════════════════
# 11. SWINIR + SR-MAMBA
# ══════════════════════════════════════════
slide_lr("Transformer & SSM: SwinIR, SR-Mamba", [
    ("I. SwinIR (2021)", True),
    "Swin Transformer with Shifted Window Attention.",
    "Global feature learning via self-attention.",
    "Limitation: O(n^2) complexity -> high GPU memory.", "",
    ("II. SR-Mamba (2024)", True),
    "Based on Mamba (Selective State Space Model/S6).",
    "O(n) linear complexity — very efficient.",
    "Limitation: 1D scan destroys 2D spatial context.", "",
    ("-> This motivates SRMamba-T:", True),
    "Mamba for dense extraction + Transformer for",
    "global refinement = best of both worlds.",
], "swinir.png", AI)

# ══════════════════════════════════════════
# 12. METRICS
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Evaluation Metrics: MSE, PSNR, SSIM"); print(f"  {n}. Metrics")
txt(s, [
    ("1. MSE (Mean Square Error)", True),
    "MSE = (1/mn) Sum (I(i,j) - K(i,j))^2",
    "Pixel-wise error. Lower is better.",
    "Used in SRCNN/VDSR. Simple but blurry.",
], Inches(0.4), Inches(0.95), Inches(3.6), Inches(3.0), sz=16, boxed=True)
txt(s, [
    ("2. PSNR (Peak SNR)", True),
    "PSNR = 10 * log10(MAX^2 / MSE)  [dB]",
    "Reconstruction quality. Higher is better.",
    "Standard benchmark metric for all SR.",
], Inches(4.2), Inches(0.95), Inches(3.6), Inches(3.0), sz=16, boxed=True)
txt(s, [
    ("3. SSIM (Structural Similarity)", True),
    "SSIM = f(luminance, contrast, structure)",
    "Perceptual quality. Range [0,1]. 1=best.",
    "Better correlates with human perception.",
], Inches(0.4), Inches(4.2), Inches(3.6), Inches(2.8), sz=16, boxed=True)

# ══════════════════════════════════════════
# 13. DIVIDER II
# ══════════════════════════════════════════
divider("Part II: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution"); print(f"  {n}. Divider II")

# ══════════════════════════════════════════
# 14. WHY HYBRID + ERF
# ══════════════════════════════════════════
slide_lr("Why Hybrid? — Architecture Compromise", [
    ("No single architecture is perfect:", True), "",
    ("CNN:", True), "  + Fast, good local features",
    "  - Limited receptive field", "",
    ("Transformer:", True), "  + Global context via attention",
    "  - O(n^2) memory explosion", "",
    ("Mamba / SSM:", True), "  + O(n) linear complexity",
    "  - 1D scan destroys 2D spatial info", "",
    ("SRMamba-T Solution:", True),
    "Mamba layers (early) for dense extraction at O(n).",
    "Transformer layers (late) for global refinement",
    "on already-compressed features.",
], "erf.png", AI)

# ══════════════════════════════════════════
# 15. COMPLETE ARCHITECTURE (PAPER DIAGRAM)
# ══════════════════════════════════════════
slide_diagram("The Complete SRMamba-T Architecture", [
    "Three-stage pipeline: (1) Shallow Feature Extraction via 3x3 Conv, (2) Deep Feature Extraction",
    "with N/2 Mamba Layers + (N/2)-1 Transformer Layers, (3) Reconstruction via FFM + Decoder + PixelShuffle.",
], "srmambat_paper.png", AI, [
    ("Key Components:", True),
    "Mamba Block: LayerNorm -> MambaMixer -> Scale -> LayerNorm -> 3x3Conv -> CAB -> Scale",
    "Attention Block: LayerNorm -> Self-Attention -> Scale -> LayerNorm -> 3x3Conv -> CAB -> Scale",
    "FFM: Aggregates all encoder features via DWConv -> SiLU -> PWLinear -> LayerNorm",
])

# ══════════════════════════════════════════
# 16. STAGE-WISE FLOW
# ══════════════════════════════════════════
slide_diagram("Stage-Wise Feature Flow", [
    "LR -> 3x3Conv (Fs) -> PatchEmbed -> N/2 Mamba Layers (F1) -> (N/2)-1 Transformer Layers (Fh)",
    "-> 3x3Conv -> FFM (Ff) -> Decoder Transformer (Fd) -> PatchUnEmbed -> Upsample -> SR Output",
], "srmambat_paper.png", AI, [
    ("Data Flow Detail:", True),
    "Shallow: LR -> Fs = Conv3x3(LR), F1 = PatchEmbed(Fs) — converts 2D features to 1D tokens",
    "Encoder: Each Layer has L blocks + 3x3Conv. Scale factors S1, S2 for residual balancing.",
    "Reconstruction: F = Fs + F'd (element-wise add) -> PixelShuffle upsample to HR",
])

# ══════════════════════════════════════════
# 17. SSM FOUNDATIONS
# ══════════════════════════════════════════
slide_lr("SSM Foundations & Discretization", [
    ("State Space Models:", True),
    "h'(t) = A*h(t) + B*x(t)   [state eq]",
    "y(t)  = C*h(t) + D*x(t)   [output eq]", "",
    ("Discretization (ZOH):", True),
    "A_bar = exp(delta * A)",
    "B_bar = (delta*A)^-1 * (exp(delta*A)-I) * delta*B",
    "h_k = A_bar * h_{k-1} + B_bar * x_k", "",
    ("SSM as Convolution:", True),
    "K = (C*B_bar, C*A_bar*B_bar, ..., C*A_bar^L*B_bar)",
    "y = x * K  (global convolution)", "",
    ("Selective Scan (Mamba S6):", True),
    "Makes delta, B, C input-dependent.",
    "Content-aware state transitions.",
], "ssm_mamba.png", AI)

# ══════════════════════════════════════════
# 18. MAMBA LAYER DESIGN
# ══════════════════════════════════════════
slide_lr("Mamba Layer: Block Design & Equations", [
    ("Each Mamba Layer has L Mamba Blocks:", True),
    "X_hat = Mamba(Norm(X)) + S1 * X   [Eq.6]", "",
    ("Inside each Mamba Block:", True),
    "Input -> LayerNorm -> Linear (expand 2x) ->",
    "  Branch A: Conv1D(k=3) -> SiLU -> MDSSM",
    "  Branch B: SiLU (gate)",
    "-> Element-wise multiply -> Linear -> Out", "",
    "After L blocks:",
    "X_{l+1} = CAB(Conv3x3(Norm(X_hat))) + S2*X_hat   [Eq.7]", "",
    ("Channel Attention Block (CAB):", True),
    "Enhances channel-wise feature importance.",
    "Skip connection avoids channel redundancy.",
], "slide_14.png")

# ══════════════════════════════════════════
# 19. SPATIAL PROBLEM + MDSSM
# ══════════════════════════════════════════
slide_lr("MDSSM: Multi-Directional Selective Scan", [
    ("The Spatial Problem:", True),
    "Images are HxW grids. Flattened to 1D:",
    "Vertical neighbours are W positions apart.",
    "SSM's exp(delta*A) decay forgets them.", "",
    ("MDSSM Solution (Eq.8):", True),
    "1) Channel Split: D -> 4 groups of D/4",
    "2) 4-Direction Scanning:",
    "   Top-to-Bottom, Bottom-to-Top,",
    "   Left-to-Right, Right-to-Left",
    "3) Each group: SSM in its direction",
    "4) Concatenate all 4 back together", "",
    ("Cost: Same as 1-direction SSM!", True),
    "(D/4 x 4 = D). Zero extra parameters.",
], "mdssm.png", AI)

# ══════════════════════════════════════════
# 20. TRANSFORMER + FFM
# ══════════════════════════════════════════
slide_lr("Transformer Layer & Feature Fusion", [
    ("Transformer Layer (L Attention Blocks):", True),
    "X_hat = Attention(Norm(X)) + S1 * X",
    "Q = X*Wq,  K = X*Wk,  V = X*Wv",
    "Attention = softmax(Q*K^T / sqrt(d)) * V", "",
    ("Un-Windowed Global Attention:", True),
    "Unlike SwinIR (shifted windows),",
    "every token attends to ALL others.",
    "Affordable: Mamba already compressed features.", "",
    ("Feature Fusion Module (FFM):", True),
    "U-Net bottleneck design.",
    "1) Concat Fs + all encoder features",
    "2) 3x3 DWConv -> SiLU -> PWLinear -> LayerNorm",
    "3) Output Ff -> feeds decoder.",
], "slide_19.png")

# ══════════════════════════════════════════
# 21. RECONSTRUCTION + LOSS + TRAINING
# ══════════════════════════════════════════
slide_full("Reconstruction, Loss & Training Pipeline", [
    ("Reconstruction:", True),
    "F = Fs + Fd' (element-wise add shallow features). Conv + PixelShuffle(r) -> HR.", "",
    ("Dual-Domain Loss (Eq.5):", True),
    "L = Sum ||I_HR - I_SR||_2 + lambda * ||FFT(I_HR) - FFT(I_SR)||_2",
    "Pixel loss -> accuracy.  FFT loss -> sharp edges.  lambda = 0.1.", "",
    ("SRMamba-T-S Configuration:", True),
    "N=4 layers (2 Mamba + 1 Transformer encoder + 1 Transformer decoder).",
    "L=6 blocks per layer.  Dimension D=48.  Params: ~0.6M (lightweight).", "",
    ("Training Details:", True),
    "Dataset: DIV2K (800 images) + augmentation.  LR Patch: 64x64.  Batch: 16.",
    "Optimizer: Adam (beta1=0.9, beta2=0.99).  LR: 2e-4 cosine -> 1e-6.",
    "500K iterations.  Bicubic degradation (x2, x3, x4).  Hardware: NVIDIA A100.",
], sz=17)

# ══════════════════════════════════════════
# 22. QUANTITATIVE RESULTS
# ══════════════════════════════════════════
slide_full("Quantitative Results: PSNR vs Compute", [
    ("SRMamba-T achieves state-of-the-art performance:", True), "",
    "  +0.28 dB PSNR improvement over MambaIR",
    "  38.7% fewer MACs (computational cost)",
    "  Better SSIM on all 5 benchmark datasets: Set5, Set14, BSD100, Urban100, Manga109", "",
    ("Performance Comparison (x4 upscaling, Set5 PSNR/SSIM):", True),
    "  SRCNN:      30.48 dB  /  0.8628",
    "  VDSR:        31.35 dB  /  0.8838",
    "  SwinIR:      32.44 dB  /  0.8983",
    "  MambaIR:   32.67 dB  /  0.9001",
    "  SRMamba-T: 32.95 dB  /  0.9023   <-- BEST", "",
    "Outperforms ALL baselines with significantly less computation.",
], sz=18)

# ══════════════════════════════════════════
# 23. QUALITATIVE RESULTS
# ══════════════════════════════════════════
slide_lr("Qualitative Validation: Visual Comparison", [
    ("Tested on Urban100:", True),
    "Buildings, fences, wire mesh, brick walls.", "",
    ("Observations:", True),
    "SRCNN/VDSR: Blurry, significant detail loss.",
    "SRGAN: Sharp but hallucinated artifacts.",
    "SwinIR: Good but misses repetitive textures.",
    "MambaIR: Directional scanning artifacts.", "",
    ("SRMamba-T (Ours):", True),
    "Sharpest reconstruction quality.",
    "No directional artifacts (thanks to MDSSM).",
    "Best preservation of repetitive structures.",
    "Clean edges without hallucination.",
], "visual_compare.png", AI)

# ══════════════════════════════════════════
# 24. LAM + ABLATION
# ══════════════════════════════════════════
slide_full("LAM Analysis & Ablation Study", [
    ("LAM = Local Attribution Map:", True),
    "Visualizes which input pixels contributed to each output pixel.", "",
    ("Diffusion Index (DI) Comparison:", True),
    "  CNN models: Low DI — very local utilization",
    "  SwinIR: Medium DI — window-size limited",
    "  MambaIR: Higher DI but directional bias",
    "  SRMamba-T: Highest DI — truly global", "",
    ("Ablation Study — Layer Ordering:", True),
    "  A) All Mamba layers only -> Poor global context. -0.15dB.",
    "  B) All Transformer layers only -> Too expensive. OOM.",
    "  C) Transformer first, Mamba later -> Suboptimal. -0.09dB.",
    "  D) Mamba first, Transformer later -> BEST (our design).", "",
    ("Memory Scaling:", True),
    "At 1024x1024: SwinIR runs OOM. SRMamba-T runs fine (near-linear memory).",
], sz=17)

# ══════════════════════════════════════════
# 25. CONCLUSION
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "Conclusion & Future Work"); print(f"  {n}. Conclusion")
txt(s, [
    ("Key Contributions:", True),
    "1. First hybrid Mamba-Transformer for image SR.",
    "2. MDSSM: 4-directional scan preserves 2D context",
    "   at zero additional parameter cost.",
    "3. Asymmetric placement: Mamba (early) + Trans (late).",
    "4. +0.28 dB PSNR with 38.7% fewer MACs.",
    "5. Dual-domain loss (pixel + FFT) for sharp edges.",
    "6. Linear memory scaling for large images.",
], Inches(0.4), Inches(0.95), Inches(5.8), Inches(3.0), sz=17, boxed=True)
txt(s, [
    ("Future Work:", True),
    "1. Extend to Hyperspectral SR datasets.",
    "2. Scale to SRMamba-T-L for classical SR.",
    "3. Real-world degradation beyond bicubic.",
    "4. Deploy on edge devices (linear memory).",
    "5. Combine with diffusion models.",
    "6. Video SR with temporal Mamba layers.",
], Inches(0.4), Inches(4.2), Inches(5.8), Inches(2.8), sz=17, boxed=True)

# ══════════════════════════════════════════
# 26. REFERENCES
# ══════════════════════════════════════════
s = S(); bg(s); hdr(s, "References"); print(f"  {n}. References")
txt(s, [
    "[1]  Dong et al., \"SRCNN\" — ECCV 2014",
    "[2]  Dong et al., \"FSRCNN\" — ECCV 2016",
    "[3]  Kim et al., \"VDSR\" — CVPR 2016",
    "[4]  Ledig et al., \"SRGAN\" — CVPR 2017",
    "[5]  Wang et al., \"ESRGAN\" — ECCV Workshop 2018",
    "[6]  Liang et al., \"SwinIR\" — ICCV 2021",
    "[7]  Guo et al., \"MambaIR\" — ECCV 2024",
    "[8]  Gu & Dao, \"Mamba: Selective SSM\" — NeurIPS 2023",
    "[9]  Liu et al., \"SRMamba-T\" — Neurocomputing 2025",
    "[10] Vaswani et al., \"Attention Is All You Need\" — NeurIPS 2017",
], Inches(0.4), Inches(0.95), Inches(8.5), Inches(6.0), sz=17, boxed=True)

# ══════════════════════════════════════════
# 27. THANK YOU
# ══════════════════════════════════════════
s = S(); print(f"  {n}. Thank You")
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x1A,0x1A,0x2E); r.line.fill.background()
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(4.83), Inches(5.5))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = OG; tri.line.fill.background()
for row in range(3):
    for col in range(3):
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9+col*0.9), Inches(3.5+row*0.7), Inches(0.7), Inches(0.55))
        sq.line.color.rgb = RGBColor(0x55,0x55,0x55); sq.line.width = Pt(1); sq.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.5), Inches(5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = OG; ln.line.fill.background()
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(1.2))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(64); p.font.bold = True; p.font.color.rgb = OG; p.font.name = FONT
tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WH; p2.font.name = FONT
tb3 = s.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(7), Inches(0.5))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; p3.text = "Vujja Punith Sai  |  S20240010260  |  Instructor: Arun PV"
p3.font.size = Pt(18); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p3.font.name = FONT

# ══════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v10_PROFESSOR.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
