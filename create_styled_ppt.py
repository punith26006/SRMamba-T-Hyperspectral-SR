import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import glob

def extract_text_from_slide(slide):
    title = ""
    bullets = []
    
    if slide.shapes.title:
        title = slide.shapes.title.text

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                bullets.append(paragraph.text.strip())
                
    return title, bullets

def generate_new_presentation(source_file, output_file):
    if not os.path.exists(source_file):
        print(f"Source file {source_file} not found.")
        return

    # Load source
    try:
        source_prs = Presentation(source_file)
    except Exception as e:
        print(f"Error reading source presentation: {e}")
        return
        
    slides_data = []
    for slide in source_prs.slides:
        title, bullets = extract_text_from_slide(slide)
        slides_data.append({"title": title, "bullets": bullets})

    # Create new presentation
    new_prs = Presentation()
    
    # We will use layout 1 (Title and Content) for most, layout 0 for title
    title_layout = new_prs.slide_layouts[0]
    content_layout = new_prs.slide_layouts[1]

    for i, data in enumerate(slides_data):
        layout = title_layout if i == 0 else content_layout
        slide = new_prs.slides.add_slide(layout)
        
        # Set Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = data["title"] if data["title"] else f"Slide {i+1}"
            title_frame = title_shape.text_frame
            for paragraph in title_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Century Gothic'
                    run.font.size = Pt(40)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
        
        # Set Bullets
        if len(slide.placeholders) > 1 and data["bullets"]:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default text
            
            for idx, bullet in enumerate(data["bullets"]):
                p = tf.add_paragraph()
                p.text = bullet
                p.alignment = PP_ALIGN.LEFT
                p.level = 0
                for run in p.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(24)
                    run.font.color.rgb = RGBColor(64, 64, 64) # Dark Gray

    new_prs.save(output_file)
    print(f"Successfully generated styled template at: {output_file}")

if __name__ == "__main__":
    # Find the most recent or main PPT to use
    search_path = os.path.join("PPTs", "*Final*.pptx")
    files = glob.glob(search_path)
    
    source_ppt = files[0] if files else os.path.join("PPTs", "SRMamba_T_BTP_Final.pptx")
    output_ppt = "New_Styled_Template.pptx"
    
    generate_new_presentation(source_ppt, output_ppt)
