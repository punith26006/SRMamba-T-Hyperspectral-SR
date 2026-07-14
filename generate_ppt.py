from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation(filename="BTP_Presentation_Template.pptx"):
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    slides_data = [
        {
            "layout": title_slide_layout,
            "title": "Automated Dataset Construction and Model Training",
            "content": ["Your Name / Team Members", "Supervisor/Guide Name", "Institution/Department", "Date"]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Introduction / Background",
            "content": [
                "Briefly introduce the domain or topic.",
                "Explain why this area is important (real-world context).",
                "Keep points short and impactful."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Problem Statement",
            "content": [
                "Clearly state the specific problem you are trying to solve.",
                "Highlight limitations or gaps in existing solutions.",
                "Use a single, clear, bold sentence summarizing the core problem."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Objectives",
            "content": [
                "Primary goal 1: (e.g., Curate a novel dataset for X)",
                "Primary goal 2: (e.g., Evaluate ML models like v3, v4, v9)",
                "What defines the success of this project?"
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Dataset Description",
            "content": [
                "Source: Where did the data come from?",
                "Size & Features: Number of samples, classes, data type.",
                "Preprocessing: Cleaning, balancing, or augmentation applied.",
                "[Insert sample image or distribution chart here]"
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Methodology / Architecture",
            "content": [
                "How did you approach the problem?",
                "Data Collection -> Preprocessing -> Model Training -> Evaluation",
                "[Insert Flowchart or Block Diagram here - Highly Recommended]"
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Model Details",
            "content": [
                "Algorithms/Models used (e.g., CNN, Random Forest, Custom).",
                "Why were these models chosen?",
                "Highlight specific hyperparameter choices or novel tweaks."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Experimental Setup & Training",
            "content": [
                "Hardware/Software used (e.g., Google Colab, PyTorch, Scikit-learn).",
                "Training details (epochs, batch size, learning rate).",
                "Train/Test split ratio."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Results & Analysis",
            "content": [
                "Summary of performance (Accuracy, Precision, Recall, F1-Score).",
                "Comparison of different builds (build_v3 vs build_final).",
                "[Insert Bar charts or readable table here]"
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Visualizing Results",
            "content": [
                "Confusion Matrix.",
                "ROC Curve or Loss vs. Epoch graph.",
                "Interesting findings (e.g., specific classes the model struggles with)."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Challenges Faced",
            "content": [
                "What went wrong? (e.g., Data imbalance, overfitting, hardware constraints).",
                "How were they overcome? (e.g., augmentation, dropout)."
            ]
        },
        {
            "layout": bullet_slide_layout,
            "title": "Conclusion & Future Scope",
            "content": [
                "Conclusion: Summary of achievements. Were objectives met?",
                "Future Scope: What could be done next? (e.g., Deployment, more data)."
            ]
        },
        {
            "layout": title_slide_layout,
            "title": "Thank You",
            "content": ["Questions?"]
        }
    ]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(slide_info["layout"])
        title = slide.shapes.title
        title.text = slide_info["title"]
        
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            for idx, bullet in enumerate(slide_info["content"]):
                if idx == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    
    save_path = os.path.join(os.getcwd(), filename)
    prs.save(save_path)
    print(f"Presentation saved to {save_path}")

if __name__ == "__main__":
    create_presentation()
