"""
FULL 46-Slide Professional BTP Presentation
- Part I: All literature review content from friend's slides (typed text)
- Part II: ALL SRMamba-T architecture diagrams (full-page with header)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
IMG = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")
FRD = os.path.join(BASE, "_temp_scripts", "slide_images_work", "friend_slides")

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DARK = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
ACC2 = RGBColor(0xE8, 0x5D, 0x04)
GOLD = RGBColor(0xFF, 0xBE, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBGR = RGBColor(0xF0, 0xF4, 0xF8)
TDRK = RGBColor(0x1A, 0x1A, 0x2E)
TBODY = RGBColor(0x33, 0x3D, 0x4D)
GREEN = RGBColor(0x2D, 0xA0, 0x5B)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

def bg(s, c=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def hdr(s, t, sub=None):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.1))
    b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), W, Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(13); p2.font.color.rgb = ACCENT; p2.font.italic = True

def txt(s, bullets, x, y, w, h, sz=15):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(b, tuple):
            p.text = b[0]; p.font.bold = b[2] if len(b) > 2 else False
        else:
            p.text = b
        p.font.size = Pt(sz); p.font.color.rgb = TBODY; p.space_after = Pt(5)

def col_hdr(s, lt, rt, y=Inches(1.4)):
    for (t, x) in [(lt, Inches(0.5)), (rt, Inches(6.9))]:
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(0.45))
        bx.fill.solid(); bx.fill.fore_color.rgb = DARK; bx.line.fill.background()
        tf = bx.text_frame; tf.paragraphs[0].text = t
        tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def card(s, title, points, x, y, w, h, color=ACCENT):
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bx.fill.solid(); bx.fill.fore_color.rgb = WHITE; bx.line.color.rgb = color; bx.line.width = Pt(1.5)
    hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
    hd.fill.solid(); hd.fill.fore_color.rgb = color; hd.line.fill.background()
    tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.4))
    tf = tb.text_frame; tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE
    bul = [f"•  {p}" for p in points]
    txt(s, bul, x + Inches(0.15), y + Inches(0.6), w - Inches(0.3), h - Inches(0.7), sz=13)

def divider(s, num, title, sub=""):
    bg(s, NAVY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.3), Inches(2.1), Inches(2.1))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.5), Inches(2.1), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.2))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(38); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub; p3.font.size = Pt(17); p3.font.color.rgb = ACCENT; p3.alignment = PP_ALIGN.CENTER

def img_slide(s, title, img_name, sub=None):
    """Full-page diagram slide with header bar"""
    bg(s)
    hdr(s, title, sub)
    path = os.path.join(IMG, img_name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(0.15), Inches(1.25), Inches(13.0), Inches(6.1))

def img_slide_2(s, title, img1, img2, sub=None):
    """Two diagrams side by side"""
    bg(s)
    hdr(s, title, sub)
    p1 = os.path.join(IMG, img1)
    p2 = os.path.join(IMG, img2)
    if os.path.exists(p1): s.shapes.add_picture(p1, Inches(0.1), Inches(1.25), Inches(6.5), Inches(6.1))
    if os.path.exists(p2): s.shapes.add_picture(p2, Inches(6.7), Inches(1.25), Inches(6.5), Inches(6.1))

n = 0
def S():
    global n; n += 1; print(f"  Slide {n}"); return prs.slides.add_slide(BL)

# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
s = S()
bg(s, NAVY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.5), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Exploration of Deep Learning Models"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "for Image Super Resolution"; p2.font.size = Pt(44); p2.font.bold = True; p2.font.color.rgb = ACCENT
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.6), Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
tb2 = s.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(3))
tf2 = tb2.text_frame
for t, sz, c, b in [("B.Tech Project Presentation", 20, GOLD, True), ("", 10, WHITE, False),
    ("Presented by:  Vujja Punith Sai", 22, WHITE, False), ("Roll No:  S20240010260", 18, RGBColor(0x99,0x99,0x99), False),
    ("Instructor:  Arun PV", 18, RGBColor(0x99,0x99,0x99), False)]:
    p = tf2.add_paragraph(); p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

# ════════════════════════════════════════════════════════════════
# SLIDE 2: CONTENTS
# ════════════════════════════════════════════════════════════════
s = S()
bg(s); hdr(s, "Contents")
left = ["Part I — Literature Review", "  1.  What is Super Resolution", "  2.  Need for Super Resolution",
    "  3.  Problem Statement", "  4.  Traditional SR Methods", "  5.  Deep Learning in SR",
    "  6.  Evolution of SR Models", "  7.  Multimodal Learning in SR",
    "  8.  CNN Models: SRCNN", "  9.  CNN Models: FSRCNN",
    "  10. CNN Models: VDSR", "  11. GAN Models: SRGAN",
    "  12. GAN Models: ESRGAN", "  13. Transformer Models: SwinIR",
    "  14. State Space Models: SR-Mamba", "  15. Evaluation Metrics (MSE, PSNR, SSIM)"]
right = ["Part II — Our Method: SRMamba-T", "  16. The Core Challenge in SISR",
    "  17. The Architecture Compromise", "  18. Hybrid Insight & ERF Analysis",
    "  19. SRMamba-T Architecture Overview", "  20. Stage-Wise Feature Flow",
    "  21. SSM Foundations & Discretization", "  22. Mamba Layer Anatomy",
    "  23. Asymmetric Mamba Mixer Block", "  24. Spatial Dependency Problem",
    "  25. Multi-Directional Selective Scan (MDSSM)", "  26. MDSSM: Preserving 2D Spatial Context",
    "  27. Transformer Layer & Attention", "  28. Feature Fusion & Reconstruction",
    "  29. Results, LAM, Ablation", "  30. Conclusion & Future Work"]
txt(s, left, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.8), sz=14)
txt(s, right, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.8), sz=14)
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.4), Inches(0.03), Inches(5.5))
dv.fill.solid(); dv.fill.fore_color.rgb = ACCENT; dv.line.fill.background()

# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER: PART I
# ════════════════════════════════════════════════════════════════
s = S(); divider(s, "I", "Literature Review", "Background & Related Work")

# ════════════════════════════════════════════════════════════════
# SLIDE 4: WHAT IS SR
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "1. What is Super Resolution?")
col_hdr(s, "Definition & Objectives", "Basic Workflow")
txt(s, [
    "Super Resolution (SR) is a Deep Learning and Computer Vision technique",
    "used to reconstruct a high-resolution image from one or more low-resolution inputs.",
    "",
    "It aims to recover:",
    "  •  Lost image details",
    "  •  Textures and edges",
    "  •  Structural information",
    "  •  Overall visual quality",
    "",
    "Main Objectives:",
    "  1. Improve Image Clarity — enhance sharpness and visibility",
    "  2. Recover Image Details — reconstruct missing fine details",
    "  3. Enhance Texture & Edges — sharper, more realistic edges",
    "  4. Improve Perceptual Quality — natural-looking images",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(5), sz=15)
txt(s, [
    "Basic SR Workflow:",
    "",
    "  Low Resolution Image",
    "          ↓",
    "  Feature Extraction",
    "          ↓",
    "  Super Resolution Model",
    "          ↓",
    "  Image Reconstruction Layer",
    "          ↓",
    "  High Resolution Image",
], Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), sz=16)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: NEED FOR SR
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "2. Need for Super Resolution")
col_hdr(s, "Problems in Low Resolution Images", "Real-World Applications")
txt(s, [
    "Low-resolution images suffer from several quality limitations:",
    "",
    "Major Problems:",
    "  •  Blur and loss of sharpness",
    "  •  Pixelation during zooming",
    "  •  Loss of fine textures and edges",
    "  •  Poor object visibility",
    "  •  Reduced structural details",
    "  •  Noise and distortion",
    "  •  Information loss during compression",
    "  •  Degraded quality in low-light conditions",
    "",
    "Industry Usage: Samsung, Google, Apple use SR in",
    "AI cameras, satellite processing, video restoration.",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(5), sz=14)
txt(s, [
    "* Satellite Imaging: terrain analysis, urban planning,",
    "     environmental monitoring",
    "* Surveillance: face recognition, vehicle number plate",
    "     enhancement",
    "* Medical Imaging: MRI and CT scan enhancement,",
    "     better disease detection, diagnostic accuracy",
    "* Mobile Camera: AI zoom, night photography,",
    "     image sharpening",
    "* Video Upscaling: converting low-quality videos to HD/4K",
    "* Remote Sensing: land use analysis, disaster monitoring,",
    "     geological mapping",
], Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), sz=14)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════
s = S(); bg(s, LBGR); hdr(s, "3. Problem Statement")
bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.3))
bx.fill.solid(); bx.fill.fore_color.rgb = WHITE; bx.line.color.rgb = ACCENT; bx.line.width = Pt(2)
txt(s, [
    "Most open-source data have coarse spatial resolution, and it is advantageous to use different",
    "sensors for effective enhancement; however, the sensors have differences in spatial resolution,",
    "spectral characteristics, noise levels, and often lack spatial or temporal overlap."
], Inches(1.2), Inches(1.9), Inches(10.9), Inches(1.8), sz=20)
txt(s, [
    "Key Challenges in Single Image Super Resolution (SISR):",
    "",
    "  •  Recovering lost high-frequency details from degraded inputs",
    "  •  Balancing reconstruction quality vs. computational cost",
    "  •  Maintaining global context while preserving local pixel fidelity",
    "  •  Existing CNN models have limited receptive fields (local only)",
    "  •  Transformer models suffer from quadratic O(n²) computational complexity",
    "  •  Pure Mamba models lose 2D spatial context due to 1D sequential scanning",
    "  •  Need for a hybrid approach that combines efficiency with global awareness",
], Inches(0.8), Inches(4.2), Inches(11.7), Inches(3), sz=16)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: TRADITIONAL METHODS
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "4. Traditional Super Resolution Methods")
txt(s, ["Before Deep Learning, SR was performed using interpolation-based image scaling methods.",
    "These methods estimate new pixel values mathematically to enlarge image resolution."], Inches(0.5), Inches(1.3), Inches(12), Inches(0.6), sz=15)
card(s, "1) Nearest Neighbour", ["Simplest interpolation technique", "Copies nearest pixel value",
    "Computationally very fast", "", "Advantages: Low cost, fast", "Limitations: Blocky, jagged edges, poor quality"],
    Inches(0.3), Inches(2.1), Inches(4.0), Inches(5.0), ACCENT)
card(s, "2) Bilinear Interpolation", ["Calculates values using neighbouring pixels", "Smoother than nearest neighbour",
    "", "Advantages: Smoother image, better visual quality", "Limitations: Loss of sharpness, blurred edges, texture loss"],
    Inches(4.6), Inches(2.1), Inches(4.0), Inches(5.0), ACC2)
card(s, "3) Bicubic Interpolation", ["Uses surrounding 16 pixels", "Produces smoother outputs",
    "", "Advantages: Better smoothness, improved edge continuity", "Limitations: Cannot reconstruct missing details, blurry textures"],
    Inches(8.9), Inches(2.1), Inches(4.0), Inches(5.0), GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: DEEP LEARNING IN SR
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "5. Deep Learning in Super Resolution")
col_hdr(s, "Why Deep Learning for SR?", "What Deep Learning Models Learn")
txt(s, [
    "Traditional interpolation methods cannot effectively",
    "recover missing image information.",
    "",
    "Deep Learning models improve SR by learning:",
    "  •  Textures",
    "  •  Edges",
    "  •  Spatial patterns",
    "  •  Feature Representation",
    "  •  Structural relationships",
    "directly from large image datasets.",
    "",
    "Main Idea: Instead of simply enlarging pixels",
    "mathematically, DL models learn the mapping from",
    "Low-Resolution images to High-Resolution images.",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(5), sz=15)
txt(s, [
    "1) Texture Pattern:",
    "   Learning realistic textures such as buildings,",
    "   roads, vegetation, facial details.",
    "",
    "2) Edge Information:",
    "   Detects object boundaries and sharp transitions.",
    "",
    "3) Image Structures:",
    "   Understands spatial relationships between",
    "   image regions.",
    "",
    "4) Feature Representation:",
    "   Extracts meaningful hierarchical image features",
    "   automatically.",
], Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), sz=15)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: EVOLUTION OF SR MODELS
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "6. Evolution of Super Resolution Models")
txt(s, [
    "Super Resolution has evolved significantly over time with advancements in Deep Learning architectures.",
    "The evolution mainly focused on improving: Reconstruction quality, Texture generation, Perceptual realism, Computational efficiency.",
], Inches(0.5), Inches(1.3), Inches(12), Inches(0.8), sz=15)
timeline = [
    ("Traditional Models", "Nearest Neighbour, Bilinear, Bicubic Interpolation", RGBColor(0x99,0x99,0x99)),
    ("CNN-Based Models", "SRCNN, FSRCNN, VDSR — Local feature extraction", ACCENT),
    ("GAN-Based Models", "SRGAN, ESRGAN — Perceptual realism", ACC2),
    ("Transformer-Based", "SwinIR — Global self-attention, shifted windows", RGBColor(0xE9,0x1E,0x63)),
    ("State Space Models", "SR-Mamba — Linear complexity, state transitions", GREEN),
]
for i, (title, desc, color) in enumerate(timeline):
    y = Inches(2.4 + i * 0.95)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), y + Inches(0.1), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    if i < len(timeline) - 1:
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.63), y + Inches(0.45), Inches(0.08), Inches(0.55))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD); ln.line.fill.background()
    tb = s.shapes.add_textbox(Inches(3.2), y, Inches(3), Inches(0.45))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = color
    tb2 = s.shapes.add_textbox(Inches(6.5), y, Inches(6), Inches(0.45))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = desc; p2.font.size = Pt(15); p2.font.color.rgb = TBODY

# ════════════════════════════════════════════════════════════════
# SLIDE 10: MULTIMODAL LEARNING
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "7. Relation to Multimodal Learning")
col_hdr(s, "Single-Modality Super Resolution", "What is Multimodal Learning?")
txt(s, [
    "Traditional SR systems mainly use:",
    "  •  A single low-resolution image",
    "  •  One data modality",
    "",
    "This is also known as Single-Modality",
    "Super Resolution.",
    "",
    "In this approach:",
    "  •  Only visual image information is used",
    "  •  Reconstruction depends entirely on",
    "     image features",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(5), sz=15)
txt(s, [
    "Multimodal Learning combines:",
    "  •  Multiple data sources",
    "  •  Different modalities",
    "",
    "To improve learning and prediction quality.",
    "",
    "Example:",
    "  (Satellite Image + Elevation Data)",
    "  → Multimodal SR",
    "  → Enhanced High Resolution Output",
    "",
    "Advantages: Better context, richer features",
    "Limitations: Data alignment challenges,",
    "  higher complexity, limited paired datasets",
], Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), sz=15)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: SRCNN
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "8. CNN-based Models: SRCNN (Super Resolution CNN)")
txt(s, [
    "SRCNN (2014) — The first major Deep Learning model proposed for image Super Resolution.",
    "  •  End-to-end learning for SR using CNNs",
    "  •  Directly learns the mapping LR → HR using CNNs",
    "  •  Instead of using traditional interpolation methods",
], Inches(0.5), Inches(1.3), Inches(12), Inches(1.0), sz=15)
card(s, "Architecture: 3 Convolution Layers", [
    "1) Feature Extraction Layer:",
    "   Extracts low-level features (edges, corners, textures)",
    "",
    "2) Non-Linear Mapping Layer:",
    "   Maps extracted LR features into HR feature representations",
    "",
    "3) Reconstruction Layer:",
    "   Reconstructs the final high-resolution image from learned feature maps",
], Inches(0.3), Inches(2.5), Inches(6.0), Inches(4.7), ACCENT)
card(s, "Key Properties", [
    "•  First end-to-end DL approach for SR",
    "•  Simple 3-layer architecture",
    "•  Operates on bicubic-upsampled input",
    "•  Loss: MSE (pixel-wise)",
    "",
    "Limitations:",
    "•  Shallow network — limited capacity",
    "•  Small receptive field",
    "•  Slow (operates in HR space)",
    "•  Cannot capture complex patterns",
], Inches(6.6), Inches(2.5), Inches(6.3), Inches(4.7), ACC2)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: FSRCNN
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "9. CNN-based Models: FSRCNN (Fast SRCNN)")
card(s, "FSRCNN Architecture", [
    "•  Faster version of SRCNN (2016)",
    "•  Key improvement: operates on LR space directly",
    "•  Uses deconvolution layer for upsampling at the end",
    "•  5 components: Feature extraction → Shrinking →",
    "   Mapping → Expanding → Deconvolution",
    "",
    "Advantages over SRCNN:",
    "•  5× faster than SRCNN",
    "•  No need for pre-upsampling",
    "•  Smaller model size",
    "•  Adaptable to different scale factors",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), ACCENT)
card(s, "Key Innovations", [
    "1) LR-space processing:",
    "   All feature extraction happens in LR space,",
    "   reducing computational cost significantly",
    "",
    "2) Shrinking + Expanding layers:",
    "   Reduce channel dimensions before mapping,",
    "   then expand back — saves parameters",
    "",
    "3) Deconvolution for upsampling:",
    "   Learned upsampling at the very end,",
    "   instead of bicubic pre-processing",
    "",
    "Limitations: Still limited receptive field",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: VDSR
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "10. CNN-based Models: VDSR (Very Deep SR)")
card(s, "VDSR Architecture (2016)", [
    "•  Very Deep Super Resolution — 20 convolution layers",
    "•  Key innovation: Global Residual Learning",
    "•  Network learns: Residual = HR - LR",
    "•  Final output: HR = LR + Learned Residual",
    "",
    "Architecture Details:",
    "•  20 conv layers (3×3 filters, 64 channels)",
    "•  No pooling layers — preserves spatial resolution",
    "•  High learning rate + gradient clipping",
    "•  Multi-scale: single model handles ×2, ×3, ×4",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), ACCENT)
card(s, "Advantages & Limitations", [
    "Advantages:",
    "•  Much deeper than SRCNN (20 vs 3 layers)",
    "•  Residual learning speeds up convergence",
    "•  Multi-scale capability (one model, many scales)",
    "•  Better reconstruction quality",
    "•  Larger receptive field (41×41 pixels)",
    "",
    "Limitations:",
    "•  Still operates in HR space (slow)",
    "•  No skip connections between layers",
    "•  Limited to local feature extraction",
    "•  Cannot capture long-range dependencies",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), ACC2)

# ════════════════════════════════════════════════════════════════
# SLIDE 14: SRGAN
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "11. GAN-based Models: SRGAN")
txt(s, [
    "GAN-Based Models introduced perceptual learning and realistic texture generation into Super Resolution.",
    "Unlike CNN models that mainly focused on pixel-wise reconstruction (MSE/PSNR), GAN models focused on perceptual realism.",
], Inches(0.5), Inches(1.3), Inches(12), Inches(0.7), sz=14)
card(s, "What is GAN?", [
    "A GAN (Generative Adversarial Network) consists of:",
    "  •  Generator: Creates HR images from LR inputs",
    "  •  Discriminator: Determines if image is real or fake",
    "",
    "They compete with each other to generate realistic outputs.",
    "Goal: Generator fools the Discriminator.",
], Inches(0.3), Inches(2.3), Inches(6.0), Inches(4.8), ACC2)
card(s, "SRGAN Architecture", [
    "•  First GAN-based SR model (2017)",
    "•  Generator: Deep ResNet with residual blocks",
    "•  Discriminator: VGG-style classifier",
    "•  Loss: Perceptual loss + Adversarial loss",
    "",
    "Key Innovation: Perceptual Loss",
    "  Uses VGG feature maps instead of pixel MSE",
    "  → Sharper textures, visually realistic results",
    "",
    "Limitation: Training instability, artifacts",
], Inches(6.6), Inches(2.3), Inches(6.3), Inches(4.8), ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: ESRGAN
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "12. GAN-based Models: ESRGAN")
card(s, "ESRGAN: Enhanced SRGAN", [
    "•  Improved version of SRGAN (2018)",
    "•  Residual-in-Residual Dense Blocks (RRDB)",
    "•  Removed batch normalization for stable training",
    "•  Relativistic Discriminator: predicts relative",
    "   quality instead of absolute real/fake",
    "",
    "Architecture Improvements:",
    "•  RRDB: Multi-level residual connections",
    "•  Dense connections within each residual block",
    "•  Network interpolation between PSNR-oriented",
    "   and GAN-oriented models",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), ACC2)
card(s, "Results & Limitations", [
    "Advantages over SRGAN:",
    "•  Better texture reconstruction",
    "•  Higher perceptual quality",
    "•  More stable training",
    "•  Sharper, more realistic outputs",
    "•  Won PIRM2018 SR Challenge",
    "",
    "Limitations of GAN-based models:",
    "•  Training can be unstable (mode collapse)",
    "•  May generate hallucinated textures",
    "•  High computational cost",
    "•  Difficult to tune hyperparameters",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 16: SWINIR
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "13. Transformer-based Models: SwinIR")
txt(s, [
    "Transformer-based models introduced self-attention mechanisms and global feature learning into SR.",
    "Unlike CNNs (local features) and GANs (perceptual realism), Transformers learn long-range dependencies and global contextual relationships.",
], Inches(0.5), Inches(1.3), Inches(12), Inches(0.7), sz=14)
card(s, "SwinIR (Swin Transformer for Image Restoration)", [
    "•  Built on Vision Transformers (ViT)",
    "•  Uses Shifted Window Attention",
    "   instead of full global attention",
    "•  Self-attention determines:",
    "   - Which image regions are important",
    "   - How distant features are related",
    "   - How global structures should be reconstructed",
    "",
    "Architecture: Shallow Feature Extraction →",
    "  Deep Feature Extraction (RSTB blocks) →",
    "  HR Reconstruction",
], Inches(0.3), Inches(2.3), Inches(6.0), Inches(4.8), RGBColor(0xE9,0x1E,0x63))
card(s, "Advantages & Limitations", [
    "Advantages:",
    "•  Global feature learning via self-attention",
    "•  Excellent texture reconstruction quality",
    "•  Better than CNNs and GANs in PSNR/SSIM",
    "•  Shifted windows reduce complexity vs full attention",
    "",
    "Limitations:",
    "•  O(n²) complexity — high GPU memory usage",
    "•  Increased computational complexity",
    "•  Longer training time",
    "•  Struggles with very large input resolutions",
    "•  Window-based attention limits receptive field",
], Inches(6.6), Inches(2.3), Inches(6.3), Inches(4.8), ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 17: SR-MAMBA
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "14. State Space Models: SR-Mamba")
card(s, "State Space Models (SSMs)", [
    "SSMs are sequence modeling architectures:",
    "  •  Efficient long-range information propagation",
    "  •  Lower memory usage than Transformers",
    "  •  Scalable computation — O(n) linear",
    "",
    "Unlike Transformers, SSMs avoid expensive",
    "global attention operations.",
    "",
    "Key properties:",
    "  •  Updates hidden states sequentially",
    "  •  Preserves long-range information",
    "  •  Efficiently models feature relationships",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), GREEN)
card(s, "SR-Mamba (Super Resolution Mamba)", [
    "•  Modern SR architecture based on Mamba/S6",
    "•  Uses State Transition Mechanisms instead of",
    "   self-attention",
    "",
    "Main Idea: Use long-range image dependencies,",
    "contextual relationships, feature interactions",
    "via state transitions.",
    "",
    "Advantages: O(n) complexity, scalable, efficient",
    "",
    "Limitations:",
    "  •  1D sequential scan destroys 2D spatial context",
    "  •  Context loss at pixel boundaries",
    "  •  Cannot fully capture 2D spatial relationships",
    "  → Need for hybrid approach: SRMamba-T!",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 18: EVALUATION METRICS (MSE, PSNR)
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "15a. Evaluation Metrics: MSE & PSNR")
card(s, "MSE (Mean Square Error)", [
    "•  Measures pixel-wise reconstruction error",
    "•  Average squared difference between pixels",
    "",
    "Formula: MSE = (1/mn) Σ (I(i,j) - K(i,j))²",
    "",
    "Where:",
    "  m, n = image dimensions",
    "  I(i,j) = pixel value of original HR image",
    "  K(i,j) = pixel value of reconstructed image",
    "",
    "•  Lower MSE = better reconstruction quality",
    "•  Commonly used in SRCNN, FSRCNN, VDSR",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), ACCENT)
card(s, "PSNR (Peak Signal-to-Noise Ratio)", [
    "•  Measures reconstruction quality in dB",
    "•  Logarithmic scale based on MSE",
    "",
    "Formula: PSNR = 10 · log₁₀ (MAX² / MSE)",
    "",
    "Where:",
    "  MAX = maximum possible pixel value (255)",
    "  MSE = mean squared error",
    "",
    "•  Higher PSNR = better image fidelity",
    "•  Inversely related to MSE",
    "•  Industry standard metric for SR evaluation",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), ACC2)

# ════════════════════════════════════════════════════════════════
# SLIDE 19: EVALUATION METRICS (SSIM)
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "15b. Evaluation Metrics: SSIM & Summary")
card(s, "SSIM (Structural Similarity Index)", [
    "•  Measures structural & perceptual quality",
    "•  Compares: luminance, contrast, structure",
    "•  More aligned with human visual perception",
    "",
    "SSIM(x,y) = f(luminance, contrast, structure)",
    "",
    "Interpretation:",
    "  SSIM → 1: high similarity, good quality",
    "  SSIM → 0: poor similarity, bad quality",
    "",
    "Advantages over MSE/PSNR:",
    "  Focuses on visual perception, not just pixel error",
], Inches(0.3), Inches(1.4), Inches(6.0), Inches(5.8), GREEN)
card(s, "Literature Conclusion", [
    "SR has evolved from:",
    "  •  Traditional interpolation methods",
    "  •  CNN-based architectures",
    "  •  GAN-based models",
    "  •  Transformer models",
    "  •  State Space models",
    "",
    "Modern architectures significantly improve:",
    "  •  Image quality and texture reconstruction",
    "  •  Perceptual realism",
    "  •  Computational efficiency",
    "",
    "Deep Learning has become a major driving force",
    "in modern image enhancement systems.",
    "",
    "→ Next: Our hybrid approach — SRMamba-T",
], Inches(6.6), Inches(1.4), Inches(6.3), Inches(5.8), DARK)

# ════════════════════════════════════════════════════════════════
# SECTION DIVIDER: PART II
# ════════════════════════════════════════════════════════════════
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution")

# ════════════════════════════════════════════════════════════════
# PART II: ALL ARCHITECTURE DIAGRAM SLIDES
# Each uses the original high-quality diagram as full-page with header
# ════════════════════════════════════════════════════════════════
diagram_slides = [
    ("slide_02.png", "16. The Core Challenge: SISR is Ill-Posed"),
    ("slide_03.png", "17. The Dimensional Perspective"),
    ("slide_04.png", "18. The Architecture Compromise in SISR"),
    ("slide_05.png", "19. Comparing CNN, Transformer & Mamba"),
    ("slide_06.png", "20. Effective Receptive Field (ERF) Analysis"),
    ("slide_07.png", "21. The Hybrid Insight: Why Combine Mamba + Transformer?"),
    ("slide_08.png", "22. Hybrid Flow Logic: Sequence Before Space"),
    ("slide_09.png", "23. The Complete SRMamba-T Architecture"),
    ("slide_10.png", "24. Stage-Wise Feature Flow"),
    ("slide_11.png", "25. Foundations: Understanding State Space Models (SSMs)"),
    ("slide_12.png", "26. SSM Discretization for Deep Learning"),
    ("slide_13.png", "27. Equivalence to Linear Convolution"),
    ("slide_14.png", "28. Anatomy and Mathematics of the Mamba Layer"),
    ("slide_15.png", "29. The Asymmetric Mamba Mixer Block"),
    ("slide_16.png", "30. The Spatial Dependency Problem"),
    ("slide_17.png", "31. Multi-Directional Selective Scan Module (MDSSM)"),
    ("slide_18.png", "32. MDSSM: Preserving 2D Spatial Context"),
    ("slide_19.png", "33. Inside the Transformer Layer"),
    ("slide_20.png", "34. Un-Windowed Global Self-Attention"),
    ("slide_21.png", "35. Feature Fusion Module (FFM)"),
    ("slide_22.png", "36. Adaptive Reconstruction & Dual-Domain Loss"),
    ("slide_23.png", "37. Training Pipeline & Setup"),
    ("slide_24.png", "38. Quantitative Results: Shattering the Computational Ceiling"),
    ("slide_25.png", "39. Qualitative Validation: Dominating Complex Textures"),
    ("slide_26.png", "40. Inside the Mind of the Model: LAM Analysis"),
    ("slide_27.png", "41. Hardware Efficiency and Memory Scaling"),
    ("slide_28.png", "42. Ablation: The Empirical Logic of Layer Ordering"),
]

for img_name, title in diagram_slides:
    s = S()
    bg(s)
    hdr(s, title)
    path = os.path.join(IMG, img_name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(0.1), Inches(1.2), Inches(13.1), Inches(6.15))

# ════════════════════════════════════════════════════════════════
# SLIDE: CONCLUSION
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "43. Conclusion & Future Work")
col_hdr(s, "Key Contributions", "Future Work")
txt(s, [
    "✓  Proposed SRMamba-T: first hybrid Mamba-Transformer",
    "   architecture for image Super Resolution",
    "",
    "✓  MDSSM: Channel-split 4-directional scanning",
    "   preserves 2D spatial context at O(n) cost",
    "",
    "✓  Asymmetric layer placement: Mamba (early) +",
    "   Transformer (late) achieves optimal efficiency",
    "",
    "✓  +0.28 dB PSNR over MambaIR with 38.7% fewer MACs",
    "",
    "✓  Dual-domain loss (Pixel + FFT) prevents",
    "   frequency-domain artifacts",
    "",
    "✓  Linear memory scaling for large resolutions",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(5), sz=15)
txt(s, [
    "•  Extend to Hyperspectral Image Super Resolution",
    "   (PaviaC, PaviaU datasets)",
    "",
    "•  Scale to larger models (SRMamba-T-L) for",
    "   classical SR tasks",
    "",
    "•  Explore multi-scale U-Net with downsampling",
    "   between encoder layers",
    "",
    "•  Investigate real-world degradation models",
    "   beyond bicubic downsampling",
    "",
    "•  Deploy on edge devices leveraging Mamba's",
    "   linear memory scaling",
    "",
    "•  Combine with diffusion models for further",
    "   perceptual quality improvement",
], Inches(6.9), Inches(2.1), Inches(5.8), Inches(5), sz=15)

# ════════════════════════════════════════════════════════════════
# SLIDE: REFERENCES
# ════════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "44. References")
txt(s, [
    "[1]  Dong et al., \"Image Super-Resolution Using Deep Convolutional Networks\" (SRCNN), ECCV 2014",
    "[2]  Dong et al., \"Accelerating the Super-Resolution CNN\" (FSRCNN), ECCV 2016",
    "[3]  Kim et al., \"Accurate Image Super-Resolution Using Very Deep Networks\" (VDSR), CVPR 2016",
    "[4]  Ledig et al., \"Photo-Realistic SISR Using a GAN\" (SRGAN), CVPR 2017",
    "[5]  Wang et al., \"ESRGAN: Enhanced Super-Resolution GANs\", ECCV Workshop 2018",
    "[6]  Liang et al., \"SwinIR: Image Restoration Using Swin Transformer\", ICCV Workshop 2021",
    "[7]  Guo et al., \"MambaIR: A Simple Baseline for Image Restoration with SSM\", ECCV 2024",
    "[8]  SRMamba-T: Hybrid Mamba-Transformer Network for Single Image Super-Resolution, 2024",
    "[9]  Gu & Dao, \"Mamba: Linear-Time Sequence Modeling with Selective State Spaces\", 2023",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017",
], Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), sz=14)

# ════════════════════════════════════════════════════════════════
# SLIDE: THANK YOU
# ════════════════════════════════════════════════════════════════
s = S(); bg(s, NAVY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
tb3 = s.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]
p3.text = "Vujja Punith Sai  |  S20240010260  |  Guide: Arun PV"
p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p3.alignment = PP_ALIGN.CENTER

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_Full_46.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
print(">>> DONE!")
