"""
FINAL 40+ Slide Presentation
- Part I: Friend's slides used as IMAGES (they have architecture diagrams)
- Part II: SRMamba-T architecture diagrams (non-redundant)
- Professional header bar template on every slide
- No redundancy
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\vujja\OneDrive\Documents\OneDrive\Desktop\btpdataset"
IMG = os.path.join(BASE, "_temp_scripts", "slide_images_work", "slide_images")
FRD = os.path.join(BASE, "_temp_scripts", "slide_images_work", "friend_slides")

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
GOLD = RGBColor(0xFF, 0xBE, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBGR = RGBColor(0xF0, 0xF4, 0xF8)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BL = prs.slide_layouts[6]
n = 0

def S():
    global n; n += 1; return prs.slides.add_slide(BL)

def bg(s, c=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def hdr(s, t):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.85))
    b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), W, Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE

def divider(s, num, title, sub=""):
    bg(s, NAVY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.3), Inches(2.1), Inches(2.1))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.5), Inches(2.1), Inches(1.7))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.2))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title
    p2.font.size = Pt(38); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf2.add_paragraph(); p3.text = sub; p3.font.size = Pt(17); p3.font.color.rgb = ACCENT; p3.alignment = PP_ALIGN.CENTER

def friend_slide(title, slide_num):
    """Full-page friend slide image with header"""
    s = S(); bg(s); hdr(s, title)
    path = os.path.join(FRD, f"slide_{slide_num:02d}_rot.png")
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(0.1), Inches(0.95), Inches(13.1), Inches(6.4))
    print(f"  Slide {n}: {title} [friend #{slide_num}]")

def friend_slide_2(title, s1, s2):
    """Two friend slides side by side"""
    s = S(); bg(s); hdr(s, title)
    p1 = os.path.join(FRD, f"slide_{s1:02d}_rot.png")
    p2 = os.path.join(FRD, f"slide_{s2:02d}_rot.png")
    if os.path.exists(p1): s.shapes.add_picture(p1, Inches(0.1), Inches(0.95), Inches(6.5), Inches(6.4))
    if os.path.exists(p2): s.shapes.add_picture(p2, Inches(6.7), Inches(0.95), Inches(6.5), Inches(6.4))
    print(f"  Slide {n}: {title} [friend #{s1}+#{s2}]")

def my_slide(title, slide_num):
    """Full-page my architecture diagram with header"""
    s = S(); bg(s); hdr(s, title)
    path = os.path.join(IMG, f"slide_{slide_num:02d}.png")
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(0.1), Inches(0.95), Inches(13.1), Inches(6.4))
    print(f"  Slide {n}: {title} [my #{slide_num}]")

def my_slide_2(title, s1, s2):
    """Two of my slides side by side"""
    s = S(); bg(s); hdr(s, title)
    p1 = os.path.join(IMG, f"slide_{s1:02d}.png")
    p2 = os.path.join(IMG, f"slide_{s2:02d}.png")
    if os.path.exists(p1): s.shapes.add_picture(p1, Inches(0.1), Inches(0.95), Inches(6.5), Inches(6.4))
    if os.path.exists(p2): s.shapes.add_picture(p2, Inches(6.7), Inches(0.95), Inches(6.5), Inches(6.4))
    print(f"  Slide {n}: {title} [my #{s1}+#{s2}]")

def txt_box(s, text, x, y, w, h, sz=16, color=RGBColor(0x33,0x3D,0x4D)):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color; p.space_after = Pt(4)

print("=" * 60)
print("BUILDING FINAL PRESENTATION")
print("=" * 60)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
s = S(); bg(s, NAVY); print(f"  Slide {n}: Title")
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.5), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Exploration of Deep Learning Models"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "for Image Super Resolution"; p2.font.size = Pt(44); p2.font.bold = True; p2.font.color.rgb = ACCENT
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.6), Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
tb2 = s.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(3))
tf2 = tb2.text_frame
for t, sz, c, b in [("B.Tech Project Presentation", 20, GOLD, True), ("", 10, WHITE, False),
    ("Presented by:  Vujja Punith Sai", 22, WHITE, False), ("Roll No:  S20240010260", 18, RGBColor(0x99,0x99,0x99), False),
    ("Instructor:  Arun PV", 18, RGBColor(0x99,0x99,0x99), False)]:
    p = tf2.add_paragraph(); p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

# ══════════════════════════════════════════════════════════════
# SLIDE 2: CONTENTS
# ══════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "Contents"); print(f"  Slide {n}: Contents")
left = ["Part I — Literature Review", "", "  1.  What is Super Resolution",
    "  2.  Need for Super Resolution", "  3.  Problem Statement",
    "  4.  Traditional SR Methods", "  5.  Deep Learning in SR",
    "  6.  Evolution of SR Models", "  7.  Multimodal Learning",
    "  8.  CNN: SRCNN (Architecture + Loss)", "  9.  CNN: FSRCNN",
    "  10. CNN: FSRCNN (Architecture)", "  11. CNN: VDSR",
    "  12. GAN: Introduction", "  13. GAN: SRGAN (Adversarial Learning)",
    "  14. GAN: SRGAN (Architecture)", "  15. GAN: ESRGAN",
    "  16. Transformer: SwinIR", "  17. SwinIR (Architecture)",
    "  18. State Space: SR-Mamba", "  19. SR-Mamba (Architecture)",
    "  20. Evaluation Metrics", "  21. Literature Conclusion"]
right = ["Part II — Our Method: SRMamba-T", "",
    "  22. The Core Challenge in SISR", "  23. Architecture Compromise",
    "  24. CNN vs Transformer vs Mamba", "  25. ERF Analysis + Hybrid Insight",
    "  26. SRMamba-T Architecture Overview", "  27. Stage-Wise Feature Flow",
    "  28. SSM Foundations", "  29. SSM Discretization",
    "  30. Mamba Layer Anatomy", "  31. Asymmetric Mamba Mixer",
    "  32. MDSSM Module", "  33. MDSSM: 2D Spatial Context",
    "  34. Transformer Layer + Attention", "  35. Feature Fusion Module",
    "  36. Reconstruction & Loss", "  37. Training Pipeline",
    "  38. Quantitative Results", "  39. Qualitative Validation",
    "  40. LAM Analysis", "  41. Hardware Efficiency + Ablation",
    "  42. Conclusion & Future Work"]
txt_box(s, left, Inches(0.4), Inches(1.1), Inches(5.8), Inches(6.2), sz=12)
txt_box(s, right, Inches(6.8), Inches(1.1), Inches(5.8), Inches(6.2), sz=12)
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.1), Inches(0.03), Inches(6))
dv.fill.solid(); dv.fill.fore_color.rgb = ACCENT; dv.line.fill.background()

# ══════════════════════════════════════════════════════════════
# SECTION DIVIDER: PART I
# ══════════════════════════════════════════════════════════════
s = S(); divider(s, "I", "Literature Review", "Background & Related Work"); print(f"  Slide {n}: Section Divider Part I")

# ══════════════════════════════════════════════════════════════
# PART I: FRIEND'S SLIDES (with architecture diagrams)
# Each friend slide is used as a full-page image with header
# ══════════════════════════════════════════════════════════════

friend_slide("1. What is Super Resolution?", 3)
friend_slide("2. Need for Super Resolution", 4)

# Problem Statement (typed text)
s = S(); bg(s, LBGR); hdr(s, "3. Problem Statement"); print(f"  Slide {n}: Problem Statement")
bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2))
bx.fill.solid(); bx.fill.fore_color.rgb = WHITE; bx.line.color.rgb = ACCENT; bx.line.width = Pt(2)
txt_box(s, [
    "Most open-source data have coarse spatial resolution, and it is advantageous to use",
    "different sensors for effective enhancement; however, the sensors have differences in",
    "spatial resolution, spectral characteristics, noise levels, and often lack spatial or",
    "temporal overlap."
], Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.6), sz=20)
txt_box(s, [
    "Key Challenges in Single Image Super Resolution (SISR):",
    "  •  Recovering lost high-frequency details from degraded inputs",
    "  •  Balancing reconstruction quality vs. computational cost",
    "  •  CNN models: limited receptive fields (local only)",
    "  •  Transformer models: O(n²) computational complexity",
    "  •  Mamba models: 1D scan loses 2D spatial context",
    "  •  Need: hybrid approach combining efficiency + global awareness",
], Inches(0.8), Inches(4.0), Inches(11.7), Inches(3.2), sz=17)

friend_slide("4. Traditional Super Resolution Methods", 5)
friend_slide("5. Deep Learning in Super Resolution", 6)
friend_slide("6. Evolution of Super Resolution Models", 7)
friend_slide_2("7. Multimodal Learning: Single-Modality vs Multi-Modality", 8, 9)

# CNN Models (3 slides with architecture diagrams from friend)
friend_slide_2("8. CNN Models: SRCNN — Architecture & 3-Layer Pipeline", 10, 11)
friend_slide("9. CNN Models: FSRCNN — Fast Super Resolution CNN", 12)
friend_slide_2("10. FSRCNN Architecture & Deconvolution", 13, 14)
friend_slide_2("11. CNN Models: VDSR — Very Deep SR (20 Layers)", 16, 17)

# GAN Models (3 slides)
friend_slide_2("12. GAN-Based Models: Why GAN & What is GAN?", 19, 20)
friend_slide_2("13. SRGAN: Architecture, Generator & Discriminator", 21, 22)
friend_slide_2("14. ESRGAN: Enhanced Super Resolution GAN", 23, 24)

# Transformer (2 slides)
friend_slide("15. Transformer-Based Models: SwinIR", 26)
friend_slide_2("16. SwinIR Architecture & Limitations", 27, 28)

# State Space Models (2 slides)  
friend_slide("17. State Space Models: SR-Mamba", 29)
friend_slide_2("18. SR-Mamba Architecture & Limitations", 30, 31)

# Evaluation Metrics (2 slides)
friend_slide("19. Evaluation Metrics: MSE & PSNR", 32)
friend_slide("20. Evaluation Metrics: SSIM & Literature Conclusion", 33)

# ══════════════════════════════════════════════════════════════
# SECTION DIVIDER: PART II
# ══════════════════════════════════════════════════════════════
s = S(); divider(s, "II", "Our Method: SRMamba-T", "Hybrid Mamba-Transformer for Image Super Resolution"); print(f"  Slide {n}: Section Divider Part II")

# ══════════════════════════════════════════════════════════════
# PART II: ALL SRMamba-T ARCHITECTURE SLIDES
# Removed slide_03 (dimensional perspective - redundant)
# Removed slide_11,12,13 (SSM foundations - already covered in SR-Mamba section)
# Combined where possible
# ══════════════════════════════════════════════════════════════

my_slide("21. The Core Challenge: SISR is Ill-Posed", 2)
my_slide("22. The Architecture Compromise in SISR", 4)
my_slide("23. Comparing CNN, Transformer & Mamba", 5)
my_slide_2("24. Effective Receptive Field & Hybrid Insight", 6, 7)
my_slide("25. Hybrid Flow: Sequence Before Space", 8)
my_slide("26. The Complete SRMamba-T Architecture", 9)
my_slide("27. Stage-Wise Feature Flow", 10)
my_slide_2("28. SSM Foundations & Discretization", 11, 12)
my_slide("29. Mamba Layer: Anatomy & Mathematics", 14)
my_slide("30. The Asymmetric Mamba Mixer Block", 15)
my_slide("31. The Spatial Dependency Problem", 16)
my_slide("32. Multi-Directional Selective Scan Module (MDSSM)", 17)
my_slide("33. MDSSM: Preserving 2D Spatial Context", 18)
my_slide("34. Inside the Transformer Layer", 19)
my_slide("35. Un-Windowed Global Self-Attention", 20)
my_slide("36. Feature Fusion Module (FFM)", 21)
my_slide("37. Adaptive Reconstruction & Dual-Domain Loss", 22)
my_slide("38. Training Pipeline & Setup", 23)
my_slide("39. Quantitative Results: PSNR vs Compute", 24)
my_slide("40. Qualitative Validation: Dominating Complex Textures", 25)
my_slide("41. LAM Analysis: Diffusion Index Comparison", 26)
my_slide_2("42. Hardware Efficiency & Ablation Study", 27, 28)

# ══════════════════════════════════════════════════════════════
# CONCLUSION
# ══════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "43. Conclusion & Future Work"); print(f"  Slide {n}: Conclusion")
# Left
txt_box(s, [
    "Key Contributions:", "",
    "✓  Proposed SRMamba-T: first hybrid Mamba-Transformer",
    "   architecture for image Super Resolution", "",
    "✓  MDSSM: Channel-split 4-directional scanning",
    "   preserves 2D spatial context at O(n) cost", "",
    "✓  Asymmetric layer placement: Mamba (early) +",
    "   Transformer (late) for optimal efficiency", "",
    "✓  +0.28 dB PSNR over MambaIR with 38.7% fewer MACs", "",
    "✓  Dual-domain loss (Pixel + FFT) prevents",
    "   frequency-domain artifacts", "",
    "✓  Linear memory scaling for large resolutions",
], Inches(0.4), Inches(1.2), Inches(5.8), Inches(6), sz=15)
txt_box(s, [
    "Future Work:", "",
    "•  Extend to Hyperspectral Image Super Resolution",
    "   (PaviaC, PaviaU datasets)", "",
    "•  Scale to larger models (SRMamba-T-L) for",
    "   classical SR tasks", "",
    "•  Explore multi-scale U-Net with downsampling", "",
    "•  Investigate real-world degradation models",
    "   beyond bicubic downsampling", "",
    "•  Deploy on edge devices leveraging Mamba's",
    "   linear memory scaling", "",
    "•  Combine with diffusion models for perceptual",
    "   quality improvement",
], Inches(6.8), Inches(1.2), Inches(5.8), Inches(6), sz=15)

# ══════════════════════════════════════════════════════════════
# REFERENCES
# ══════════════════════════════════════════════════════════════
s = S(); bg(s); hdr(s, "44. References"); print(f"  Slide {n}: References")
txt_box(s, [
    "[1]  Dong et al., \"Image Super-Resolution Using Deep Convolutional Networks\" (SRCNN), ECCV 2014",
    "[2]  Dong et al., \"Accelerating the Super-Resolution CNN\" (FSRCNN), ECCV 2016",
    "[3]  Kim et al., \"Accurate Image Super-Resolution Using Very Deep Networks\" (VDSR), CVPR 2016",
    "[4]  Ledig et al., \"Photo-Realistic SISR Using a GAN\" (SRGAN), CVPR 2017",
    "[5]  Wang et al., \"ESRGAN: Enhanced Super-Resolution GANs\", ECCV Workshop 2018",
    "[6]  Liang et al., \"SwinIR: Image Restoration Using Swin Transformer\", ICCV Workshop 2021",
    "[7]  Guo et al., \"MambaIR: A Simple Baseline for Image Restoration with SSM\", ECCV 2024",
    "[8]  Gu & Dao, \"Mamba: Linear-Time Sequence Modeling with Selective State Spaces\", 2023",
    "[9]  SRMamba-T: Hybrid Mamba-Transformer Network for Single Image Super-Resolution, 2024",
    "[10] Vaswani et al., \"Attention Is All You Need\", NeurIPS 2017",
    "[11] Dosovitskiy et al., \"An Image is Worth 16x16 Words: Transformers for Image Recognition\", 2020",
], Inches(0.4), Inches(1.2), Inches(12.3), Inches(5.8), sz=14)

# ══════════════════════════════════════════════════════════════
# THANK YOU
# ══════════════════════════════════════════════════════════════
s = S(); bg(s, NAVY); print(f"  Slide {n}: Thank You")
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "Thank You!"
p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
tb3 = s.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1))
tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]
p3.text = "Vujja Punith Sai  |  S20240010260  |  Guide: Arun PV"
p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p3.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
out = os.path.join(BASE, "PPTs", "SRMamba_T_BTP_v3.pptx")
prs.save(out)
print(f"\n>>> Saved: {out}")
print(f">>> Total slides: {n}")
print(">>> DONE!")
