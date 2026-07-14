import os
import glob
from pptx import Presentation

def extract_text_from_slide(slide):
    title = ""
    bullets = []
    
    if slide.shapes.title:
        title = slide.shapes.title.text

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                bullets.append(paragraph.text.strip())
                
    return title, bullets

def generate_html_presentation(source_file, output_html):
    if not os.path.exists(source_file):
        print(f"Source file {source_file} not found.")
        return

    try:
        source_prs = Presentation(source_file)
    except Exception as e:
        print(f"Error reading source presentation: {e}")
        return

    # Start HTML with Reveal.js CDN
    html_content = """<!doctype html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <title>BTP Presentation</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
    
    <!-- Reveal.js Core CSS -->
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.css">
    
    <!-- Choose a Theme (e.g., black, white, league, beige, sky, night, serif, simple, solarized) -->
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/theme/night.min.css" id="theme">
    
    <style>
        .reveal h1, .reveal h2, .reveal h3 {
            color: #42b883; /* Custom green accent color */
            text-transform: none;
        }
        .reveal ul {
            font-size: 0.9em;
            line-height: 1.4;
        }
        .reveal li {
            margin-bottom: 15px;
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
"""

    # Generate Slides
    for i, slide in enumerate(source_prs.slides):
        title, bullets = extract_text_from_slide(slide)
        
        # Determine slide transition (alternate for a cool effect)
        transition = "slide" if i % 2 == 0 else "zoom"
        
        html_content += f'            <section data-transition="{transition}">\n'
        
        if title:
            # Main title on first slide gets h2, others get h3
            tag = "h2" if i == 0 else "h3"
            html_content += f'                <{tag}>{title}</{tag}>\n'
            
        if bullets:
            html_content += '                <ul>\n'
            for bullet in bullets:
                # Escape html chars
                bullet_clean = bullet.replace('<', '&lt;').replace('>', '&gt;')
                html_content += f'                    <li>{bullet_clean}</li>\n'
            html_content += '                </ul>\n'
            
        html_content += '            </section>\n'

    # Close HTML and add scripts
    html_content += """        </div>
    </div>

    <!-- Reveal.js Core JS -->
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            controls: true,
            progress: true,
            center: true,
            transition: 'slide', // none/fade/slide/convex/concave/zoom
            backgroundTransition: 'fade'
        });
    </script>
</body>
</html>
"""

    with open(output_html, "w", encoding="utf-8") as f:
        f.write(html_content)
        
    print(f"Successfully generated HTML presentation at: {output_html}")

if __name__ == "__main__":
    search_path = os.path.join("PPTs", "*Final*.pptx")
    files = glob.glob(search_path)
    
    source_ppt = files[0] if files else os.path.join("PPTs", "SRMamba_T_BTP_Final.pptx")
    output_html = "Web_Presentation.html"
    
    generate_html_presentation(source_ppt, output_html)
